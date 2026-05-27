#!/usr/bin/env python3
"""Export selected PPT slide tables to a landscape Word document (copy-paste)."""

from __future__ import annotations

import argparse
from pathlib import Path

from docx import Document
from docx.enum.section import WD_ORIENT
from docx.shared import Inches
from pptx import Presentation


def parse_slides(spec: str) -> list[int]:
    out: list[int] = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        out.append(int(part) - 1)
    return out


def export_tables(ppt_path: Path, out_path: Path, slide_indices: list[int], title: str) -> None:
    prs = Presentation(str(ppt_path))
    doc = Document()

    sec = doc.sections[0]
    sec.orientation = WD_ORIENT.LANDSCAPE
    sec.page_width, sec.page_height = sec.page_height, sec.page_width
    sec.left_margin = Inches(0.5)
    sec.right_margin = Inches(0.5)
    sec.top_margin = Inches(0.5)
    sec.bottom_margin = Inches(0.5)

    doc.add_heading(title, level=1)
    doc.add_paragraph(f"来源：{ppt_path.name}")

    for si in slide_indices:
        if si < 0 or si >= len(prs.slides):
            continue
        doc.add_heading(f"第 {si + 1} 页", level=2)
        tables = [sh.table for sh in prs.slides[si].shapes if sh.has_table]
        for ti, tbl in enumerate(tables, 1):
            doc.add_paragraph(f"表格 {ti}")
            rows, cols = len(tbl.rows), len(tbl.columns)
            wt = doc.add_table(rows=rows, cols=cols)
            wt.style = "Table Grid"
            for r in range(rows):
                for c in range(cols):
                    wt.cell(r, c).text = tbl.cell(r, c).text.strip()
            doc.add_paragraph("")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))
    print(f"Wrote {out_path}")


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--ppt", type=Path, required=True)
    ap.add_argument("--out", type=Path, required=True)
    ap.add_argument("--slides", type=str, default="1,2,4", help="1-based slide numbers, comma-separated")
    ap.add_argument("--title", type=str, default="PPT 表格导出（横版，便于复制）")
    args = ap.parse_args()
    export_tables(args.ppt, args.out, parse_slides(args.slides), args.title)


if __name__ == "__main__":
    main()

