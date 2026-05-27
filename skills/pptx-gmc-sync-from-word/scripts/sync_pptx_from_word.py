#!/usr/bin/env python3
"""Sync PPT slide tables from Word GMC / n / P-values.

TVAX immunogenicity deck structure:
- Slide 1 uses PPS sources (including M4 supplement)
- Slide 2 uses FAS sources (including M4 P-column supplement)
- Slide 4 table 1 uses C3/C2 analysis-set split
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from docx import Document
from pptx import Presentation

ARM_RE = re.compile(r"\((A1|A2|B1|B2|C1|C2|C3)\)")
MONTH_RE = re.compile(r"\d+")


def arm_code(text: str) -> str | None:
    m = ARM_RE.search(text)
    return m.group(1) if m else None


def month_num(text: str) -> int | None:
    m = MONTH_RE.search(text)
    return int(m.group()) if m else None


def table_rows(doc: Document, index0: int) -> list[list[str]]:
    t = doc.tables[index0]
    return [[c.text.strip() for c in r.cells] for r in t.rows]


def build_month_arm_map(rows: list[list[str]]) -> dict[tuple[int, str], dict[str, str]]:
    out: dict[tuple[int, str], dict[str, str]] = {}
    for vals in rows[1:]:
        if len(vals) < 13:
            continue
        mm = month_num(vals[0])
        a = arm_code(vals[1])
        if mm is None or a is None:
            continue
        out[(mm, a)] = {
            "n": vals[2],
            "gmc": vals[3],
            "p_a2b2": vals[6],
            "p_a2c1": vals[10],
            "p_b2c1": vals[12],
        }
    return out


def supplement_pps_month4_b_groups(doc: Document, table_index0: int) -> tuple[dict[tuple[int, str], dict[str, str]], str | None]:
    pps_updates: dict[tuple[int, str], dict[str, str]] = {}
    p_m4_p: str | None = None
    for vals in table_rows(doc, table_index0)[2:]:
        if not vals or vals[0].strip() != "PPS":
            continue
        mm = month_num(vals[1])
        if mm != 4:
            continue
        pps_updates[(4, "B1")] = {"n": vals[2], "gmc": vals[3]}
        pps_updates[(4, "B2")] = {"n": vals[5], "gmc": vals[6]}
        if len(vals) > 8:
            p_m4_p = vals[8]
    return pps_updates, p_m4_p


def supplement_fas_m4_p(doc: Document, table_index0: int) -> str | None:
    for vals in table_rows(doc, table_index0)[2:]:
        if not vals or vals[0].strip() != "FAS":
            continue
        if arm_code(vals[1]) != "A2":
            continue
        return vals[8] if len(vals) > 8 else None
    return None


def build_c3c2_map(doc: Document, table_index0: int) -> dict[tuple[str, int], dict[str, str]]:
    out: dict[tuple[str, int], dict[str, str]] = {}
    for vals in table_rows(doc, table_index0)[2:]:
        grp = vals[0].strip() if vals else ""
        if grp not in ("PPS", "FAS"):
            continue
        mm = month_num(vals[1])
        if mm is None:
            continue
        out[(grp, mm)] = {
            "c3_n": vals[2],
            "c3_gmc": vals[3],
            "c2_n": vals[5],
            "c2_gmc": vals[6],
            "p": vals[8],
        }
    return out


def set_cell(table, r: int, c: int, value: str | None) -> bool:
    if value is None or value == "":
        return False
    cell = table.cell(r, c)
    if cell.text.strip() == value:
        return False
    cell.text = value
    return True


def normalize_gmc_headers(table) -> None:
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            txt = table.cell(r, c).text
            if "校正" in txt and "GMC" in txt:
                table.cell(r, c).text = txt.replace("校正\nGMC", "GMC").replace("校正GMC", "GMC")


def update_main_slide_table(table, data: dict[tuple[int, str], dict[str, str]]) -> int:
    changes = 0
    arms = [("A1", 1, 2), ("A2", 3, 4), ("B1", 5, 6), ("B2", 7, 8), ("C1", 9, 10)]
    for r in range(3, len(table.rows)):
        mm = month_num(table.cell(r, 0).text)
        if mm is None:
            continue
        for a, nc, gc in arms:
            dct = data.get((mm, a))
            if not dct:
                continue
            if set_cell(table, r, nc, dct.get("n")):
                changes += 1
            if set_cell(table, r, gc, dct.get("gmc")):
                changes += 1

        if mm == 4:
            for c in (11, 12, 13):
                if set_cell(table, r, c, "/"):
                    changes += 1
            continue

        dct_a2 = data.get((mm, "A2"))
        if dct_a2:
            if set_cell(table, r, 11, dct_a2.get("p_a2b2")):
                changes += 1
            if table.cell(r, 9).text.strip() not in ("", "/"):
                if set_cell(table, r, 12, dct_a2.get("p_a2c1")):
                    changes += 1
                if set_cell(table, r, 13, dct_a2.get("p_b2c1")):
                    changes += 1
    return changes


def update_summary_table(table, data: dict[tuple[int, str], dict[str, str]], p_m4: str | None) -> int:
    changes = 0
    for r in range(3, len(table.rows)):
        ma = month_num(table.cell(r, 1).text)
        mb = month_num(table.cell(r, 6).text)
        mc = month_num(table.cell(r, 11).text)

        slots = [
            (ma, "A1", 2, 3),
            (ma, "A2", 4, 5),
            (mb, "B1", 7, 8),
            (mb, "B2", 9, 10),
            (mc, "C1", 12, 13),
        ]
        for mm, a, nc, gc in slots:
            if mm is None:
                continue
            dct = data.get((mm, a))
            if not dct:
                continue
            if set_cell(table, r, nc, dct.get("n")):
                changes += 1
            if set_cell(table, r, gc, dct.get("gmc")):
                changes += 1

        if mb == 4 and p_m4:
            if set_cell(table, r, 14, p_m4):
                changes += 1
            continue

        if mc and data.get((mc, "A2")):
            dct_a2 = data[(mc, "A2")]
            for col, key in ((14, "p_a2b2"), (15, "p_a2c1"), (16, "p_b2c1")):
                if set_cell(table, r, col, dct_a2.get(key)):
                    changes += 1
    return changes


def update_slide4_table1(table, cmap: dict[tuple[str, int], dict[str, str]]) -> int:
    changes = 0
    for r in range(3, len(table.rows)):
        label = table.cell(r, 0).text.strip()
        grp = "FAS" if label.startswith("FAS") or (label == "" and r >= 9) else "PPS"
        mm = month_num(table.cell(r, 1).text)
        if mm is None:
            continue
        dct = cmap.get((grp, mm))
        if not dct:
            continue
        for c, key in ((3, "c3_n"), (4, "c3_gmc"), (5, "c2_n"), (6, "c2_gmc"), (7, "p")):
            if set_cell(table, r, c, dct.get(key)):
                changes += 1
    return changes


def sync(
    word_path: Path,
    ppt_path: Path,
    out_path: Path,
    *,
    pps_table: int,
    fas_table: int,
    c3c2_table: int,
    pps_m4_supp_table: int,
    fas_m4_p_supp_table: int,
) -> int:
    doc = Document(str(word_path))
    prs = Presentation(str(ppt_path))

    total = 0
    pps = build_month_arm_map(table_rows(doc, pps_table - 1))
    fas = build_month_arm_map(table_rows(doc, fas_table - 1))

    pps_m4_updates, pps_m4_p = supplement_pps_month4_b_groups(doc, pps_m4_supp_table - 1)
    pps.update(pps_m4_updates)
    fas_m4_p = supplement_fas_m4_p(doc, fas_m4_p_supp_table - 1)
    cmap = build_c3c2_map(doc, c3c2_table - 1)

    s1_tables = [sh.table for sh in prs.slides[0].shapes if sh.has_table]
    if len(s1_tables) >= 2:
        normalize_gmc_headers(s1_tables[0])
        normalize_gmc_headers(s1_tables[1])
        total += update_main_slide_table(s1_tables[0], pps)
        total += update_summary_table(s1_tables[1], pps, pps_m4_p)

    s2_tables = [sh.table for sh in prs.slides[1].shapes if sh.has_table]
    if len(s2_tables) >= 2:
        normalize_gmc_headers(s2_tables[0])
        normalize_gmc_headers(s2_tables[1])
        total += update_main_slide_table(s2_tables[0], fas)
        total += update_summary_table(s2_tables[1], fas, fas_m4_p)

    s4_tables = [sh.table for sh in prs.slides[3].shapes if sh.has_table]
    if s4_tables:
        normalize_gmc_headers(s4_tables[0])
        total += update_slide4_table1(s4_tables[0], cmap)

    prs.save(str(out_path))
    return total


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--word", type=Path, required=True)
    ap.add_argument("--ppt", type=Path, required=True)
    ap.add_argument("--out", type=Path, default=None)

    ap.add_argument("--pps-table", type=int, default=8)
    ap.add_argument("--fas-table", type=int, default=9)
    ap.add_argument("--c3c2-table", type=int, default=13)
    ap.add_argument("--pps-m4-supp-table", type=int, default=7)
    ap.add_argument("--fas-m4-p-supp-table", type=int, default=4)

    args = ap.parse_args()
    out = args.out or args.ppt.with_name(f"{args.ppt.stem}-按Word更新.pptx")
    n = sync(
        args.word,
        args.ppt,
        out,
        pps_table=args.pps_table,
        fas_table=args.fas_table,
        c3c2_table=args.c3c2_table,
        pps_m4_supp_table=args.pps_m4_supp_table,
        fas_m4_p_supp_table=args.fas_m4_p_supp_table,
    )
    print(f"Wrote {out} ({n} cell updates)")


if __name__ == "__main__":
    main()

