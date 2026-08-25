"""
Rebuild baseline slide (第13页) with:
1. Two-tier headers matching reference format (0,1月程序 / 低剂量组(A1) etc.)
2. Removed: 检验方法, P值, Mean(SD), Min/Max
3. Renamed: GM → GMC
4. Fixed layout: no overlapping tables
"""

import copy
import glob
import os

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ======== Configuration ========
FONT_NAME = "微软雅黑"
COLOR_HEADER = RGBColor(0xC0, 0x00, 0x00)
COLOR_DATA = RGBColor(0x00, 0x00, 0x00)

# Heights (EMU integers)
HEADER_H = 170000
DATA_H = 135000

# Horizontal positions (EMU integers)
LEFT_X = 80000
RIGHT_X = 4980000

# Column widths (EMU integers)
L_COL0 = 800000  # 指标/既往接种
L_COL_GRP = 780000  # each group column (left tables)
R_COL0 = 800000  # 指标/既往接种
R_COL_GRP = 980000  # each group column (right tables)


# ======== Helper Functions ========


def set_cell(
    cell, text, size=Pt(7), bold=False, color=COLOR_DATA, align=PP_ALIGN.CENTER
):
    """Set cell text with formatting (handles multi-line via \\n)."""
    cell.text = str(text)
    for para in cell.text_frame.paragraphs:
        para.alignment = align
        para.space_before = Pt(0)
        para.space_after = Pt(0)
        for run in para.runs:
            run.font.size = size
            run.font.name = FONT_NAME
            run.font.bold = bold
            run.font.color.rgb = color
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn("a:ea"))
            if ea is None:
                ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", FONT_NAME)
    cell.margin_top = Emu(10000)
    cell.margin_bottom = Emu(10000)
    cell.margin_left = Emu(15000)
    cell.margin_right = Emu(15000)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_border(cell, color="000000", width="6350"):
    """Set thin black borders on all sides."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for bn in ["lnL", "lnR", "lnT", "lnB"]:
        old = tcPr.find(qn(f"a:{bn}"))
        if old is not None:
            tcPr.remove(old)
        el = etree.SubElement(tcPr, qn(f"a:{bn}"))
        el.set("w", width)
        el.set("cap", "flat")
        el.set("cmpd", "sng")
        sf = etree.SubElement(el, qn("a:solidFill"))
        sc = etree.SubElement(sf, qn("a:srgbClr"))
        sc.set("val", color)
        etree.SubElement(el, qn("a:prstDash")).set("val", "solid")


def add_label(slide, text, x, y, w, h=165000, size=Pt(8), color=COLOR_HEADER):
    """Add a section label text box."""
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.size = size
        run.font.name = FONT_NAME
        run.font.bold = True
        run.font.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", FONT_NAME)


def make_table(slide, nrows, ncols, x, y, col_widths, header_rows=2):
    """Create a table and return the pptx table object."""
    total_w = sum(col_widths)
    total_h = HEADER_H * header_rows + DATA_H * (nrows - header_rows)
    shape = slide.shapes.add_table(
        nrows, ncols, Emu(x), Emu(y), Emu(total_w), Emu(total_h)
    )
    tbl = shape.table
    # Clear default table style
    for attr in list(tbl._tbl.tblPr.attrib.keys()):
        del tbl._tbl.tblPr.attrib[attr]
    # Set column widths
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Emu(w)
    # Set row heights
    for ri, row in enumerate(tbl.rows):
        row.height = Emu(HEADER_H if ri < header_rows else DATA_H)
    return tbl


def borders_all(tbl, nrows, ncols):
    """Apply borders to every cell (safe with merged cells)."""
    done = set()
    for ri in range(nrows):
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            cid = id(cell._tc)
            if cid not in done:
                set_border(cell)
                done.add(cid)


# ---- Left table header (18-59, 6 cols, 2-tier) ----
def header_left(tbl, col0_label, n_vals=None):
    """
    Row 0: col0 [merge r0-r1] | 0,1月程序 [merge c1-c2] | 0,2月程序 [merge c3-c4] | 0,1,6月程序
    Row 1:                    | 低剂量组(A1) N=xx | 高剂量组(A2) N=xx | ... | 阳性对照组(C1) N=xx
    """
    n = n_vals or [150, 150, 150, 150, 150]
    # Col 0: merge rows 0-1
    tbl.cell(0, 0).merge(tbl.cell(1, 0))
    set_cell(tbl.cell(0, 0), col0_label, bold=True, color=COLOR_HEADER)
    # 0,1月程序 spanning cols 1-2
    tbl.cell(0, 1).merge(tbl.cell(0, 2))
    set_cell(tbl.cell(0, 1), "0,1月程序", bold=True, color=COLOR_HEADER)
    # 0,2月程序 spanning cols 3-4
    tbl.cell(0, 3).merge(tbl.cell(0, 4))
    set_cell(tbl.cell(0, 3), "0,2月程序", bold=True, color=COLOR_HEADER)
    # 0,1,6月程序 col 5
    set_cell(tbl.cell(0, 5), "0,1,6月程序", bold=True, color=COLOR_HEADER)
    # Row 1: group names with N
    groups = [
        f"低剂量组(A1)\nN={n[0]}",
        f"高剂量组(A2)\nN={n[1]}",
        f"低剂量组(B1)\nN={n[2]}",
        f"高剂量组(B2)\nN={n[3]}",
        f"阳性对照组(C1)\nN={n[4]}",
    ]
    for ci, g in enumerate(groups):
        set_cell(tbl.cell(1, ci + 1), g, bold=True, color=COLOR_DATA)


# ---- Right table header (60+, 3 cols, 2-tier) ----
def header_right(tbl, col0_label, n_c3=75, n_c2=75):
    """
    Row 0: col0 [merge r0-r1] | 0,1,6月程序 [merge c1-c2]
    Row 1:                    | 高剂量组(C3) N=xx | 阳性对照组(C2) N=xx
    """
    tbl.cell(0, 0).merge(tbl.cell(1, 0))
    set_cell(tbl.cell(0, 0), col0_label, bold=True, color=COLOR_HEADER)
    tbl.cell(0, 1).merge(tbl.cell(0, 2))
    set_cell(tbl.cell(0, 1), "0,1,6月程序", bold=True, color=COLOR_HEADER)
    set_cell(tbl.cell(1, 1), f"高剂量组(C3)\nN={n_c3}", bold=True, color=COLOR_DATA)
    set_cell(tbl.cell(1, 2), f"阳性对照组(C2)\nN={n_c2}", bold=True, color=COLOR_DATA)


def fill_data_rows(tbl, data_list):
    """Fill data rows. data_list = [(row_idx, [col_values]), ...]"""
    for ri, row_data in data_list:
        for ci, val in enumerate(row_data):
            a = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            set_cell(tbl.cell(ri, ci), val, align=a)


def add_baseline_subheaders(tbl, ncols, rate_row=2, gmc_row=5):
    """Add merged sub-header rows for 阳性率 and GMC sections."""
    tbl.cell(rate_row, 0).merge(tbl.cell(rate_row, ncols - 1))
    set_cell(
        tbl.cell(rate_row, 0),
        "免前抗-HBs阳性率",
        bold=True,
        color=COLOR_HEADER,
        align=PP_ALIGN.LEFT,
    )
    tbl.cell(gmc_row, 0).merge(tbl.cell(gmc_row, ncols - 1))
    set_cell(
        tbl.cell(gmc_row, 0),
        "免前抗-HBs GMC(mIU/ml)",
        bold=True,
        color=COLOR_HEADER,
        align=PP_ALIGN.LEFT,
    )


# ======== Main Script ========

files = glob.glob(
    r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\*.pptx"
)
pptx_path = [f for f in files if not os.path.basename(f).startswith("~")][0]
print(f"Opening: {pptx_path}")
prs = Presentation(pptx_path)
print(f"Total slides: {len(prs.slides)}")

# ---- Verify and clear slide 13 (index 12) ----
slide = prs.slides[12]
spTree = slide.shapes._spTree
removed = 0
for child in list(spTree):
    tag = child.tag
    if any(
        tag.endswith(t) for t in ["}sp", "}graphicFrame", "}grpSp", "}pic", "}cxnSp"]
    ):
        spTree.remove(child)
        removed += 1
print(f"Cleared {removed} shapes from slide 13")

# ---- Copy title bar + line from slide 14 (index 13) ----
ref_slide = prs.slides[13]
for shape in ref_slide.shapes:
    if shape.name in ("Rectangle 2", "直接连接符 2"):
        spTree.append(copy.deepcopy(shape._element))

# Set title text
for shape in slide.shapes:
    if shape.name == "Rectangle 2" and shape.has_text_frame:
        runs = list(shape.text_frame.paragraphs[0].runs)
        if runs:
            runs[0].text = "2期基线情况"
            for r in runs[1:]:
                r.text = ""
        break

# ================================================================
# Vertical position calculations
# ================================================================
Y1_LABEL = 590000
Y1_SUB = 755000
Y1_TBL = 890000
T1_H = HEADER_H * 2 + DATA_H * 3  # 745000

Y2_LABEL = Y1_TBL + T1_H + 60000  # 1695000
Y2_SUB = Y2_LABEL + 165000  # 1860000
Y2_TBL = Y2_SUB + 130000  # 1990000
T2_H = HEADER_H * 2 + DATA_H * 6  # 1150000

Y3_LABEL = Y2_TBL + T2_H + 60000  # 3200000
Y3_TBL = Y3_LABEL + 165000  # 3365000
# T3 ends at: 3365000 + 1150000 = 4515000 < 5143500 ✓

print(f"Layout: Section1 table ends at {Y1_TBL + T1_H}")
print(f"Layout: Section2 table ends at {Y2_TBL + T2_H}")
print(f"Layout: Section3 table ends at {Y3_TBL + T2_H}")
print(f"Slide height: {prs.slide_height}")

# Column widths
LW = [L_COL0] + [L_COL_GRP] * 5  # [800000, 780000 x5] = 4700000
RW = [R_COL0, R_COL_GRP, R_COL_GRP]  # [800000, 980000, 980000] = 2760000

# ================================================================
# SECTION 1: 既往有无接种乙型肝炎疫苗
# ================================================================
add_label(slide, "一、既往有无接种乙型肝炎疫苗", LEFT_X, Y1_LABEL, 4500000)
add_label(
    slide,
    "18-59岁人群(FAS)",
    LEFT_X,
    Y1_SUB,
    2000000,
    h=130000,
    size=Pt(7),
    color=COLOR_DATA,
)
add_label(
    slide,
    "60岁及以上人群(FAS)",
    RIGHT_X,
    Y1_SUB,
    2000000,
    h=130000,
    size=Pt(7),
    color=COLOR_DATA,
)

# Table 1a: 18-59 vaccination (5 rows x 6 cols)
t1a = make_table(slide, 5, 6, LEFT_X, Y1_TBL, LW)
header_left(t1a, "既往接种")
fill_data_rows(
    t1a,
    [
        (2, ["有 n(%)", "13(8.67)", "13(8.67)", "12(8.00)", "11(7.33)", "13(8.67)"]),
        (
            3,
            [
                "无 n(%)",
                "111(74.00)",
                "110(73.33)",
                "121(80.67)",
                "113(75.33)",
                "109(72.67)",
            ],
        ),
        (
            4,
            [
                "不详 n(%)",
                "26(17.33)",
                "27(18.00)",
                "17(11.33)",
                "26(17.33)",
                "28(18.67)",
            ],
        ),
    ],
)
borders_all(t1a, 5, 6)

# Table 1b: 60+ vaccination (5 rows x 3 cols)
t1b = make_table(slide, 5, 3, RIGHT_X, Y1_TBL, RW)
header_right(t1b, "既往接种")
fill_data_rows(
    t1b,
    [
        (2, ["有 n(%)", "3(4.00)", "1(1.33)"]),
        (3, ["无 n(%)", "67(89.33)", "70(93.33)"]),
        (4, ["不详 n(%)", "5(6.67)", "4(5.33)"]),
    ],
)
borders_all(t1b, 5, 3)

# ================================================================
# SECTION 2: 免前基线抗-HBs阳性率与GMC
# ================================================================
add_label(slide, "二、免前基线抗-HBs阳性率与GMC", LEFT_X, Y2_LABEL, 5000000)
add_label(
    slide,
    "18-59岁人群(FAS)",
    LEFT_X,
    Y2_SUB,
    2000000,
    h=130000,
    size=Pt(7),
    color=COLOR_DATA,
)
add_label(
    slide,
    "60岁及以上人群(FAS)",
    RIGHT_X,
    Y2_SUB,
    2000000,
    h=130000,
    size=Pt(7),
    color=COLOR_DATA,
)

# Table 2a: 18-59 baseline (8 rows x 6 cols)
# Rows: [header0, header1, 阳性率subhdr, 阳性n%, 95%CI, GMCsubhdr, GMC, Median]
t2a = make_table(slide, 8, 6, LEFT_X, Y2_TBL, LW)
header_left(t2a, "指标")
add_baseline_subheaders(t2a, 6)
fill_data_rows(
    t2a,
    [
        (3, ["阳性 n(%)", "0(0.00)", "1(0.67)", "0(0.00)", "1(0.67)", "0(0.00)"]),
        (
            4,
            [
                "95%CI(%)",
                "0.00, 2.43",
                "0.02, 3.66",
                "0.00, 2.43",
                "0.02, 3.66",
                "0.00, 2.43",
            ],
        ),
        (
            6,
            [
                "GMC(95%CI)",
                "1.18(1.11,1.26)",
                "1.26(1.17,1.35)",
                "1.12(1.06,1.18)",
                "1.18(1.08,1.29)",
                "1.16(1.09,1.23)",
            ],
        ),
        (7, ["Median", "1.000", "1.000", "1.000", "1.000", "1.000"]),
    ],
)
borders_all(t2a, 8, 6)

# Table 2b: 60+ baseline (8 rows x 3 cols)
t2b = make_table(slide, 8, 3, RIGHT_X, Y2_TBL, RW)
header_right(t2b, "指标")
add_baseline_subheaders(t2b, 3)
fill_data_rows(
    t2b,
    [
        (3, ["阳性 n(%)", "0(0.00)", "0(0.00)"]),
        (4, ["95%CI(%)", "0.00, 4.80", "0.00, 4.80"]),
        (6, ["GMC(95%CI)", "1.15(1.07,1.25)", "1.20(1.09,1.32)"]),
        (7, ["Median", "1.000", "1.000"]),
    ],
)
borders_all(t2b, 8, 3)

# ================================================================
# SECTION 3: 有乙肝疫苗接种史的18-59岁人群
# ================================================================
add_label(
    slide,
    "三、有乙肝疫苗接种史的18-59岁人群免前基线抗-HBs阳性率与GMC(FAS)",
    LEFT_X,
    Y3_LABEL,
    8900000,
)

# Table 3: vaccinated 18-59 (8 rows x 6 cols)
t3 = make_table(slide, 8, 6, LEFT_X, Y3_TBL, LW)
header_left(t3, "指标", n_vals=[13, 13, 12, 11, 13])
add_baseline_subheaders(t3, 6)
fill_data_rows(
    t3,
    [
        (3, ["阳性 n(%)", "0(0.00)", "0(0.00)", "0(0.00)", "1(9.09)", "0(0.00)"]),
        (
            4,
            [
                "95%CI(%)",
                "0.00, 24.71",
                "0.00, 24.71",
                "0.00, 26.46",
                "0.23, 41.28",
                "0.00, 24.71",
            ],
        ),
        (
            6,
            [
                "GMC(95%CI)",
                "1.23(0.90,1.67)",
                "1.12(0.95,1.32)",
                "1.13(0.86,1.50)",
                "1.97(0.69,5.66)",
                "1.15(0.94,1.41)",
            ],
        ),
        (7, ["Median", "1.000", "1.000", "1.000", "1.000", "1.000"]),
    ],
)
borders_all(t3, 8, 6)

# ================================================================
# Save
# ================================================================
# Try saving directly, if locked save to temp file
try:
    prs.save(pptx_path)
    print(f"\nSaved: {pptx_path}")
except PermissionError:
    # File is locked (probably open in PowerPoint)
    temp_path = pptx_path.replace(".pptx", "_updated.pptx")
    prs.save(temp_path)
    print(f"\nOriginal file is locked. Saved to: {temp_path}")
    print("Please close the original PPT, then rename the _updated file.")
print(f"Total slides: {len(prs.slides)}")
print("Done! Slide 13 rebuilt successfully.")
