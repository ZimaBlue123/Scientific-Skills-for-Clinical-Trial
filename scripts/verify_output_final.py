"""V29 生成结果自检: 页数/分页/🏦图标/Study1新表述/PMID标注/ACTRN.aspx/新闻稿链接/动脉网移除/标题字体
用法: python verify_v29_output.py
"""

import re

from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx import Presentation

BASE = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial"
PPTX = BASE + r"\review_materials\CpG_Vaccine_Safety_Summary-V29-20260827.pptx"
DOCX = BASE + r"\review_materials\CpG_Vaccine_Safety_Summary-V29-20260827.docx"

CHICTR_EXPECT = {
    "ChiCTR2600119810": "https://www.chictr.org.cn/showproj.html?proj=311880",
    "ChiCTR2600118487": "https://www.chictr.org.cn/showproj.html?proj=301189",
    "ChiCTR2500108408": "https://www.chictr.org.cn/showproj.html?proj=280568",
}
ACTRN_EXPECT = "https://anzctr.org.au/ACTRN12624000064505.aspx"  # V29: 修正 .aspx 后缀
NEWS_EXPECT = "https://www.eatg.org/?p=20095"  # V29: Uvax Bio 中期分析#1 新闻稿(EATG全文)

issues = []

# ---------------- PPTX ----------------
prs = Presentation(PPTX)
slides = list(prs.slides)
print(f"[PPTX] 总页数: {len(slides)} (期望 10)")
if len(slides) != 10:
    issues.append(f"PPTX 页数 {len(slides)} != 10")


def slide_tables(slide):
    return [shp for shp in slide.shapes if shp.has_table]


# 第 4 页 (index 3): 应仅 3 行 HEPLISAV-B
tbl4 = slide_tables(slides[3])[0].table
n4 = len(tbl4.rows)
print(f"[PPTX] 第4页 表格行数: {n4} (期望 4 = 表头+3行数据)")
if n4 != 4:
    issues.append(f"第4页表格 {n4} 行 != 4")
col0_4 = [tbl4.cell(r, 0).text.split("\n")[0] for r in range(1, n4)]
for r, t in enumerate(col0_4, 1):
    print(f"  - 行{r}: {t}")
bad4 = [t for t in col0_4 if "HEPLISAV" not in t]
if bad4:
    issues.append(f"第4页含非 HEPLISAV 行: {bad4}")

# 第 5 页 (index 4): 首行应为 CYFENDUS, 共 5 行数据
tbl5 = slide_tables(slides[4])[0].table
n5 = len(tbl5.rows)
print(f"[PPTX] 第5页 表格行数: {n5} (期望 6 = 表头+5行数据)")
if n5 != 6:
    issues.append(f"第5页表格 {n5} 行 != 6")
first5 = tbl5.cell(1, 0).text.split("\n")[0]
print(f"  - 第5页首行: {first5}")
if "CYFENDUS" not in first5:
    issues.append(f"第5页首行不是 CYFENDUS: {first5}")

# ---- V29 自检: 详表第一列图标 (倒数第2行 🏦 / 最后一行 🎯) ----
print("\n[PPTX] 详表第一列图标检查 (第4-9页):")
for si in range(3, 9):
    for shp in slide_tables(slides[si]):
        tbl = shp.table
        for r in range(1, len(tbl.rows)):
            lines = tbl.cell(r, 0).text.split("\n")
            sp = lines[-2] if len(lines) >= 2 else ""
            ind = lines[-1] if len(lines) >= 1 else ""
            ok = sp.startswith("🏦") and ind.startswith("🎯")
            if not ok:
                issues.append(
                    f"第{si + 1}页 行{r} 第一列图标缺失: 申办者行={sp[:14]!r} 适应症行={ind[:14]!r}"
                )
                print(f"  ✗ 第{si + 1}页 行{r}: {sp[:20]!r} / {ind[:20]!r}")
    print(f"  - 第{si + 1}页: 图标检查完成")

# ---- V29 自检: Study 1 新表述 (去累赘) + PMID 标注 + 怡道/Uvax 行 ----
print("\n[PPTX] Study 1 新表述 / PMID / 怡道参考列 / Uvax 标注检查:")
study1_ok = False
pmid_ok = False
yidao_ref_ok = False
uvax_tag_ok = False
for si, slide in enumerate(slides):
    for shp in slide.shapes:
        if shp.has_table:
            tbl = shp.table
            if len(tbl.columns) < 5:
                continue  # 概览表仅4列
            for r in range(1, len(tbl.rows)):
                sec_text = tbl.cell(r, 3).text
                if "汇总发生率23%-39%" in sec_text:
                    # V29 新表述: "HEPLISAV-B第1/2剂后38.5%/34.8% vs 对照Engerix-B第1/2/3剂后33.6%/24.7%/20.2%"
                    if (
                        "HEPLISAV-B第1/2剂后38.5%/34.8%" in sec_text
                        and "对照Engerix-B第1/2/3剂后33.6%/24.7%/20.2%" in sec_text
                        and "HEPLISAV-B第1/2剂后17.4%/13.8%" in sec_text
                        and "HEPLISAV-B第1/2剂后16.9%/12.8%" in sec_text
                    ):
                        study1_ok = True
                    else:
                        issues.append(f"第{si + 1}页 Study 1 新表述未落位")
                if "源自中期分析#1" in sec_text and "2024-11-19" in sec_text:
                    uvax_tag_ok = True
                ref_text = tbl.cell(r, 4).text
                if "PMID: 23727002" in ref_text:
                    if "Study 2/HBV-16" in ref_text and "HBV-13" not in ref_text:
                        pmid_ok = True
                    else:
                        issues.append(f"第{si + 1}页 PMID 23727002 标注未修正: {ref_text[-80:]}")
                if "NDA受理号: CXSS2500108" in ref_text:
                    yidao_ref_ok = True
print(f"  - Study 1 新表述(去等号累赘): {'✓ 已落位' if study1_ok else '✗ 缺失'}")
print(f"  - PMID 23727002 标注: {'✓ 保留 Study 2/HBV-16' if pmid_ok else '✗ 未修正'}")
print(f"  - 怡道参考列 NDA受理号: {'✓ CXSS2500108' if yidao_ref_ok else '✗ 缺失'}")
print(f"  - Uvax 中期分析#1 标注: {'✓ 已落位' if uvax_tag_ok else '✗ 缺失'}")
if not study1_ok:
    issues.append("Study 1 新表述缺失")
if not pmid_ok:
    issues.append("PMID 23727002 标注未修正")
if not yidao_ref_ok:
    issues.append("怡道参考列 NDA受理号 缺失")
if not uvax_tag_ok:
    issues.append("Uvax 中期分析#1 标注缺失")

# ---- V29 自检: ACTRN 链接 (.aspx) + Uvax 新闻稿链接 + 动脉网已移除 ----
print("\n[PPTX] ACTRN / 新闻稿链接 / 动脉网移除检查:")
acctrn_url = None
news_url = None
all_text = []
for si, slide in enumerate(slides):
    for shp in slide.shapes:
        if shp.has_table:
            for r in shp.table.rows:
                for c in r.cells:
                    for p in c.text_frame.paragraphs:
                        line = p.text
                        if line:
                            all_text.append(line)
                        for run in p.runs:
                            t = run.text
                            url = (
                                run.hyperlink.address
                                if run.hyperlink and run.hyperlink.address
                                else None
                            )
                            if "ACTRN12624000064505" in t:
                                acctrn_url = url
                            if "Uvax Bio新闻稿" in t:
                                news_url = url
        elif shp.has_text_frame:
            for p in shp.text_frame.paragraphs:
                if p.text:
                    all_text.append(p.text)
                for run in p.runs:
                    if "ACTRN12624000064505" in (run.text or ""):
                        acctrn_url = (
                            run.hyperlink.address
                            if run.hyperlink and run.hyperlink.address
                            else None
                        )
                    if "Uvax Bio新闻稿" in (run.text or ""):
                        news_url = (
                            run.hyperlink.address
                            if run.hyperlink and run.hyperlink.address
                            else None
                        )
full_text = "\n".join(all_text)
print(f"  - ACTRN 链接: {acctrn_url}")
print(f"  - Uvax 新闻稿链接: {news_url}")
print(f"  - 输出含'动脉网': {'是!' if '动脉网' in full_text else '否 ✓'}")
if acctrn_url != ACTRN_EXPECT:
    issues.append(f"ACTRN 链接异常: 期望 {ACTRN_EXPECT}, 实际 {acctrn_url}")
if news_url != NEWS_EXPECT:
    issues.append(f"Uvax 新闻稿链接异常: 期望 {NEWS_EXPECT}, 实际 {news_url}")
if "动脉网" in full_text:
    issues.append("PPTX 输出仍含'动脉网'")

# 全 PPTX 扫描 ChiCTR 超链接
ppt_chi_text = []
for si, slide in enumerate(slides):
    for shp in slide.shapes:
        if shp.has_table:
            for r in shp.table.rows:
                for c in r.cells:
                    for p in c.text_frame.paragraphs:
                        for run in p.runs:
                            t = run.text
                            if "ChiCTR" in t:
                                m = re.search(r"(ChiCTR\d+)", t)
                                if m:
                                    url = (
                                        run.hyperlink.address
                                        if run.hyperlink and run.hyperlink.address
                                        else None
                                    )
                                    ppt_chi_text.append((si + 1, m.group(1), url))
        elif shp.has_text_frame:
            for p in shp.text_frame.paragraphs:
                for run in p.runs:
                    if "ChiCTR" in (run.text or ""):
                        m = re.search(r"(ChiCTR\d+)", run.text)
                        if m:
                            url = (
                                run.hyperlink.address
                                if run.hyperlink and run.hyperlink.address
                                else None
                            )
                            ppt_chi_text.append((si + 1, m.group(1), url))
print("\n[PPTX] ChiCTR 文本/链接扫描:")
for page, code, url in ppt_chi_text:
    exp = CHICTR_EXPECT.get(code)
    status = "OK" if url == exp else f"异常(期望 {exp}, 实际 {url})"
    print(f"  - 第{page}页 {code}: {status}")
    if url != exp:
        issues.append(f"PPTX 第{page}页 {code} 链接异常: {url}")

# 标题字体样式检查
print("\n[PPTX] 每页大标题:")
for si, slide in enumerate(slides):
    for shp in slide.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip():
            txt = shp.text_frame.text.strip().split("\n")[0]
            runs = shp.text_frame.paragraphs[0].runs
            if runs:
                f = runs[0].font
                sz = f.size.pt if f.size else "?"
                col = f.color.rgb if f.color and f.color.type else "?"
                print(f"  - 第{si + 1}页: 「{txt[:22]}」 {sz}pt 加粗{bool(f.bold)} 颜色{col}")
                if sz != 20 or not f.bold or str(col) != "C00000":
                    issues.append(f"第{si + 1}页标题样式异常: {sz}pt/{f.bold}/{col}")
            break

# ---------------- DOCX ----------------
print("\n[DOCX] 超链接扫描:")
doc = Document(DOCX)
docx_rels = [rel.target_ref for rel in doc.part.rels.values() if rel.reltype == RT.HYPERLINK]
print(f"  - 超链接总数: {len(docx_rels)}")
for code, exp in CHICTR_EXPECT.items():
    hit = exp in docx_rels
    print(f"  - {code} -> {exp}: {'已嵌入' if hit else '缺失!'}")
    if not hit:
        issues.append(f"DOCX 缺少 {code} 链接 {exp}")
print(f"  - ACTRN(.aspx) -> {ACTRN_EXPECT}: {'已嵌入' if ACTRN_EXPECT in docx_rels else '缺失!'}")
print(f"  - Uvax新闻稿 -> {NEWS_EXPECT}: {'已嵌入' if NEWS_EXPECT in docx_rels else '缺失!'}")
if ACTRN_EXPECT not in docx_rels:
    issues.append("DOCX 缺少 ACTRN .aspx 链接")
if NEWS_EXPECT not in docx_rels:
    issues.append("DOCX 缺少 Uvax 新闻稿链接")

# DOCX 全文不含"动脉网"
doc_full = "\n".join(p.text for p in doc.paragraphs)
for t in doc.tables:
    for row in t.rows:
        for c in row.cells:
            doc_full += "\n" + c.text
print(f"[DOCX] 全文含'动脉网': {'是!' if '动脉网' in doc_full else '否 ✓'}")
if "动脉网" in doc_full:
    issues.append("DOCX 输出仍含'动脉网'")

# DOCX 图标 + Study1 + PMID + 怡道/Uvax 检查 (表3)
t3 = doc.tables[2]
print("\n[DOCX] 详表图标 / Study 1 / PMID / 怡道 / Uvax 检查:")
dx_icon_ok = True
dx_study1 = False
dx_pmid = False
dx_yidao = False
dx_uvax = False
for r in range(1, len(t3.rows)):
    lines = t3.rows[r].cells[0].text.split("\n")
    if len(lines) >= 2:
        if not (lines[-2].startswith("🏦") and lines[-1].startswith("🎯")):
            dx_icon_ok = False
            issues.append(f"DOCX 详表 行{r + 1} 图标缺失: {lines[-2][:14]!r}/{lines[-1][:14]!r}")
    sec = t3.rows[r].cells[3].text
    if (
        "汇总发生率23%-39%" in sec
        and "HEPLISAV-B第1/2剂后38.5%/34.8%" in sec
        and "对照Engerix-B第1/2/3剂后33.6%/24.7%/20.2%" in sec
    ):
        dx_study1 = True
    if "源自中期分析#1" in sec and "2024-11-19" in sec:
        dx_uvax = True
    ref = t3.rows[r].cells[4].text
    if "PMID: 23727002" in ref:
        dx_pmid = "Study 2/HBV-16" in ref and "HBV-13" not in ref
        if not dx_pmid:
            issues.append(f"DOCX PMID 23727002 标注未修正: {ref[-80:]}")
    if "NDA受理号: CXSS2500108" in ref:
        dx_yidao = True
print(f"  - 图标: {'✓ 24行全部带图标' if dx_icon_ok else '✗ 存在缺失'}")
print(f"  - Study 1 新表述: {'✓' if dx_study1 else '✗'}")
print(f"  - PMID 标注: {'✓' if dx_pmid else '✗'}")
print(f"  - 怡道 NDA受理号: {'✓' if dx_yidao else '✗'}")
print(f"  - Uvax 中期分析标注: {'✓' if dx_uvax else '✗'}")
if not dx_icon_ok:
    issues.append("DOCX 详表图标缺失")
if not dx_study1:
    issues.append("DOCX Study 1 新表述缺失")
if not dx_pmid:
    issues.append("DOCX PMID 标注未修正")
if not dx_yidao:
    issues.append("DOCX 怡道 NDA受理号 缺失")
if not dx_uvax:
    issues.append("DOCX Uvax 中期分析标注缺失")

print("\n[汇总]")
if issues:
    print("发现问题:")
    for i in issues:
        print(f"  ✗ {i}")
else:
    print("  ✓ 全部自检项通过")
