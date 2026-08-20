"""Script to convert the specific Word summary table into a 16:9 Dark Red PPT with hyperlinks, simplified headers, and no title underline."""

import re
import sys

try:
    import pptx
except ImportError:
    import subprocess

    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

try:
    import docx
except ImportError:
    import subprocess

    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-docx"])
    import docx

import pptx.enum.shapes
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def create_ppt_from_docx(doc_path, ppt_path):
    print(f"Reading document: {doc_path}")
    doc = docx.Document(doc_path)

    table = doc.tables[0]
    headers = [c.text for c in table.rows[0].cells]

    # Header Simplification
    for i in range(len(headers)):
        h = headers[i].replace("\n", " ")
        if "临床试验基本信息" in h:
            headers[i] = "临床试验基本信息"
        elif "核心参考文献" in h:
            headers[i] = "核心参考文献"
        else:
            headers[i] = h

    row_data = []
    for row in table.rows[1:]:
        row_data.append([c.text for c in row.cells])

    print(f"Extracted {len(row_data)} rows of data.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_RED = RGBColor(192, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    LIGHT_GRAY = RGBColor(245, 245, 245)
    HEADER_TEXT = RGBColor(255, 255, 255)
    LINK_BLUE = RGBColor(5, 99, 193)

    blank_layout = prs.slide_layouts[6]

    def format_cell(
        cell, text, size=8.5, bold=False, text_color=BLACK, bg_color=None, is_ref=False
    ):
        # Clear default empty paragraph
        cell.text = ""

        # Aggressively reduce margins to maximize text area
        cell.margin_top = Pt(2)
        cell.margin_bottom = Pt(2)
        cell.margin_left = Pt(5)
        cell.margin_right = Pt(5)

        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

        p = cell.text_frame.paragraphs[0]
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        p.line_spacing = 1.0

        clean_text = text.replace("**", "")

        if not is_ref:
            run = p.add_run()
            run.text = clean_text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = text_color
            run.font.name = "Microsoft YaHei"
        else:
            # Parse links for Reference Column
            pattern = re.compile(r"(PMID:\s*)(\d+)|(DOI:\s*)(10\.\S+)")
            last_idx = 0
            for match in pattern.finditer(clean_text):
                # Text before match
                if match.start() > last_idx:
                    run = p.add_run()
                    run.text = clean_text[last_idx : match.start()]
                    run.font.size = Pt(size)
                    run.font.bold = bold
                    run.font.color.rgb = text_color
                    run.font.name = "Microsoft YaHei"

                # The match (Hyperlink)
                run = p.add_run()
                if match.group(1):  # PMID
                    run.text = match.group(1) + match.group(2)
                    run.hyperlink.address = (
                        f"https://pubmed.ncbi.nlm.nih.gov/{match.group(2)}/"
                    )
                elif match.group(3):  # DOI
                    run.text = match.group(3) + match.group(4)
                    run.hyperlink.address = f"https://doi.org/{match.group(4)}"

                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = LINK_BLUE  # Standard link blue
                run.font.name = "Microsoft YaHei"

                last_idx = match.end()

            # Text after last match
            if last_idx < len(clean_text):
                run = p.add_run()
                run.text = clean_text[last_idx:]
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = text_color
                run.font.name = "Microsoft YaHei"

    # 1. Title Slide
    title_slide = prs.slides.add_slide(blank_layout)
    rect = title_slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.5),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK_RED
    rect.line.color.rgb = DARK_RED

    title_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(2)
    )
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "CpG佐剂预防性疫苗\n核心临床试验与安全性数据汇总"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(44)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = DARK_RED
    p.runs[0].font.name = "Microsoft YaHei"

    # 2. Content Slides
    col_widths = [Inches(1.8), Inches(1.3), Inches(3.2), Inches(5.6), Inches(1.3)]
    left_margin = Inches(0.1)
    top_margin = Inches(0.9)  # Maximize space

    chunk_size = 2
    total_chunks = (len(row_data) + chunk_size - 1) // chunk_size

    for slide_idx in range(total_chunks):
        chunk = row_data[slide_idx * chunk_size : (slide_idx + 1) * chunk_size]

        slide = prs.slides.add_slide(blank_layout)

        # Slide Title (Without underline)
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.05), Inches(12.33), Inches(0.6)
        )
        tf = title_shape.text_frame
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.add_paragraph()
        p.text = f"CpG佐剂预防性疫苗安全性汇总 ({slide_idx + 1}/{total_chunks})"
        p.runs[0].font.size = Pt(26)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = DARK_RED
        p.runs[0].font.name = "Microsoft YaHei"

        # Divider Line REMOVED based on user request.

        # Add Table. Height = 1.0 forces auto-expand.
        table_shape = slide.shapes.add_table(
            len(chunk) + 1, 5, left_margin, top_margin, sum(col_widths), Inches(1.0)
        )
        tbl = table_shape.table

        # Compress header row height
        tbl.rows[0].height = Inches(0.35)

        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

        # Write headers
        for c_idx, h_text in enumerate(headers):
            cell = tbl.cell(0, c_idx)
            format_cell(
                cell,
                h_text,
                size=11,
                bold=True,
                text_color=HEADER_TEXT,
                bg_color=DARK_RED,
            )

        # Write data rows
        for r_idx, r_data in enumerate(chunk):
            bg_color = WHITE if r_idx % 2 == 0 else LIGHT_GRAY
            for c_idx, c_text in enumerate(r_data):
                cell = tbl.cell(r_idx + 1, c_idx)
                tbl.rows[r_idx + 1].height = Inches(0.5)
                # Apply hyperlink parser only to the final column (index 4)
                is_ref_column = c_idx == 4
                format_cell(
                    cell,
                    c_text,
                    size=8.5,
                    bold=False,
                    text_color=BLACK,
                    bg_color=bg_color,
                    is_ref=is_ref_column,
                )

    prs.save(ppt_path)
    print(f"Successfully saved PPT to: {ppt_path}")


if __name__ == "__main__":
    import traceback

    try:
        doc_path = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\CpG_Vaccine_Safety_Summary-V5-20260820.docx"
        # Output V4
        ppt_path = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\CpG_Vaccine_Safety_Summary_PPT_V4.pptx"
        create_ppt_from_docx(doc_path, ppt_path)
    except Exception:
        print("An error occurred:")
        traceback.print_exc()
