"""
generate_heplisav_safety_v5.py
==============================
V5 生成脚本：基于 V4，仅更新 Slide 7 中 Phase 2 (NCT00511095) 行的「安全性数据汇总」列。

修改内容（用户 2026-08-28 确认的推荐版文字）
--------------------------------------------
1. 征集性反应：调整语序——「第1/2剂」提到最前，再分别说明局部/全身发生率；
2. 非征集性AE：具体事件补充百分比（头痛31/207=15.0%、关节痛16/207=7.7%、背痛12/207=5.8%），
   并注明 CT.gov 仅列发生率＞5%事件；
3. SAE：因 CT.gov 未作因果判定，删除 7 项事件明细，仅保留总量 2/207例(1.0%)；
4. 免疫原性佐证：保留数据，微调措辞。

数据来源：ClinicalTrials.gov API（NCT00511095，2026-08-28 二次核验）

输出
----
review_materials/60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v5.pptx
"""

import copy

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

SRC = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v4.pptx"
DST = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v5.pptx"

# 用户确认的推荐版文字（Phase 2 安全性数据汇总，SAE 明细已删除）
NEW_SAFETY_LINES = [
    ("（源自CT.gov结果，2017-04-12发布）", None),
    (
        "▪ 征集性反应：第1/2剂后局部反应发生率38.2%/37.9%，全身反应发生率37.2%/33.0%，以轻中度为主",
        None,
    ),
    (
        "▪ 非征集性AE：59/207例（28.5%），以头痛31例（15.0%）、关节痛16例（7.7%）、背痛12例（5.8%）为主（CT.gov仅列发生率＞5%事件）",
        None,
    ),
    ("▪ SAE：2/207例（1.0%），CT.gov未作因果判定", None),
    ("▪ 免疫原性佐证：第28周血清保护率95.4%，抗-HBs GMC 349.4 mIU/mL，总体耐受性良好", None),
]


# ---------------------------------------------------------------------------
# 表格单元格级工具函数（与 V4 脚本保持一致）
# ---------------------------------------------------------------------------


def _find_first_rPr(txBody):
    """返回 txBody 中第一个 run 的 rPr 深拷贝（用作字体模板）。"""
    for p in txBody.findall(qn("a:p")):
        for r in p.findall(qn("a:r")):
            rPr = r.find(qn("a:rPr"))
            if rPr is not None:
                return copy.deepcopy(rPr)
    return None


def _make_default_rPr():
    """构造默认 rPr：Arial 10pt。"""
    from pptx.oxml import parse_xml

    xml = (
        '<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'sz="1000" lang="zh-CN">'
        '<a:latin typeface="Arial"/>'
        '<a:ea typeface="Arial"/>'
        "</a:rPr>"
    )
    return parse_xml(xml)


def set_cell_text(cell, lines):
    """
    整格重写文本。lines 为 [(text, url_or_None), ...]，
    保留单元格原有段落/run 格式（取第一个 run 的 rPr 为模板）。
    """
    part = cell.part
    txBody = cell.text_frame._txBody
    ps = txBody.findall(qn("a:p"))
    tmpl_rPr = _find_first_rPr(txBody)
    if tmpl_rPr is None:
        tmpl_rPr = _make_default_rPr()
    for p in ps:
        txBody.remove(p)
    for text, url in lines:
        p = txBody.makeelement(qn("a:p"), {})
        r = p.makeelement(qn("a:r"), {})
        rPr = copy.deepcopy(tmpl_rPr)
        for h in rPr.findall(qn("a:hlinkClick")):
            rPr.remove(h)
        r.append(rPr)
        t = p.makeelement(qn("a:t"), {})
        t.text = text
        r.append(t)
        p.append(r)
        if url:
            _set_para_hyperlink(p, url, part)
        txBody.append(p)


def _set_para_hyperlink(para, url, part):
    """为段落第一个 run 添加外链超链接。"""
    run = para.find(qn("a:r"))
    if run is None:
        return
    rPr = run.find(qn("a:rPr"))
    if rPr is None:
        rPr = run.makeelement(qn("a:rPr"), {})
        run.insert(0, rPr)
    for old in rPr.findall(qn("a:hlinkClick")):
        rPr.remove(old)
    rId = part.relate_to(url, RT.HYPERLINK, is_external=True)
    hlink = rPr.makeelement(qn("a:hlinkClick"), {})
    hlink.set(qn("r:id"), rId)
    hlink.set("tooltip", url)
    rPr.append(hlink)


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------


def find_phase2_row(table):
    """在表格中定位 Phase 2 (NCT00511095) 行，依据 col4 参考信息文本。"""
    for ri in range(len(table.rows)):
        ref = table.cell(ri, 4).text_frame.text
        if "NCT00511095" in ref:
            return ri
    raise RuntimeError("未找到 Phase 2 (NCT00511095) 行")


def build_v5():
    prs = Presentation(SRC)
    slide = prs.slides[6]  # Slide 7（HEPLISAV-B：Phase 1/2 页）

    table = None
    for shp in slide.shapes:
        if shp.has_table:
            table = shp.table
            break
    if table is None:
        raise RuntimeError("Slide 7 未找到表格")

    row_idx = find_phase2_row(table)
    print("定位到 Phase 2 行: row %d" % row_idx)

    # 更新前内容
    old_text = table.cell(row_idx, 3).text_frame.text
    print("----- 更新前 (col3) -----")
    print(old_text)

    set_cell_text(table.cell(row_idx, 3), NEW_SAFETY_LINES)

    # 更新后内容
    new_text = table.cell(row_idx, 3).text_frame.text
    print("----- 更新后 (col3) -----")
    print(new_text)

    prs.save(DST)
    print("已保存:", DST)
    print("总幻灯片数:", len(prs.slides))
    return DST


# ---------------------------------------------------------------------------
# 验证
# ---------------------------------------------------------------------------


def verify(path):
    prs = Presentation(path)
    print("\n===== 验证输出 =====")
    print("总幻灯片数:", len(prs.slides))
    for idx in (6, 7, 8):
        slide = prs.slides[idx]
        title = ""
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip() and not shp.has_table:
                title = shp.text_frame.text
                break
        print("Slide %d: %r" % (idx + 1, title))
    # Slide 7 表格各行首列 + Phase 2 行 col3
    slide = prs.slides[6]
    for shp in slide.shapes:
        if shp.has_table:
            tbl = shp.table
            for ri in range(len(tbl.rows)):
                c0 = tbl.cell(ri, 0).text_frame.text.replace("\n", "|")
                c2 = tbl.cell(ri, 2).text_frame.text.replace("\n", "|")[:40]
                print("  R%d: %r | %r" % (ri, c0[:30], c2))
    # 超链接完整性：Phase 2 行 col4
    for shp in slide.shapes:
        if shp.has_table:
            tbl = shp.table
            for ri in range(len(tbl.rows)):
                cell = tbl.cell(ri, 4)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        rPr = run._r.find(qn("a:rPr"))
                        if rPr is not None and rPr.find(qn("a:hlinkClick")) is not None:
                            print("  Slide7 R%d 含超链接: %r" % (ri, run.text))


if __name__ == "__main__":
    out = build_v5()
    verify(out)
    print("\n完成。")
