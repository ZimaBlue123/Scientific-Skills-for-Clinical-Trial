import sys

from pptx import Presentation


def extract_text_from_pptx(pptx_path, output_path):
    try:
        prs = Presentation(pptx_path)
        with open(output_path, "w", encoding="utf-8") as f:
            for i, slide in enumerate(prs.slides):
                f.write(f"--- Slide {i+1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        f.write(shape.text + "\n")
    except Exception as e:
        print(f"Error reading {pptx_path}: {e}")


if __name__ == "__main__":
    if len(sys.argv) > 2:
        extract_text_from_pptx(sys.argv[1], sys.argv[2])
    else:
        print("Please provide the path to the PPTX file and output txt file.")
