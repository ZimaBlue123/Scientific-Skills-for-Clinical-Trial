#!/usr/bin/env python3
"""Generate a 2-page English PPTX overview for TVAX-006 (recombinant zoster vaccine, CHO).

Page 1: Product Overview
Page 2: Clinical Development Strategy

Style: 16:9 (13.333 x 7.5 in), pure white background, dark red accent RGB(192,0,0).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Palette
# ----------------------------------------------------------------------------
DARK_RED = RGBColor(0xC0, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
BORDER_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _no_line(shape):
    shape.line.fill.background()


def _line(shape, rgb, width_pt=1.5):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=1.5, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        _set_fill(sp, fill)
    else:
        sp.fill.background()
    if line is not None:
        _line(sp, line, line_w)
    else:
        _no_line(sp)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, runs, size=11, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=2, line_spacing=1.0):
    """runs: str, or list of paragraphs; each paragraph is str or list of (text, bold, color)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if isinstance(runs, str):
        runs = [runs]

    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, bold, color)]
        for text, b, c in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = b
            r.font.color.rgb = c
    return tb


def add_header(slide, title_text):
    """Dark red banner + white bold title."""
    add_rect(slide, 0, 0, 13.333, 0.9, fill=DARK_RED)
    add_text(slide, 0.6, 0, 12.1, 0.9, title_text, size=28, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def add_page_number(slide, page_label):
    """Small gray page number at bottom-right (no footer strip)."""
    add_text(slide, 11.55, 7.05, 1.6, 0.35, page_label, size=10, bold=False, color=MID_GRAY,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def section_title(slide, x, y, w, text):
    add_text(slide, x, y, w, 0.35, text, size=16, bold=True, color=DARK_RED,
             align=PP_ALIGN.LEFT)


# ============================================================================
# PAGE 1 — Product Overview
# ============================================================================
def build_page1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Product Overview")

    # Subtitle
    add_text(slide, 0.6, 1.02, 12.1, 0.5,
             "Recombinant Zoster Vaccine (CHO Cell)  —  TVAX-006",
             size=20, bold=True, color=DARK_GRAY)

    # ---- Left column: Product Specifications ----
    section_title(slide, 0.6, 1.72, 5.7, "Product Specifications")
    specs = [
        ("Product / code", "Recombinant Zoster Vaccine (CHO Cell) — TVAX-006"),
        ("Active ingredient", "VZV glycoprotein E (gE), 50 μg per dose"),
        ("Expression system", "Chinese Hamster Ovary (CHO) cells"),
        ("Adjuvant", "TVA01 adjuvant"),
        ("Dosage form", "Injection (gE lyophilized powder + TVA01 adjuvant suspension)"),
        ("Strength", "0.5 mL per human dose"),
        ("Route", "Intramuscular injection (deltoid muscle)"),
        ("Indication", "Prevention of herpes zoster (HZ)"),
        ("Target population", "Adults aged ≥ 40 years"),
        ("Registration", "China (confirmed); other regions TBD"),
    ]
    paras = []
    for label, value in specs:
        paras.append([(label + ":  ", True, BLACK), (value, False, DARK_GRAY)])
    add_text(slide, 0.6, 2.16, 5.75, 3.6, paras, size=12, line_spacing=1.05, space_after=6)

    # ---- Right column: Key Composition ----
    section_title(slide, 6.6, 1.72, 6.2, "Key Composition")
    comp_boxes = [
        ("Antigen", "Varicella-zoster virus glycoprotein E (gE)", MSO_SHAPE.HEXAGON),
        ("Adjuvant", "TVA01 adjuvant system", MSO_SHAPE.TEAR),
        ("Strength", "gE 50 μg / 0.5 mL per dose", MSO_SHAPE.LIGHTNING_BOLT),
    ]
    y = 2.12
    for tag, desc, icon_shape in comp_boxes:
        add_rect(slide, 6.6, y, 6.2, 0.82, fill=DARK_RED)
        # monochrome icon (white glyph on dark-red box)
        add_rect(slide, 6.88, y + 0.19, 0.44, 0.44, fill=WHITE, shape=icon_shape)
        add_text(slide, 7.5, y + 0.08, 5.15, 0.66,
                 [[(tag + "  ", True, WHITE), (desc, False, WHITE)]],
                 size=13, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        y += 0.98

    # ---- Bottom: Immunization Schedule ----
    section_title(slide, 0.6, 5.5, 12.0, "Immunization Schedule")
    # two dose circles + arrow
    cy = 6.12
    add_rect(slide, 0.6, cy, 0.5, 0.5, fill=DARK_RED, shape=MSO_SHAPE.OVAL)
    add_text(slide, 0.6, cy, 0.5, 0.5, "1", size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 1.3, cy + 0.2, 0.6, 0.12, fill=DARK_RED, shape=MSO_SHAPE.RIGHT_ARROW)
    add_rect(slide, 2.1, cy, 0.5, 0.5, fill=DARK_RED, shape=MSO_SHAPE.OVAL)
    add_text(slide, 2.1, cy, 0.5, 0.5, "2", size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, 2.95, cy + 0.02, 9.8, 0.55,
             [[("Two doses (0.5 mL each), intramuscular.  ", True, BLACK),
               ("Dose 2 at Month 2 (2–6 months after Dose 1).", False, DARK_GRAY)]],
             size=14, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    # storage line
    add_text(slide, 0.6, 6.72, 12.0, 0.3,
             "Storage: 2–8 °C; tentative shelf life 36 months",
             size=10.5, color=MID_GRAY)

    add_page_number(slide, "1 / 2")


# ============================================================================
# PAGE 2 — Clinical Development Strategy
# ============================================================================
def build_page2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Clinical Development Strategy")

    phases = [
        {
            "title": "Phase I — Australia",
            "status": "Completed",
            "rows": [
                "Single-center, randomized, double-blind, placebo- & Shingrix®-controlled",
                "Healthy adults 30–70 y",
                "N = 100",
                "Endpoint: Safety & tolerability",
            ],
        },
        {
            "title": "Phase I — China",
            "status": "Completed",
            "rows": [
                "Randomized, blinded, placebo-controlled",
                "Adults ≥ 40 y",
                "N = 100",
                "Endpoint: Safety & tolerability",
            ],
        },
        {
            "title": "Phase II — China",
            "status": "Ongoing",
            "rows": [
                "Randomized, Ganwei®/Shingrix®-controlled",
                "Adults ≥ 40 y",
                "N = 420",
                "Endpoint: Immunogenicity & safety",
            ],
        },
        {
            "title": "Phase III — China",
            "status": "Planned",
            "rows": [
                "Multicenter, randomized, double-blind, placebo-controlled, superiority",
                "Adults ≥ 40 y",
                "N = 16,000",
                "Endpoint: Vaccine efficacy (HZ)",
            ],
        },
    ]

    box_w = 2.82
    gap = 0.33
    x0 = 0.54
    top = 1.25
    box_h = 3.55

    for i, ph in enumerate(phases):
        x = x0 + i * (box_w + gap)
        # phase card
        add_rect(slide, x, top, box_w, box_h, fill=WHITE, line=DARK_RED, line_w=1.5)
        # title band
        add_rect(slide, x, top, box_w, 0.52, fill=DARK_RED)
        add_text(slide, x + 0.12, top, box_w - 0.24, 0.52, ph["title"],
                 size=13.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # status chip
        add_text(slide, x + 0.14, top + 0.6, box_w - 0.28, 0.3, ph["status"],
                 size=11.5, bold=True, color=DARK_RED)
        # info rows
        paras = []
        for row in ph["rows"]:
            paras.append([("•  ", True, DARK_RED), (row, False, DARK_GRAY)])
        add_text(slide, x + 0.14, top + 0.92, box_w - 0.28, box_h - 1.0, paras,
                 size=11.5, line_spacing=1.02, space_after=5)

        # connector arrow
        if i < 3:
            ax = x + box_w + 0.02
            add_rect(slide, ax, top + box_h / 2 - 0.11, gap - 0.04, 0.22,
                     fill=DARK_RED, shape=MSO_SHAPE.RIGHT_ARROW)

    # ---- Bottom summary: Safety / Immunogenicity ----
    summ_top = 5.02
    summ_h = 1.42
    col_w = 6.0

    add_rect(slide, x0, summ_top, col_w, summ_h, fill=LIGHT_GRAY)
    add_text(slide, x0 + 0.15, summ_top + 0.08, col_w - 0.3, 0.3, "Safety (Phase I & II)",
             size=13, bold=True, color=DARK_RED)
    add_text(slide, x0 + 0.15, summ_top + 0.42, col_w - 0.3, summ_h - 0.5,
             [
                 "•  Well tolerated; AEs mostly Grade 1–2 and transient",
                 "•  No vaccine-related SAEs / AESIs; no AE-related withdrawals",
                 "•  Favorable tolerability in adults aged ≥ 60 years",
             ], size=11.5, line_spacing=1.05, space_after=3)

    x2 = x0 + col_w + 0.25
    add_rect(slide, x2, summ_top, col_w, summ_h, fill=LIGHT_GRAY)
    add_text(slide, x2 + 0.15, summ_top + 0.08, col_w - 0.3, 0.3,
             "Immunogenicity (Phase I & II)", size=13, bold=True, color=DARK_RED)
    add_text(slide, x2 + 0.15, summ_top + 0.42, col_w - 0.3, summ_h - 0.5,
             [
                 "•  Robust anti-gE / anti-VZV humoral + CD4+ T-cell responses",
                 "•  ≥ 50 y: comparable to Shingrix®; 40–49 y: superior to live-attenuated control",
                 "•  Peak ~1 month post Dose 2; durable through 6 months",
             ], size=11.5, line_spacing=1.05, space_after=3)

    # ---- Current status sentence ----
    add_rect(slide, x0, 6.48, 12.25, 0.56, fill=DARK_RED)
    add_text(slide, x0 + 0.2, 6.48, 11.85, 0.56,
             "Phase I (Australia & China) completed; Phase II ongoing under immunogenicity "
             "persistence and long-term safety follow-up — now at EOP2, preparing for CDE pre-Phase III communication.",
             size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    add_page_number(slide, "2 / 2")


# ============================================================================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_page1(prs)
    build_page2(prs)

    out = r"E:/Cursor Project/2-Scientific-Skills-for-Clinical_Trial/review_materials/006 PPT/TVAX-006_Clinical_Development_Overview_20260831.pptx"
    prs.save(out)
    print("SAVED:", out)
    print("slides:", len(prs.slides))


if __name__ == "__main__":
    main()
