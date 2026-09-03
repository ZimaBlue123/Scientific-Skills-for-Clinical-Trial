# -*- coding: utf-8 -*-
"""
TVAX-009 糖尿病亚组免疫原性结果（仅 PPS）—— 4 页 16:9 PPT 生成器。

将 Word 版（PPS 4 表）转成 4 张幻灯片，版式沿用 Word 的「组别作行、时间点作列」。

规格（用户指定）：
- 画布 16:9 宽屏：13.333 × 7.5 英寸（12192000 × 6858000 EMU），全局纯白背景 #FFFFFF。
- 品牌色：主强调色（表头背景、幻灯片标题字体）暗红 RGB(192,0,0) = C00000。

从零构建（add_slide blank layout + 新建表格 XML），无 deepcopy，规避 creationId 重复。
4 页：
  1. 抗-HBs 阳转率（PPS）绝对时间点（10 列 × 15 行）
  2. 抗-HBs 阳转率（PPS）相对时间点（5 列 × 15 行）
  3. 抗-HBs GMC（PPS）绝对时间点（10 列 × 15 行）
  4. 抗-HBs GMC（PPS）相对时间点（5 列 × 15 行）
"""
import json
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BASE = 'E:/Cursor Project/2-Scientific-Skills-for-Clinical_Trial/review_materials'
OUT = BASE + '/TVAX-009-糖尿病亚组免疫原性结果-PPS-4页PPT-20260903.pptx'

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
BRAND = 'C00000'        # 暗红 RGB(192,0,0)
WHITE = 'FFFFFF'
BLACK = '000000'
GREY = 'F2F2F2'

with open(BASE + '/diabetes_immuno_result.json', encoding='utf-8') as f:
    DATA = json.load(f)

PPS = DATA['pps']

GROUPS = [
    ('A1', '低剂量组(A1)', '0,1月程序', None),
    ('A2', '高剂量组(A2)', '0,1月程序', None),
    ('B1', '低剂量组(B1)', '0,2月程序', None),
    ('B2', '高剂量组(B2)', '0,2月程序', None),
    ('C1', '阳性对照组(C1)', '0,1,6月程序', '18-59岁'),
    ('C3', '0,1,6月高剂量组(C3)', '0,1,6月程序', '≥60岁'),
    ('C2', '阳性对照组(C2)', '0,1,6月程序', '≥60岁'),
]
PROGRAMS = [
    ('0,1月程序', ['A1', 'A2']),
    ('0,2月程序', ['B1', 'B2']),
    ('0,1,6月程序', ['C1', 'C3', 'C2']),
]
PROG_OF = {code: prog for code, _, prog, _ in GROUPS}
LABEL = {code: lab for code, lab, _, _ in GROUPS}
ANNO = {code: a for code, _, _, a in GROUPS}

PROG_TIMEPOINTS = {
    '0,1月程序': ['M1', 'M2', 'M3', 'M6', 'M7', 'M8'],
    '0,2月程序': ['M1', 'M2', 'M3', 'M4', 'M7', 'M8'],
    '0,1,6月程序': ['M1', 'M2', 'M6', 'M7', 'M8'],
}
ALL_TP = ['M1', 'M2', 'M3', 'M4', 'M6', 'M7', 'M8']
REL_MAP = {
    '全免后1个月': {'0,1月程序': 'M2', '0,2月程序': 'M3', '0,1,6月程序': 'M7'},
    '全免后2个月': {'0,1月程序': 'M3', '0,2月程序': 'M4', '0,1,6月程序': 'M8'},
}


def _esc(s):
    return s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')


LN_T = ('<a:lnT w="12700" cap="flat" cmpd="sng" algn="ctr">'
        '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
        '<a:prstDash val="solid"/><a:round/>'
        '<a:headEnd type="none" w="med" len="med"/>'
        '<a:tailEnd type="none" w="med" len="med"/></a:lnT>')


def _ln_b(solid):
    if solid:
        return ('<a:lnB w="12700" cap="flat" cmpd="sng" algn="ctr">'
                '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
                '<a:prstDash val="solid"/><a:round/>'
                '<a:headEnd type="none" w="med" len="med"/>'
                '<a:tailEnd type="none" w="med" len="med"/></a:lnB>')
    return '<a:lnB><a:noFill/></a:lnB>'


def make_tc(text, bold=False, grid_span=None, row_span=None, lnB_solid=False, lnT=True,
            fill=None, font_size=900, text_color=BLACK):
    attrs = ''
    if grid_span:
        attrs += ' gridSpan="%d"' % grid_span
    if row_span:
        attrs += ' rowSpan="%d"' % row_span
    b = ' b="1"' if bold else ''
    s = '<a:tc xmlns:a="%s"%s><a:txBody><a:bodyPr/><a:lstStyle/>' % (A_NS, attrs)
    for line in text.split('\n'):
        s += ('<a:p><a:pPr indent="0" algn="ctr">'
              '<a:lnSpc><a:spcPct val="100000"/></a:lnSpc><a:buNone/></a:pPr>'
              '<a:r><a:rPr lang="zh-CN" altLang="en-US" sz="%d" kern="0"%s>'
              '<a:solidFill><a:srgbClr val="%s"/></a:solidFill><a:effectLst/>'
              '<a:latin typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
              '<a:ea typeface="微软雅黑" panose="020B0503020204020204" pitchFamily="34" charset="-122"/>'
              '<a:cs typeface="Times New Roman" panose="02020603050405020304" charset="0"/>'
              '<a:sym typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
              '</a:rPr><a:t>%s</a:t></a:r></a:p>' % (font_size, b, text_color, _esc(line)))
    s += '</a:txBody>'
    s += '<a:tcPr marL="36000" marR="0" marT="0" marB="0" anchor="ctr">'
    s += '<a:lnL><a:noFill/></a:lnL><a:lnR><a:noFill/></a:lnR>'
    if lnT:
        s += LN_T
    s += _ln_b(lnB_solid)
    if fill:
        s += '<a:solidFill><a:srgbClr val="%s"/></a:solidFill>' % fill
    else:
        s += '<a:noFill/>'
    s += '</a:tcPr></a:tc>'
    return s


def table_str(col_widths, rows):
    s = '<a:tbl xmlns:a="%s"><a:tblPr firstRow="0" firstCol="0" bandRow="0"/><a:tblGrid>' % A_NS
    for w in col_widths:
        s += '<a:gridCol w="%d"/>' % w
    s += '</a:tblGrid>'
    total_h = 0
    for h, cells in rows:
        s += '<a:tr h="%d">%s</a:tr>' % (h, ''.join(cells))
        total_h += h
    s += '</a:tbl>'
    return s, total_h


def sero_val(code, tag, tp):
    rec = PPS[code][tag][tp]
    if rec['N_阳转'] == 0 and rec['pct'] is None:
        return '—'
    return '%d/%d\n(%.2f%%)' % (rec['n_阳转'], rec['N_阳转'], rec['pct'])


def gmc_val(code, tag, tp):
    rec = PPS[code][tag][tp]
    if rec['N_gmc'] == 0 and rec['gmc'] is None:
        return '—'
    return '%.2f\n(N=%d)' % (rec['gmc'], rec['N_gmc'])


def group_label(code):
    lab = LABEL[code]
    if ANNO[code]:
        lab += '\n(%s)' % ANNO[code]
    return lab


def build_absolute(mode):
    col_w = [914400, 1645920, 640080] + [1215000] * 7  # 程序/组别/病史 + M1..M8
    header_h = 365760
    data_h = 345600
    rows = []
    # 表头（10 格）
    cells = [make_tc('程序', bold=True, lnB_solid=True, fill=BRAND, text_color=WHITE),
             make_tc('组别', bold=True, lnB_solid=True, fill=BRAND, text_color=WHITE),
             make_tc('糖尿病病史', bold=True, lnB_solid=True, fill=BRAND, text_color=WHITE)]
    for tp in ALL_TP:
        cells.append(make_tc(tp, bold=True, lnB_solid=True, fill=BRAND, text_color=WHITE))
    rows.append((header_h, cells))

    # 数据（rowSpan 占据的列在后续行不生成单元格）
    for prog, codes in PROGRAMS:
        n = len(codes) * 2
        for ci, code in enumerate(codes):
            for tag in ('有', '无'):
                cells = []
                if ci == 0 and tag == '有':
                    cells.append(make_tc(prog, row_span=n, fill=None))
                if tag == '有':
                    cells.append(make_tc(group_label(code), row_span=2, fill=None))
                cells.append(make_tc(tag, fill=GREY))
                last_row = (prog == PROGRAMS[-1][0] and code == PROGRAMS[-1][1][-1] and tag == '无')
                for tp in ALL_TP:
                    if tp in PROG_TIMEPOINTS[PROG_OF[code]]:
                        v = sero_val(code, tag, tp) if mode == 'sero' else gmc_val(code, tag, tp)
                    else:
                        v = '—'
                    cells.append(make_tc(v, lnB_solid=last_row))
                rows.append((data_h, cells))
    xml, h = table_str(col_w, rows)
    return xml, sum(col_w), h


def build_relative(mode):
    col_w = [914400, 1645920, 640080, 4251960, 4251960]  # 程序/组别/病史 + 全免后1/2个月
    header_h = 365760
    data_h = 345600
    rows = []
    cells = [make_tc('程序', bold=True, lnB_solid=True, fill=BRAND, text_color=WHITE),
             make_tc('组别', bold=True, lnB_solid=True, fill=BRAND, text_color=WHITE),
             make_tc('糖尿病病史', bold=True, lnB_solid=True, fill=BRAND, text_color=WHITE),
             make_tc('全免后1个月', bold=True, lnB_solid=True, fill=BRAND, text_color=WHITE),
             make_tc('全免后2个月', bold=True, lnB_solid=True, fill=BRAND, text_color=WHITE)]
    rows.append((header_h, cells))

    for prog, codes in PROGRAMS:
        n = len(codes) * 2
        for ci, code in enumerate(codes):
            for tag in ('有', '无'):
                cells = []
                if ci == 0 and tag == '有':
                    cells.append(make_tc(prog, row_span=n, fill=None))
                if tag == '有':
                    cells.append(make_tc(group_label(code), row_span=2, fill=None))
                cells.append(make_tc(tag, fill=GREY))
                last_row = (prog == PROGRAMS[-1][0] and code == PROGRAMS[-1][1][-1] and tag == '无')
                for tp_label in ('全免后1个月', '全免后2个月'):
                    m = REL_MAP[tp_label][PROG_OF[code]]
                    v = sero_val(code, tag, m) if mode == 'sero' else gmc_val(code, tag, m)
                    cells.append(make_tc(v, lnB_solid=last_row))
                rows.append((data_h, cells))
    xml, h = table_str(col_w, rows)
    return xml, sum(col_w), h


def add_textbox(slide, text, left, top, width, height, size_pt, bold, color, fill=None,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.name = '微软雅黑'
        run.font.color.rgb = RGBColor.from_string(color)
    if fill:
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        tb.fill.background()
    tb.line.fill.background()
    return tb


def add_table_frame(slide, xml_str, left, top, width, height):
    placeholder = slide.shapes.add_table(1, 1, Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)))
    gf = placeholder._element
    tbl_el = etree.fromstring(xml_str)
    old_tbl = gf.find('.//' + qn('a:tbl'))
    old_tbl.getparent().replace(old_tbl, tbl_el)
    return gf


def set_white_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def build_slide(prs, title, desc, xml_str, tw, th):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)

    # 标题（暗红）
    add_textbox(slide, title, 274320, 228600, 11430000, 411480,
                size_pt=20, bold=True, color=BRAND, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # 表格
    add_table_frame(slide, xml_str, 228600, 777240, tw, th)

    # 底部说明
    add_textbox(slide, desc, 228600, 6035040, 11704320, 685800,
                size_pt=10, bold=False, color='333333', align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)


def main():
    specs = [
        ('有、无糖尿病病史人群接种后抗-HBs阳转率（PPS）——绝对时间点',
         '本分析纳入基线合并糖尿病病史的受试者共 30 例（18-59 岁 18 例、≥60 岁 12 例）。'
         '有糖尿病病史人群完成全程 2 剂/3 剂接种后阳转率总体达较高水平，0,2月程序两组在 M3 均达 100.00%，'
         '0,1,6月高剂量组及两个阳性对照组在 M7 均达 100.00%。',
         build_absolute('sero')),
        ('有、无糖尿病病史人群接种后抗-HBs阳转率（PPS）——相对时间点',
         '全免后 1 个月、2 个月，各组有、无糖尿病病史人群的阳转率均维持在较高水平。'
         '注：全免后 1/2 个月在各程序对应时间点——0,1月程序 M2/M3，0,2月程序 M3/M4，0,1,6月程序 M7/M8。',
         build_relative('sero')),
        ('有、无糖尿病病史人群接种后抗-HBs GMC（PPS）——绝对时间点',
         '各组有、无糖尿病病史人群完成全程 2 剂/3 剂接种后各时间点抗体 GMC 均呈上升趋势。'
         '抗-HBs 浓度低于检测下限（<2.00 mIU/mL）的样本按 LLOQ/2（1.00 mIU/mL）计入 GMC 计算。',
         build_absolute('gmc')),
        ('有、无糖尿病病史人群接种后抗-HBs GMC（PPS）——相对时间点',
         '全免后 1 个月、2 个月，各组有、无糖尿病病史人群的抗体 GMC 均维持在一定水平。'
         '注：全免后 1/2 个月在各程序对应时间点——0,1月程序 M2/M3，0,2月程序 M3/M4，0,1,6月程序 M7/M8。',
         build_relative('gmc')),
    ]

    prs = Presentation()
    prs.slide_width = 12192000   # 13.333 in
    prs.slide_height = 6858000   # 7.5 in

    for title, desc, (xml_str, tw, th) in specs:
        build_slide(prs, title, desc, xml_str, tw, th)

    prs.save(OUT)
    print('saved:', OUT)
    print('total slides:', len(prs.slides))


if __name__ == '__main__':
    main()
