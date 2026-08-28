"""V8 输出自检：结构、内容、链接、溢出估算"""

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

P = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v8.pptx"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

prs = Presentation(P)
print(f"== 总页数: {len(prs.slides._sldIdLst)} (v7 为 23) ==")
print(f"slide size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} in")

# 每页标题
print("\n== 各页标题 ==")
for i, slide in enumerate(prs.slides):
    titles = []
    for sh in slide.shapes:
        if sh.has_text_frame and not sh.has_table:
            t = sh.text_frame.text.strip()
            if t:
                titles.append(t[:60])
    print(f"  Slide {i + 1}: {titles[0] if titles else '(no textbox)'}")


def link_map(slide):
    rels = {}
    for rid, rel in slide.part.rels.items():
        if "hyperlink" in rel.reltype:
            rels[rid] = rel.target_ref
    return rels


def show_table(slide, label, max_cell=500):
    for sh in slide.shapes:
        if sh.has_table:
            tb = sh.table
            lm = link_map(slide)
            print(
                f"\n== {label}: rows={len(tb.rows)} cols={len(tb.columns)} top={Emu(sh.top).inches:.2f}in"
            )
            print(f"   row heights(声明): {[round(Emu(r.height).inches, 2) for r in tb.rows]}")
            print(
                f"   declared total={sum(Emu(r.height).inches for r in tb.rows):.2f}in bottom_edge={Emu(sh.top).inches + sum(Emu(r.height).inches for r in tb.rows):.2f}in"
            )
            for ri, row in enumerate(tb.rows):
                print(f"  --- row {ri} ---")
                for ci, cell in enumerate(row.cells):
                    parts = []
                    for para in cell.text_frame.paragraphs:
                        runs = []
                        for run in para.runs:
                            hlink = run._r.find("{%s}hlinkClick" % A_NS)
                            if hlink is not None:
                                rid = hlink.get("{%s}id" % R_NS)
                                runs.append(f"[{run.text}->{lm.get(rid, rid)}]")
                            else:
                                runs.append(run.text)
                        parts.append("".join(runs))
                    txt = " || ".join(parts)
                    if len(txt) > max_cell:
                        txt = txt[:max_cell] + f"...[共{len(''.join(parts))}字]"
                    print(f"    c{ci}: {txt}")


# Slide 8
show_table(prs.slides[7], "Slide 8")

# Slide 9
show_table(prs.slides[8], "Slide 9")

# 链接完整性：统计所有 slide 的 hlinkClick 有效性
print("\n== 链接完整性检查 ==")
total_bad = 0
for i, slide in enumerate(prs.slides):
    lm = link_map(slide)
    n_hlink = 0
    bad = 0
    for hlink in slide.shapes._spTree.iter(qn("a:hlinkClick")):
        rid = hlink.get("{%s}id" % R_NS)
        n_hlink += 1
        if rid not in lm:
            bad += 1
            print(f"  Slide {i + 1}: BROKEN hlink rId={rid} (not in rels)")
    if bad:
        total_bad += bad
    print(f"  Slide {i + 1}: {n_hlink} links, {bad} broken")
print(f"  TOTAL broken: {total_bad}")

# 孤儿 rel 检查（未被引用的超链接 rel）
print("\n== 孤儿 rel 检查 ==")
for i, slide in enumerate(prs.slides):
    used = set()
    for hlink in slide.shapes._spTree.iter(qn("a:hlinkClick")):
        rid = hlink.get("{%s}id" % R_NS)
        if rid:
            used.add(rid)
    orphans = [
        rid
        for rid, rel in slide.part.rels.items()
        if "hyperlink" in rel.reltype and rid not in used
    ]
    if orphans:
        print(f"  Slide {i + 1}: ORPHAN rels {orphans}")
print("  (无输出=全部干净)")

# 溢出估算：基于文本量粗估渲染高度
print("\n== 溢出估算（粗估） ==")
for si in (7, 8):
    slide = prs.slides[si]
    for sh in slide.shapes:
        if sh.has_table:
            tb = sh.table
            top = Emu(sh.top).inches
            page_h = 7.5
            for ri, row in enumerate(tb.rows):
                # 估算每行渲染高度：按最长单元格文本行数
                max_chars_per_line = 38  # 10pt 中文在 ~5.8in 列
                max_lines = 0
                for cell in row.cells:
                    n_chars = len(cell.text.replace(chr(10), ""))
                    n_paras = len(cell.text_frame.paragraphs)
                    est = 0
                    for para in cell.text_frame.paragraphs:
                        ln = max(1, (len(para.text) + max_chars_per_line - 1) // max_chars_per_line)
                        est += ln
                    max_lines = max(max_lines, est)
                est_h = max_lines * 0.17 + 0.1  # 每行0.17in + 边距
                print(
                    f"  Slide {si + 1} row{ri}: 声明{Emu(row.height).inches:.2f}in 粗估渲染~{est_h:.2f}in"
                )
