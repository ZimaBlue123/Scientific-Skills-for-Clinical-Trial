"""
diagnose_heplisav_links.py
==========================
诊断 V5 Slide 7/8 超链接问题：
1) V3 原 slide 7 表格 col4 各行的链接 URL（作为正确基准）
2) V5 Slide 7 FIH 行 col4 链接实际指向（疑似错位）
3) V5 Slide 8 part rels 中是否存在 hlinkClick 引用的 rId（疑似悬空）
"""

from pptx import Presentation
from pptx.oxml.ns import qn

V3 = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v3.pptx"
V5 = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v5.pptx"


def para_url(para, part):
    for run in para.runs:
        rPr = run._r.find(qn("a:rPr"))
        if rPr is None:
            continue
        hlink = rPr.find(qn("a:hlinkClick"))
        if hlink is None:
            continue
        rId = hlink.get(qn("r:id"))
        if not rId:
            continue
        try:
            return rId, part.rels[rId].target_ref
        except KeyError:
            return rId, "!! rId 不存在于 part.rels !!"
    return None, None


def dump_table_links(path, slide_idx, label):
    prs = Presentation(path)
    slide = prs.slides[slide_idx]
    print("\n===== %s (Slide %d) =====" % (label, slide_idx + 1))
    for shp in slide.shapes:
        if not shp.has_table:
            continue
        part = shp.part
        tbl = shp.table
        print("  part rels 全部条目:")
        for rId, rel in sorted(part.rels.items(), key=lambda kv: (len(kv[0]), kv[0])):
            print(
                "    %s -> %s (%s)"
                % (rId, rel.reltype.split("/")[-1], getattr(rel, "target_ref", "?"))
            )
        for ri in range(len(tbl.rows)):
            cell = tbl.cell(ri, 4)
            for para in cell.text_frame.paragraphs:
                rId, url = para_url(para, part)
                if rId:
                    print("  R%d col4: %-38r rId=%s -> %s" % (ri, para.text[:38], rId, url))


if __name__ == "__main__":
    dump_table_links(V3, 6, "V3 原 slide 7（正确基准）")
    dump_table_links(V5, 6, "V5 Slide 7")
    dump_table_links(V5, 7, "V5 Slide 8")
