"""
修复 HEPLISAV-B 安全性汇总 PPT 的坏链接，输出 v6
==================================================
输入: v5 (用户调整版, 23页)
输出: v6

修复内容:
1. Slide 7 (HEPLISAV-B Phase 1/2 页): FIH 行 2 段链接错位
   - 'NCT00095160'  文本实际链接到 NCT01999699  → 改为 https://clinicaltrials.gov/study/NCT00095160
   - 'PMID: 12744879' 文本实际链接到 NCT00511095 → 改为 https://pubmed.ncbi.nlm.nih.gov/12744879/
2. Slide 8 (HEPLISAV-B Phase 3 页): 13 段(16 run)链接 target 为 rId 占位字符串
   → 按 run 文本→V3 权威 URL 映射重建（映射来源: V3 源文件 slide7.xml 的 rels）

方法: 对目标 slide 的每个 a:r/rPr/a:hlinkClick, 按 run 文本查映射表,
      用 part.relate_to() 建立真实外部超链接关系并重写 r:id,
      然后清理未使用的 rId 占位关系。
"""

import re
import sys

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

SRC = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v5.pptx"
DST = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v6.pptx"

# 权威链接映射: run 文本 → 正确 URL（源自 V3 源文件 slide7.xml 的 rels）
LINK_MAP = {
    # Slide 7 FIH 行
    "NCT00095160": "https://clinicaltrials.gov/study/NCT00095160",
    "PMID: 12744879": "https://pubmed.ncbi.nlm.nih.gov/12744879/",
    # Slide 8 Phase 3 注册综合数据
    "NCT00435812": "https://clinicaltrials.gov/study/NCT00435812",
    "NCT01005407": "https://clinicaltrials.gov/study/NCT01005407",
    "NCT02117934": "https://clinicaltrials.gov/study/NCT02117934",
    "FDA Package Insert": "https://www.fda.gov/media/108745/download",
    "FDA Clinical Review": "https://www.fda.gov/media/109808/download",
    "PMID: 29628151": "https://pubmed.ncbi.nlm.nih.gov/29628151/",
    "PMID: 23727002": "https://pubmed.ncbi.nlm.nih.gov/23727002/",
    # Slide 8 Phase 3 特殊人群
    "NCT00985426": "https://clinicaltrials.gov/study/NCT00985426",
    "NCT01282762": "https://clinicaltrials.gov/study/NCT01282762",
    "PMID: 25576215": "https://pubmed.ncbi.nlm.nih.gov/25576215/",
    "PMID: 37085451": "https://pubmed.ncbi.nlm.nih.gov/37085451/",
}


def fix_links_in_slide(prs, slide, link_map, label):
    """遍历 slide 所有 run 级 hlinkClick, 按 run 文本映射重建 target。
    返回 (修复数, 未匹配列表)"""
    part = slide.part
    fixed = 0
    unmatched = []
    for r_elem in slide._element.iter(qn("a:r")):
        rPr = r_elem.find(qn("a:rPr"))
        if rPr is None:
            continue
        hl = rPr.find(qn("a:hlinkClick"))
        if hl is None:
            continue
        rtext = "".join(t.text or "" for t in r_elem.findall(qn("a:t")))
        if rtext in link_map:
            url = link_map[rtext]
            new_rid = part.relate_to(url, RT.HYPERLINK, is_external=True)
            hl.set(qn("r:id"), new_rid)
            # 移除可能存在的 action 占位属性
            for attr in ("action", "tgtFrame"):
                if hl.get(attr) is not None:
                    del hl.attrib[attr]
            fixed += 1
            print("  [%s] %-22r -> %s" % (label, rtext, url))
        else:
            unmatched.append(rtext)
    return fixed, unmatched


def cleanup_dangling_rels(slide, label):
    """删除 slide part 中 target 为 rIdN 占位的超链接关系"""
    removed = 0
    rels = slide.part.rels
    for rid in list(rels.keys()):
        rel = rels[rid]
        if rel.reltype.endswith("/hyperlink") and rel.is_external:
            tgt = rel.target_ref if hasattr(rel, "target_ref") else str(rel.target)
            if re.fullmatch(r"rId\d+", str(tgt).strip()):
                slide.part.drop_rel(rid)
                removed += 1
    if removed:
        print("  [%s] 清理占位关系 %d 条" % (label, removed))
    return removed


def main():
    prs = Presentation(SRC)
    print("打开 v5: %d 页" % len(prs.slides))

    # --- 修复 Slide 7 (index 6) ---
    slide7 = prs.slides[6]
    print("\n===== Slide 7 (Phase 1/2 页) =====")
    fixed7, unm7 = fix_links_in_slide(prs, slide7, LINK_MAP, "S7")
    if unm7:
        print("  S7 未匹配(保留原链接): %r" % unm7)
    cleanup_dangling_rels(slide7, "S7")

    # --- 修复 Slide 8 (index 7) ---
    slide8 = prs.slides[7]
    print("\n===== Slide 8 (Phase 3 页) =====")
    fixed8, unm8 = fix_links_in_slide(prs, slide8, LINK_MAP, "S8")
    if unm8:
        print("  S8 未匹配(保留原链接): %r" % unm8)
    cleanup_dangling_rels(slide8, "S8")

    prs.save(DST)
    print("\n已保存: %s" % DST)
    print("修复统计: Slide7=%d, Slide8=%d" % (fixed7, fixed8))


if __name__ == "__main__":
    main()
