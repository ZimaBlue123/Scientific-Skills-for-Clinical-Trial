"""
CpG佐剂预防性疫苗安全性汇总 Word→PPT 生成器 (v1.0)
- 源文件: review_materials/CpG_Vaccine_Safety_Summary-V5-20260820.docx (单总表 5 列)
- 策略: "一页一苗" — 每页详细展示 1 款疫苗或 1 个主题块
- 版式: 16:9 宽屏, 纯白背景 #FFFFFF, 暗红强调 #C00000
- 内容: 无损保留 5 列全部文字信息, 安全性列按原文层次化排版
"""

import re

from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---------------- 常量 ----------------
SRC = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\CpG_Vaccine_Safety_Summary-V5-20260820.docx"
OUT = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\CpG_Vaccine_Safety_Summary_PPT_20260820.pptx"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x00, 0x00)  # 暗红
RED_LT = RGBColor(0xFB, 0xEA, 0xEA)  # 淡红底
DARK = RGBColor(0x33, 0x33, 0x33)  # 正文深灰
GRAY = RGBColor(0x7F, 0x7F, 0x7F)  # 次要灰
LGRAY = RGBColor(0xF4, 0xF4, 0xF4)  # 卡片底
BORDER = RGBColor(0xDD, 0xDD, 0xDD)
LINKC = RGBColor(0x1F, 0x4E, 0x79)  # 链接深蓝

FONT = "Microsoft YaHei"
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
M = Inches(0.35)  # 页面边距
TITLE_H = Inches(0.62)  # 标题栏高
PILL_H = Inches(0.52)  # 概览胶囊行高
BODY_Y = Inches(1.66)  # 主体区上缘
BODY_H = Inches(5.12)  # 主体区高
REF_Y = Inches(6.88)  # 参考文献条上缘
REF_H = Inches(0.46)  # 参考文献条高


# ---------------- 工具函数 ----------------
def set_font(run, size=10, bold=False, color=DARK, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", FONT)


def add_box(slide, x, y, w, h, fill=LGRAY, line=BORDER, radius=0.045, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, runs, first=False, align=PP_ALIGN.LEFT, line=1.06, before=0, after=0):
    """runs: list of (text, dict(size=,bold=,color=,italic=,link=))"""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line
    if before:
        p.space_before = Pt(before)
    if after:
        p.space_after = Pt(after)
    for txt, kw in runs:
        r = p.add_run()
        r.text = txt
        set_font(
            r,
            kw.get("size", 10),
            kw.get("bold", False),
            kw.get("color", DARK),
            kw.get("italic", False),
        )
        if kw.get("link"):
            r.hyperlink.address = kw["link"]
    return p


def norm_lines(s):
    return [ln.strip() for ln in s.split("\n") if ln.strip()]


# ---------------- 数据读取 ----------------
doc = Document(SRC)
title = doc.paragraphs[0].text.strip()
intro = [p.text for p in doc.paragraphs[1:6] if p.text.strip()]
tb = doc.tables[0]
rows = []
for row in tb.rows[1:]:
    rows.append([c.text for c in row.cells])


def parse_field(col2, key):
    """从"键: 值"行提取值"""
    for ln in norm_lines(col2):
        if ln.startswith(key + ":"):
            return ln[len(key) + 1 :].strip()
    return ""


def split_vaccine(c0):
    """列1 -> (主名, 佐剂说明, 适应症行列表)"""
    lines = norm_lines(c0)
    main = lines[0] if lines else ""
    sub = ""
    idx = 1
    if idx < len(lines) and lines[idx].startswith("("):
        sub = lines[idx]
        idx += 1
    while idx < len(lines) and not lines[idx].startswith("适应症"):
        sub = (sub + " " + lines[idx]).strip() if sub else lines[idx]
        idx += 1
    ind = [ln for ln in lines[idx:] if ln.strip()]
    return main, sub, ind


# ---------------- PPT 构建 ----------------
prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def add_header(slide, no, main, sub, ind_lines, page_total):
    """标题栏: 序号徽标 + 疫苗名 + 佐剂说明 + 右侧适应症"""
    # 顶部暗红标题栏
    add_box(
        slide,
        M,
        Inches(0.24),
        SLIDE_W - 2 * M,
        TITLE_H,
        fill=RED,
        line=None,
        radius=0.09,
    )
    # 序号徽标
    bad = add_box(
        slide,
        M + Inches(0.14),
        Inches(0.34),
        Inches(0.42),
        Inches(0.42),
        fill=WHITE,
        line=None,
        radius=0.5,
    )
    tf = bad.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{no:02d}"
    set_font(r, 15, True, RED)
    # 疫苗主名 + 佐剂
    tf = add_text(
        slide,
        M + Inches(0.70),
        Inches(0.24),
        Inches(8.0),
        TITLE_H,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    p = para(tf, [(main, dict(size=20, bold=True, color=WHITE))], first=True, line=1.0)
    if sub:
        p = para(
            tf,
            [(sub, dict(size=11, bold=False, color=RGBColor(0xFF, 0xD9, 0xD9)))],
            line=1.0,
        )
    # 右侧适应症
    tf = add_text(
        slide,
        M + Inches(8.8),
        Inches(0.24),
        Inches(4.0),
        TITLE_H,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    for i, ln in enumerate(ind_lines):
        p = para(
            tf,
            [(ln, dict(size=10.5, bold=(i == 0), color=WHITE))],
            first=(i == 0),
            align=PP_ALIGN.RIGHT,
            line=1.05,
        )
    # 右上角页码
    tf = add_text(
        slide, SLIDE_W - M - Inches(0.9), Inches(0.78), Inches(0.9), Inches(0.3)
    )
    para(
        tf,
        [(f"{no} / {page_total}", dict(size=9, color=GRAY))],
        first=True,
        align=PP_ALIGN.RIGHT,
    )


def add_pill(slide, items, x, y, w, h):
    """概览胶囊行: items = [(标签, 值)]"""
    n = len(items)
    gap = Inches(0.14)
    cw = (w - gap * (n - 1)) / n
    for i, (lab, val) in enumerate(items):
        cx = x + i * (cw + gap)
        add_box(slide, cx, y, cw, h, fill=WHITE, line=BORDER, radius=0.12)
        tf = add_text(
            slide,
            cx + Inches(0.12),
            y + Inches(0.06),
            cw - Inches(0.24),
            h - Inches(0.12),
            anchor=MSO_ANCHOR.MIDDLE,
        )
        para(
            tf,
            [
                (lab + "  ", dict(size=9, bold=True, color=RED)),
                (val, dict(size=10, bold=True, color=DARK)),
            ],
            first=True,
            line=1.0,
        )


def add_ref_bar(slide, refs, x, y, w, h):
    """底部参考文献条: 每组 ref 一行, 组内 ' | ' 分隔, PMID/DOI 附链接"""
    add_box(slide, x, y, w, h, fill=LGRAY, line=None, radius=0.10)
    tf = add_text(
        slide,
        x + Inches(0.16),
        y + Inches(0.04),
        w - Inches(0.32),
        h - Inches(0.08),
        anchor=MSO_ANCHOR.MIDDLE,
    )
    for i, ref in enumerate(refs):
        segs = split_ref(ref)
        runs = []
        if i == 0:
            runs.append(("核心参考文献  ", dict(size=9, bold=True, color=RED)))
        for j, (seg, link) in enumerate(segs):
            if j:
                runs.append((" | ", dict(size=9, color=GRAY)))
            runs.append((seg, dict(size=9, color=LINKC if link else DARK, link=link)))
        para(tf, runs, first=(i == 0), line=1.05)


def split_ref(ref):
    """把 'PMID: 123\\nDOI: xxx' 拆成 (显示文本, 链接) 段; 显示文本与原文字符完全一致"""
    out = []
    for ln in ref.split("\n"):
        ln = ln.strip()
        if not ln:
            continue
        m = re.match(r"^(PMID:\s*)(\d+)(.*)$", ln)
        if m:
            link = "https://pubmed.ncbi.nlm.nih.gov/" + m.group(2) + "/"
            out.append((m.group(1) + m.group(2), link))
            if m.group(3).strip():
                out.append((m.group(3), None))
            continue
        m = re.match(r"^(DOI:\s*)(\S+)(.*)$", ln)
        if m:
            out.append((m.group(1) + m.group(2), "https://doi.org/" + m.group(2)))
            if m.group(3).strip():
                out.append((m.group(3), None))
            continue
        out.append((ln, None))
    return out


# ---------- 页面 1: 封面 ----------
s = new_slide()
# 顶部暗红条
add_box(s, Inches(0), Inches(0), SLIDE_W, Inches(0.18), fill=RED, line=None, radius=0)
add_box(
    s, Inches(0), Inches(7.32), SLIDE_W, Inches(0.18), fill=RED, line=None, radius=0
)
# 封面标题: 保留原文完整文字
tf = add_text(s, M, Inches(1.15), SLIDE_W - 2 * M, Inches(1.6))
para(
    tf,
    [("CpG 佐剂预防性疫苗：", dict(size=38, bold=True, color=RED))],
    first=True,
    line=1.05,
)
para(
    tf,
    [
        (
            "核心临床试验与安全性数据汇总 (终局完善版)",
            dict(size=24, bold=True, color=DARK),
        )
    ],
    line=1.05,
)
para(
    tf,
    [
        (
            "FDA (ClinicalTrials.gov)  ×  NMPA (中国临床试验注册与信息公示平台)",
            dict(size=14, color=GRAY),
        )
    ],
    line=1.1,
    before=8,
)

# 标注说明卡片
add_box(
    s,
    M,
    Inches(3.30),
    SLIDE_W - 2 * M,
    Inches(2.60),
    fill=LGRAY,
    line=None,
    radius=0.03,
)
tf = add_text(
    s, M + Inches(0.30), Inches(3.52), SLIDE_W - 2 * M - Inches(0.60), Inches(2.2)
)
intro2 = intro[1] if len(intro) > 1 else ""
for i, ln in enumerate(norm_lines(intro2)):
    if i == 0:
        para(tf, [(ln, dict(size=12, bold=True, color=RED))], first=True, after=6)
    else:
        para(
            tf,
            [("▪ ", dict(size=10.5, color=RED)), (ln, dict(size=10.5, color=DARK))],
            line=1.12,
            after=3,
        )

tf = add_text(s, M, Inches(6.30), SLIDE_W - 2 * M, Inches(0.9))
para(
    tf,
    [
        (
            "转化自 V5 终局完善版  ·  一页一苗  ·  16:9 宽屏  ·  内容无损",
            dict(size=11.5, bold=True, color=DARK),
        )
    ],
    first=True,
    line=1.1,
)
para(tf, [("2026-08-20", dict(size=10.5, color=GRAY))], line=1.1)

# ---------- 页面 2: 总览 ----------
s = new_slide()
add_box(s, M, Inches(0.24), SLIDE_W - 2 * M, TITLE_H, fill=RED, line=None, radius=0.09)
tf = add_text(
    s, M + Inches(0.30), Inches(0.24), Inches(9.0), TITLE_H, anchor=MSO_ANCHOR.MIDDLE
)
para(
    tf,
    [
        (
            "内容总览 — 10 个主题块（9 款疫苗 + 1 个研发管线块）",
            dict(size=18, bold=True, color=WHITE),
        )
    ],
    first=True,
)

labels = ["疫苗名称", "适应症", "注册平台 & 编号", "数据状态"]
st = [
    "数据已公开",
    "数据已公开",
    "数据已公开",
    "数据已公开",
    "数据已公开",
    "数据已公开",
    "数据已公开",
    "数据已公开",
    "数据已公开",
    "数据未公开（进行中/审评）",
]
overview = []
for i, row in enumerate(rows):
    c0, c1 = row[0], row[1]
    main, sub, ind = split_vaccine(c0)
    plats = norm_lines(c1)
    plat = " / ".join(plats[:2])
    overview.append((main, ind[0] if ind else "", plat, st[i]))

tbl_shape = s.shapes.add_table(
    len(overview) + 1,
    4,
    M + Inches(0.30),
    Inches(1.30),
    SLIDE_W - 2 * M - Inches(0.60),
    Inches(5.30),
)
tbl = tbl_shape.table
tbl.columns[0].width = Inches(4.0)
tbl.columns[1].width = Inches(4.2)
tbl.columns[2].width = Inches(3.0)
tbl.columns[3].width = Inches(1.6)
for ci, lab in enumerate(labels):
    cell = tbl.cell(0, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RED
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(lab, dict(size=11, bold=True, color=WHITE))], first=True)
for ri, row in enumerate(overview, start=1):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if ri % 2 else LGRAY
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bold = ci == 0
        color = RED if (ci == 3 and "未公开" in val) else DARK
        para(
            tf,
            [
                (
                    str(ri) + ". " + val if ci == 0 else val,
                    dict(size=10, bold=bold, color=color),
                )
            ],
            first=True,
            line=1.0,
        )

# ---------- 内容页 (行1-9): 一页一苗 ----------
PAGE_TOTAL = 13
for idx, row in enumerate(rows[:9]):
    c0, c1, c2, c3, c4 = row
    s = new_slide()
    main, sub, ind = split_vaccine(c0)
    add_header(s, idx + 1, main, sub, ind, PAGE_TOTAL)

    # 概览胶囊: 注册平台 / 分期 / 样本量 / 免疫程序
    plats = norm_lines(c1)
    plat = plats[0] if plats else ""
    phase = parse_field(c2, "分期")
    size = parse_field(c2, "样本量")
    prog = parse_field(c2, "免疫程序")
    pill_items = [
        ("注册平台", plat),
        ("分期", phase),
        ("样本量", size),
        ("免疫程序", prog),
    ]
    add_pill(s, pill_items, M, Inches(0.98), SLIDE_W - 2 * M, PILL_H)

    # 左卡片: 注册平台&编号(完整) + 临床试验基本信息
    lw = Inches(4.30)
    add_box(s, M, BODY_Y, lw, BODY_H, fill=LGRAY, line=None, radius=0.03)
    tf = add_text(
        s,
        M + Inches(0.16),
        BODY_Y + Inches(0.10),
        lw - Inches(0.32),
        BODY_H - Inches(0.2),
    )
    para(
        tf,
        [("注册平台 & 编号", dict(size=12.5, bold=True, color=RED))],
        first=True,
        after=4,
    )
    for ln in norm_lines(c1):
        is_id = bool(re.match(r"^[A-Z]{2,4}\d+", ln))
        para(
            tf,
            [(ln, dict(size=10, bold=is_id, color=RED if is_id else DARK))],
            line=1.12,
            after=2,
        )
    para(
        tf,
        [("临床试验基本信息", dict(size=12.5, bold=True, color=RED))],
        before=8,
        after=5,
    )
    for ln in norm_lines(c2):
        m = re.match(
            r"^(分期|设计|样本量|适应人群|试验分组|免疫程序|研究终点|适应症):\s*(.*)$",
            ln,
        )
        if m:
            para(
                tf,
                [
                    (m.group(1) + ": ", dict(size=9.5, bold=True, color=RED)),
                    (m.group(2), dict(size=9.5, color=DARK)),
                ],
                line=1.12,
                after=3.5,
            )
        else:
            para(tf, [(ln, dict(size=9.5, color=DARK))], line=1.12, after=3.5)

    # 右卡片: 安全性数据汇总
    rw = Inches(8.18)
    rx = M + lw + Inches(0.15)
    add_box(s, rx, BODY_Y, rw, BODY_H, fill=WHITE, line=BORDER, radius=0.03)
    tf = add_text(
        s,
        rx + Inches(0.16),
        BODY_Y + Inches(0.10),
        rw - Inches(0.32),
        BODY_H - Inches(0.2),
    )
    para(
        tf,
        [
            ("安全性数据汇总  ", dict(size=12.5, bold=True, color=RED)),
            ("(逐项标注 AE/ADR)", dict(size=9, color=GRAY)),
        ],
        first=True,
        after=6,
    )
    safe_lines = norm_lines(c3)
    # 字号自适应
    n_heavy = sum(1 for ln in safe_lines if ln.startswith("■"))
    fsize = 10.5 if len(safe_lines) <= 12 else (10 if len(safe_lines) <= 16 else 9.5)
    for ln in safe_lines:
        if ln.startswith("【"):
            para(
                tf,
                [(ln, dict(size=fsize + 0.5, bold=True, color=RED))],
                line=1.08,
                before=5,
                after=3,
            )
        elif ln.startswith("■"):
            # 行内 (AE)/(ADR) 高亮, 文字与原文一致
            parts = re.split(r"(\(AE\)|\(ADR\)|\(SADR\))", ln)
            runs = []
            for seg in parts:
                if not seg:
                    continue
                if re.fullmatch(r"\(AE\)|\(ADR\)|\(SADR\)", seg):
                    runs.append((seg, dict(size=fsize, bold=True, color=RED)))
                else:
                    runs.append((seg, dict(size=fsize, color=DARK)))
            para(tf, runs, line=1.15, after=3)
        else:
            para(
                tf,
                [
                    ("   ", dict(size=fsize, color=GRAY)),
                    (ln, dict(size=fsize - 0.5, color=GRAY, italic=True)),
                ],
                line=1.12,
                after=2,
            )

    # 底部参考文献条
    refs = [r for r in c4.split("\n\n") if r.strip()]
    add_ref_bar(s, refs, M, REF_Y, SLIDE_W - 2 * M, REF_H)

# ---------- 页面 12: 管线块 (行10) ----------
row = rows[9]
c0, c1, c2, c3, c4 = row
s = new_slide()
main, sub, ind = split_vaccine(c0)
add_header(s, 10, main, sub, ind, PAGE_TOTAL)
add_pill(
    s,
    [
        ("状态", "数据未公开（进行中 / NMPA 审评审批阶段）"),
        ("覆盖", "带状疱疹、乙肝等新一代 CpG 佐剂疫苗"),
    ],
    M,
    Inches(0.98),
    SLIDE_W - 2 * M,
    PILL_H,
)

lw = Inches(7.60)
add_box(s, M, BODY_Y, lw, BODY_H, fill=LGRAY, line=None, radius=0.03)
tf = add_text(
    s, M + Inches(0.16), BODY_Y + Inches(0.10), lw - Inches(0.32), BODY_H - Inches(0.2)
)
para(
    tf,
    [("注册平台 & 编号", dict(size=12.5, bold=True, color=RED))],
    first=True,
    after=4,
)
for ln in norm_lines(c1):
    is_id = bool(re.match(r"^[A-Z]{2,4}\d+", ln))
    para(
        tf,
        [(ln, dict(size=10, bold=is_id, color=RED if is_id else DARK))],
        line=1.12,
        after=2,
    )
para(
    tf, [("研发管线项目一览", dict(size=12.5, bold=True, color=RED))], before=8, after=5
)
for ln in norm_lines(c2):
    if ln.startswith("■"):
        para(
            tf,
            [
                ("■ ", dict(size=9.5, bold=True, color=RED)),
                (ln[2:], dict(size=9.5, color=DARK)),
            ],
            line=1.15,
            after=4,
        )
    else:
        para(tf, [(ln, dict(size=9.5, color=DARK))], line=1.15, after=4)

rw = Inches(4.88)
rx = M + lw + Inches(0.15)
add_box(s, rx, BODY_Y, rw, BODY_H, fill=WHITE, line=BORDER, radius=0.03)
tf = add_text(
    s, rx + Inches(0.14), BODY_Y + Inches(0.10), rw - Inches(0.28), BODY_H - Inches(0.2)
)
first_done = [False]
for ln in norm_lines(c3):
    if ln.startswith("【"):
        para(
            tf,
            [(ln, dict(size=11.5, bold=True, color=RED))],
            first=not first_done[0],
            before=0 if not first_done[0] else 6,
            after=4,
        )
    elif ln.startswith("■"):
        para(
            tf,
            [
                ("■ ", dict(size=9.5, bold=True, color=RED)),
                (ln[2:], dict(size=9.5, color=DARK)),
            ],
            first=not first_done[0],
            line=1.15,
            after=4,
        )
    else:
        para(
            tf,
            [(ln, dict(size=9.5, color=DARK))],
            first=not first_done[0],
            line=1.15,
            after=4,
        )
    first_done[0] = True
para(tf, [("数据来源", dict(size=11.5, bold=True, color=RED))], before=8, after=3)
for ln in norm_lines(c4):
    para(tf, [(ln, dict(size=9, color=GRAY))], line=1.12, after=1)

add_ref_bar(
    s,
    ["数据来源: ClinicalTrials.gov 及 NMPA/CDE 最新公示数据"],
    M,
    REF_Y,
    SLIDE_W - 2 * M,
    REF_H,
)

# ---------- 页面 13: 尾页 ----------
s = new_slide()
add_box(s, Inches(0), Inches(0), SLIDE_W, Inches(0.18), fill=RED, line=None, radius=0)
add_box(
    s, Inches(0), Inches(7.32), SLIDE_W, Inches(0.18), fill=RED, line=None, radius=0
)
tf = add_text(s, M, Inches(2.2), SLIDE_W - 2 * M, Inches(1.2))
para(
    tf,
    [("数据完整性与可追溯性说明", dict(size=24, bold=True, color=RED))],
    first=True,
    line=1.1,
)
para(
    tf,
    [
        (
            "全部安全性数值均标注 (AE) 或 (ADR)，PMID / DOI 已附直达链接，可回溯至原始文献。",
            dict(size=13, color=DARK),
        )
    ],
    line=1.2,
    before=6,
)

add_box(
    s, M, Inches(3.70), SLIDE_W - 2 * M, Inches(2.0), fill=LGRAY, line=None, radius=0.03
)
tf = add_text(
    s, M + Inches(0.30), Inches(3.92), SLIDE_W - 2 * M - Inches(0.60), Inches(1.6)
)
for ln in norm_lines(intro[0])[:2]:
    para(
        tf,
        [("▪ ", dict(size=10.5, color=RED)), (ln, dict(size=10.5, color=DARK))],
        first=(ln == norm_lines(intro[0])[0]),
        line=1.15,
        after=4,
    )
    if ln == norm_lines(intro[0])[0]:
        break
tf = add_text(s, M, Inches(6.1), SLIDE_W - 2 * M, Inches(0.8))
para(
    tf,
    [("谢谢观看  ·  基于 V5 终局完善版转制", dict(size=15, bold=True, color=DARK))],
    first=True,
    line=1.1,
)
para(
    tf,
    [
        (
            "2026-08-20  ·  一页一苗 · 内容无损  ·  白底 #FFFFFF / 暗红 #C00000",
            dict(size=10.5, color=GRAY),
        )
    ],
    line=1.1,
)

prs.save(OUT)
print(f"✅ PPT 已生成: {OUT}")
print(f"   页数: {len(prs.slides._sldIdLst)}")
