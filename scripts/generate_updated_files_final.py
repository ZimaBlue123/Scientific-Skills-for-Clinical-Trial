"""V29 生成脚本 — CpG佐剂预防性疫苗安全性汇总 (PPTX + DOCX, 双格式超链接)
流程: V13原文件docx -> V29docx (概览2表 + 详表5列表 + 总结)
      V29-temp.pptx (win32com骨架) -> V29pptx (概览2页 + 详表6页 + 总结页)
超链接: 登记号(NCT/CTR/CTRI/PACTR/jRCT/ChiCTR/ACTRN) + 文献(PMID/DOI) + FDA官方文件
V28 新增: 详表第一列图标 (🧑💼申办者 / 🎯适应症) + Study 1 三值归属说明 + PMID 23727002 标注修正
V29 变更: 申办者图标改 🏦 银行大楼; ACTRN 链接补 .aspx 修正; 移除动脉网路由; 新增 Uvax Bio 新闻稿路由(EATG全文)
"""

import copy
import os
import re

import docx.opc.constants
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.table import Table
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from update_data_v29 import (
    details_data,
    marketed_data,
    marketed_toc,
    pipeline_data,
    pipeline_toc,
    summary_blocks,
)

BASE = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial"
OUT_DOCX = os.path.join(BASE, r"review_materials\CpG_Vaccine_Safety_Summary-V29-20260827.docx")
OUT_PPTX = os.path.join(BASE, r"review_materials\CpG_Vaccine_Safety_Summary-V29-20260827.pptx")
TEMP_PPTX = os.path.join(
    BASE, r"review_materials\CpG_Vaccine_Safety_Summary-V29-20260827-temp.pptx"
)
BASE_DOCX = os.path.join(
    BASE, r"review_materials\CpG_Vaccine_Safety_Summary-V13-20260820-原文件.docx"
)

# ---------------- 超链接路由 (登记号 / 文献 / 官方文件) ----------------
LINK_PATTERN = (
    r"(PMID:\s*\d+|DOI:\s*\S+|NCT\d{8}|CTR\d{8}|CTRI/\S+|PACTR\d+"
    r"|jRCT\w+|ChiCTR\w+|ACTRN\d+|NMPA CDE[^\n,]*|FDA [^,\n\)]+|https?://\S+)"
)

# ChiCTR 登记号 -> 详情页链接 (2026-08-27 逐一核验; showproj 为当前版详情页, hvshowproject 为历史版)
CHICTR_LINKS = {
    "ChiCTR2600119810": "https://www.chictr.org.cn/showproj.html?proj=311880",  # 华普 HP-2001 带状疱疹 I期
    "ChiCTR2600118487": "https://www.chictr.org.cn/showproj.html?proj=301189",  # 简达带状疱疹 I期
    "ChiCTR2500108408": "https://www.chictr.org.cn/showproj.html?proj=280568",  # 远大 TVAX-009 乙肝 II期
}


def resolve_link(part, product_text=""):
    """返回该文本片段对应的超链接 URL, 无则 None (仅保留可验证直链)
    已核验(2026-08-27): FDA media 4条全通; NCT/PMID/DOI/CTRI/PACTR/ACTRN 官方平台可直链。
    ChiCTR 仅对已核验到详情页(showproj)的 3 个登记号嵌链接; CTR/jRCT 无稳定可直达具体试验页 -> 仅文字。
    """
    m = re.match(r"PMID:\s*(\d+)", part)
    if m:
        return f"https://pubmed.ncbi.nlm.nih.gov/{m.group(1)}/"
    m = re.match(r"DOI:\s*(\S+)", part)
    if m:
        return f"https://doi.org/{m.group(1)}"
    m = re.match(r"(NCT\d{8})", part)
    if m:
        return f"https://clinicaltrials.gov/study/{m.group(1)}"
    m = re.match(r"(CTRI/\S+)", part)
    if m:
        return f"https://trialsearch.who.int/Trial2.aspx?TrialID={m.group(1)}"
    m = re.match(r"(PACTR\d+)", part)
    if m:
        return f"https://trialsearch.who.int/Trial2.aspx?TrialID={m.group(1)}"
    m = re.match(r"ACTRN(\d+)", part)
    if m:
        return f"https://anzctr.org.au/ACTRN{m.group(1)}.aspx"
    m = re.match(r"(ChiCTR\w+)", part)
    if m:
        return CHICTR_LINKS.get(m.group(1))
    if re.match(r"https?://\S+", part):
        return part
    if part.startswith("Uvax Bio新闻稿"):
        return (
            "https://www.eatg.org/?p=20095"  # Uvax Bio 2024-11-19 中期分析#1 新闻稿全文(EATG转载)
        )
    if part.startswith("FDA") and ("Package Insert" in part or "Clinical Review" in part):
        if "HEPLISAV" in product_text:
            return (
                "https://www.fda.gov/media/109808/download"
                if "Review" in part
                else "https://www.fda.gov/media/108745/download"
            )
        if "CYFENDUS" in product_text or "AV7909" in product_text:
            return (
                "https://www.fda.gov/media/171141/download"
                if "Review" in part
                else "https://www.fda.gov/media/189503/download"
            )
        return "https://www.fda.gov/vaccines-blood-biologics/vaccines"
    # CTR/jRCT/NMPA CDE/企业公告等: 无稳定具体页直链 -> 不嵌入链接(仅文字)
    return None


# ---------------- DOCX 部分 ----------------
def add_docx_hyperlink(paragraph, text, url):
    part = paragraph.part
    r_id = part.relate_to(url, docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    new_run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    rFonts = OxmlElement("w:rFonts")
    rFonts.set(qn("w:ascii"), "微软雅黑")
    rFonts.set(qn("w:eastAsia"), "微软雅黑")
    rPr.append(rFonts)
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "0563C1")
    rPr.append(color)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rPr.append(u)
    new_run.append(rPr)
    t = OxmlElement("w:t")
    t.text = text
    t.set(qn("xml:space"), "preserve")
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)


ICON_SPONSOR = "🏦 "  # 申办者: 蓝色银行大楼(用户指定 2026-08-27)
ICON_INDICATION = "🎯 "  # 适应症: 箭靶


def apply_col0_icons(lines):
    """详表第一列: 倒数第2行(申办者)前加小人图标, 最后一行(适应症)前加箭靶图标。
    兼容 4 行(名称/分期/申办者/适应症)与 3 行(CYFENDUS/IndoVac, 名称/申办者/适应症)两种结构。
    """
    if len(lines) >= 2:
        lines[-2] = ICON_SPONSOR + lines[-2]
        lines[-1] = ICON_INDICATION + lines[-1]
    return lines


def set_docx_cell_rich_text(cell, text, product_text="", first_col_icons=False):
    """富文本写入 docx 单元格: 按行分段, 登记号/文献自动嵌入超链接
    first_col_icons=True 时在首列加申办者/适应症彩色图标(仅详表数据行启用)"""
    for p_extra in cell.paragraphs[1:]:
        p_extra._element.getparent().remove(p_extra._element)
    p = cell.paragraphs[0]
    for r in list(p.runs):
        r._element.getparent().remove(r._element)
    lines = text.split("\n")
    if first_col_icons:
        lines = apply_col0_icons(lines)
    paras = [p]
    for _ in lines[1:]:
        paras.append(cell.add_paragraph())
    for para, line in zip(paras, lines):
        parts = re.split(LINK_PATTERN, line)
        for part in parts:
            if not part:
                continue
            url = resolve_link(part, product_text)
            if url:
                add_docx_hyperlink(para, part, url)
            else:
                para.add_run(part)


def set_docx_cell_text(cell, text):
    if not cell.paragraphs:
        cell.add_paragraph()
    p = cell.paragraphs[0]
    style_run = p.runs[0] if p.runs else p.add_run()
    for r in p.runs[1:]:
        r._element.getparent().remove(r._element)
    style_run.text = text
    for p_extra in cell.paragraphs[1:]:
        p_extra._element.getparent().remove(p_extra._element)


def add_col_docx(t, doc):
    grid = t._tbl.find(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tblGrid")
    if grid is not None:
        cols = grid.findall(
            ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}gridCol"
        )
        if cols:
            grid.append(copy.deepcopy(cols[-1]))
    for row in t.rows:
        new_tc = copy.deepcopy(row.cells[-1]._tc)
        row._tr.append(new_tc)
    return Table(t._tbl, doc)


def adjust_docx_table(doc, t, data_len):
    while len(t.rows) < data_len + 1:
        new_tr = copy.deepcopy(t.rows[-1]._tr)
        t._tbl.append(new_tr)
        t = Table(t._tbl, doc)
    while len(t.rows) > data_len + 1:
        t._tbl.remove(t.rows[-1]._tr)
        t = Table(t._tbl, doc)
    return t


def process_docx():
    doc = Document(BASE_DOCX)

    # 标题: (V13) -> (V29)
    if doc.paragraphs and "(V13)" in doc.paragraphs[0].text:
        _set_para_text(doc.paragraphs[0], doc.paragraphs[0].text.replace("(V13)", "(V29)"))

    # 概览标题段
    for p in doc.paragraphs:
        if "【概览】" in p.text:
            _set_para_text(p, "【概览】CpG 预防性疫苗核心管线目录\n表1：已上市的产品")
            break

    # ---- 表1: 已上市概览 ----
    t1 = doc.tables[0]
    t1 = adjust_docx_table(doc, t1, len(marketed_toc))
    for c, h in enumerate(["疫苗名称", "适应症", "申办者", "在研阶段/状态"]):
        set_docx_cell_text(t1.rows[0].cells[c], h)
    for r, row_data in enumerate(marketed_toc):
        for c, text in enumerate(row_data):
            set_docx_cell_text(t1.rows[r + 1].cells[c], text)

    # ---- 表2: 在研概览 (复制表1结构, 紧随其后插入) ----
    new_tbl = copy.deepcopy(t1._tbl)  # 干净 4 列模板
    label_p = OxmlElement("w:p")
    label_r = OxmlElement("w:r")
    label_t = OxmlElement("w:t")
    label_t.text = "表2：在研的产品"
    label_r.append(label_t)
    label_p.append(label_r)
    t1._tbl.addnext(label_p)
    label_p.addnext(new_tbl)
    t2 = Table(new_tbl, doc)
    t2 = adjust_docx_table(doc, t2, len(pipeline_toc))
    for c, h in enumerate(["疫苗名称", "适应症", "申办者", "在研阶段/状态"]):
        set_docx_cell_text(t2.rows[0].cells[c], h)
    for r, row_data in enumerate(pipeline_toc):
        for c, text in enumerate(row_data):
            set_docx_cell_text(t2.rows[r + 1].cells[c], text)

    # ---- 表3: 详表 (原 V13 详表, 5 列, 富文本 + 超链接) ----
    t3 = doc.tables[2]
    assert t3.rows[0].cells[0].text.strip().startswith("疫苗信息"), (
        f"表3 表头异常: {t3.rows[0].cells[0].text[:20]}"
    )
    if len(t3.rows[0].cells) < 5:
        t3 = add_col_docx(t3, doc)
    t3 = adjust_docx_table(doc, t3, len(details_data))
    det_headers = ["疫苗信息", "在研阶段或状态", "临床试验基本信息", "安全性数据汇总", "参考信息"]
    for c, h in enumerate(det_headers):
        set_docx_cell_text(t3.rows[0].cells[c], h)
    for r, row_data in enumerate(details_data):
        for c, text in enumerate(row_data):
            set_docx_cell_rich_text(
                t3.rows[r + 1].cells[c], text, product_text=row_data[0], first_col_icons=(c == 0)
            )

    # ---- 总结段 ----
    summary_text = "\n\n".join(
        [summary_blocks["1"], summary_blocks["2"], summary_blocks["3"], summary_blocks["4"]]
    )
    for p in doc.paragraphs:
        if p.text.strip().startswith("通过汇总"):
            _set_para_text(p, summary_text)
            break

    doc.save(OUT_DOCX)
    print(f"DOCX saved: {OUT_DOCX}")


def _set_para_text(p, text):
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    for r in p.runs[1:]:
        r._element.getparent().remove(r._element)


# ---------------- PPTX 部分 ----------------
def set_pptx_cell_rich_text(cell, text, font_size_pt=None, product_text="", first_col_icons=False):
    tf = cell.text_frame
    tf.word_wrap = True
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p_elem = tf.paragraphs[i]._p
        p_elem.getparent().remove(p_elem)

    lines = text.split("\n")
    if first_col_icons:
        lines = apply_col0_icons(lines)
    for l_idx, line in enumerate(lines):
        if l_idx == 0:
            p = tf.paragraphs[0]
            for r_elem in list(p._p.r_lst):
                p._p.remove(r_elem)
        else:
            p = tf.add_paragraph()

        parts = re.split(LINK_PATTERN, line)
        for part in parts:
            if not part:
                continue
            r = p.add_run()
            r.text = part
            if font_size_pt:
                r.font.size = Pt(font_size_pt)
            url = resolve_link(part, product_text)
            if url:
                r.hyperlink.address = url


def add_pptx_row(table):
    new_tr = copy.deepcopy(table._tbl.tr_lst[-1])
    for tc in new_tr.tc_lst:
        if tc.txBody is not None:
            for p in tc.txBody.p_lst:
                for r in p.r_lst:
                    r.text = ""
    table._tbl.append(new_tr)


def get_pptx_cell(table, r, c):
    from pptx.table import _Cell

    return _Cell(table._tbl.tr_lst[r].tc_lst[c], table)


def _style_title(shape, text, size_pt=20):
    # 规范化文本(合并换行/空白), 避免模板首行为空行导致样式落空
    text = " ".join(text.split())
    shape.text_frame.text = text
    tf = shape.text_frame
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p_elem = tf.paragraphs[i]._p
        p_elem.getparent().remove(p_elem)
    for para in tf.paragraphs:
        for r in para.runs:
            r.font.color.rgb = RGBColor(192, 0, 0)
            r.font.name = "微软雅黑"
            r.font.bold = True
            r.font.size = Pt(size_pt)


def apply_title_style(slide, text, keyword="概览", size_pt=20):
    for s in slide.shapes:
        if s.has_text_frame and s.text and keyword in s.text:
            _style_title(s, text, size_pt)
            break


def add_section_label(slide, text, top_in, left_in=0.5):
    txBox = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(4), Inches(0.4))
    txBox.text_frame.text = text
    para = txBox.text_frame.paragraphs[0]
    para.font.size = Pt(12)
    para.font.bold = True
    para.font.color.rgb = RGBColor(192, 0, 0)


def fill_detail_slide(slide, chunk, title_text, font_pt=10):
    # 大标题统一 20pt 暗红加粗(与概览页一致)
    apply_title_style(slide, title_text, keyword="安全性汇总", size_pt=20)
    t_det = None
    for s in slide.shapes:
        if s.has_table:
            t_det = s.table
            s.top = int(Inches(1.0))  # 标题下方起始
            s.height = int(Inches(6.2))  # 可用高度(16:9 页高7.5")
            break
    if t_det is None:
        raise RuntimeError("detail slide has no table")

    headers = ["疫苗信息", "在研阶段或状态", "临床试验基本信息", "安全性数据汇总", "参考信息"]
    for c, h in enumerate(headers):
        set_pptx_cell_rich_text(get_pptx_cell(t_det, 0, c), h, font_size_pt=10.5)

    while len(t_det.rows) < len(chunk) + 1:
        add_pptx_row(t_det)
    while len(t_det.rows) > len(chunk) + 1:
        t_det._tbl.remove(t_det._tbl.tr_lst[-1])

    for r, row_data in enumerate(chunk):
        for c, text in enumerate(row_data):
            set_pptx_cell_rich_text(
                get_pptx_cell(t_det, r + 1, c),
                text,
                font_size_pt=font_pt,
                product_text=row_data[0],
                first_col_icons=(c == 0),
            )


def process_pptx():
    prs = Presentation(TEMP_PPTX)

    # ---- 概览页 1 (已上市 + 在研前半) ----
    s1 = prs.slides[1]
    apply_title_style(s1, "【概览】CpG预防性疫苗核心管线目录")

    tables_s1 = [s for s in s1.shapes if s.has_table]
    add_section_label(s1, "■ 已上市产品：", 0.8)
    add_section_label(s1, "■ 在研产品：", 3.9)

    t_marketed = tables_s1[0].table
    tables_s1[0].top = int(Inches(1.2))
    tables_s1[0].height = int(Inches(2.5))

    t_pipe1 = tables_s1[1].table
    tables_s1[1].top = int(Inches(4.3))
    tables_s1[1].height = int(Inches(2.5))

    while len(t_marketed.rows) < len(marketed_toc) + 1:
        add_pptx_row(t_marketed)
    while len(t_marketed.rows) > len(marketed_toc) + 1:
        t_marketed._tbl.remove(t_marketed._tbl.tr_lst[-1])
    for c, h in enumerate(["疫苗名称", "适应症", "申办者", "在研阶段/状态"]):
        set_pptx_cell_rich_text(get_pptx_cell(t_marketed, 0, c), h, font_size_pt=10.5)
    for r, row_data in enumerate(marketed_toc):
        for c, text in enumerate(row_data):
            set_pptx_cell_rich_text(get_pptx_cell(t_marketed, r + 1, c), text, font_size_pt=10)

    pipe_split = 6
    pipe1_data = pipeline_toc[:pipe_split]
    pipe2_data = pipeline_toc[pipe_split:]

    while len(t_pipe1.rows) < len(pipe1_data) + 1:
        add_pptx_row(t_pipe1)
    while len(t_pipe1.rows) > len(pipe1_data) + 1:
        t_pipe1._tbl.remove(t_pipe1._tbl.tr_lst[-1])
    for c, h in enumerate(["疫苗名称", "适应症", "申办者", "在研阶段/状态"]):
        set_pptx_cell_rich_text(get_pptx_cell(t_pipe1, 0, c), h, font_size_pt=10.5)
    for r, row_data in enumerate(pipe1_data):
        for c, text in enumerate(row_data):
            set_pptx_cell_rich_text(get_pptx_cell(t_pipe1, r + 1, c), text, font_size_pt=10)

    # ---- 概览页 2 (在研续) ----
    s2 = prs.slides[2]
    apply_title_style(s2, "【概览】CpG预防性疫苗核心管线目录")
    add_section_label(s2, "■ 在研产品（续）：", 0.8)

    t_pipe2 = None
    for s in s2.shapes:
        if s.has_table:
            t_pipe2 = s.table
            s.top = int(Inches(1.2))
            s.height = int(Inches(5.5))
            break
    while len(t_pipe2.rows) < len(pipe2_data) + 1:
        add_pptx_row(t_pipe2)
    while len(t_pipe2.rows) > len(pipe2_data) + 1:
        t_pipe2._tbl.remove(t_pipe2._tbl.tr_lst[-1])
    for c, h in enumerate(["疫苗名称", "适应症", "申办者", "在研阶段/状态"]):
        set_pptx_cell_rich_text(get_pptx_cell(t_pipe2, 0, c), h, font_size_pt=10.5)
    for r, row_data in enumerate(pipe2_data):
        for c, text in enumerate(row_data):
            set_pptx_cell_rich_text(get_pptx_cell(t_pipe2, r + 1, c), text, font_size_pt=10)

    # ---- 详表页: 已上市板块 (8条 -> 2页: 第4页仅3行HEPLISAV-B, 第5页5行=CYFENDUS起) ----
    detail_start = 3
    slide_idx = detail_start
    marketed_chunks = [marketed_data[:3], marketed_data[3:]]
    for ci, chunk in enumerate(marketed_chunks):
        title = "【已上市产品】CpG佐剂预防性疫苗安全性汇总" + ("（续）" if ci > 0 else "")
        fill_detail_slide(prs.slides[slide_idx], chunk, title)
        slide_idx += 1

    # ---- 详表页: 在研板块 (16条, 4行/页 -> 4页) ----
    pipeline_chunks = [pipeline_data[i : i + 4] for i in range(0, len(pipeline_data), 4)]
    for ci, chunk in enumerate(pipeline_chunks):
        title = "【在研产品】CpG佐剂预防性疫苗安全性汇总" + ("（续）" if ci > 0 else "")
        fill_detail_slide(prs.slides[slide_idx], chunk, title)
        slide_idx += 1

    # ---- 总结页 ----
    s_sum = prs.slides[len(prs.slides) - 1]
    summary_text = "\n\n".join(
        [summary_blocks["1"], summary_blocks["2"], summary_blocks["3"], summary_blocks["4"]]
    )
    # 大标题统一 20pt 暗红加粗(与其余页一致)
    for s in s_sum.shapes:
        if s.has_text_frame and s.text and "总体安全性总结" in s.text:
            _style_title(s, s.text, 20)
            break
    for s in s_sum.shapes:
        if s.has_text_frame and s.text and "1. " in s.text:
            s.text = summary_text
            break

    prs.save(OUT_PPTX)
    print(f"PPTX saved: {OUT_PPTX} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    process_docx()
    process_pptx()
