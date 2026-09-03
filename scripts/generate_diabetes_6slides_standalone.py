# -*- coding: utf-8 -*-
"""
TVAX-009 糖尿病亚组免疫原性 6 页 —— 独立全新 PPT 生成器（无克隆 / 无 deepcopy）。

背景：此前用 deepcopy 克隆参考页导致 a16:creationId GUID 跨页重复，PPT 打不开。
本脚本从零构建全新 Presentation，表格用「新生成的 XML」直接嵌入 graphicFrame，
不复制任何现有形状，从根上规避 XML 腐坏问题。

输出：只含 6 页（阳转率/GMC × PPS/FAS × 绝对/相对时间点），可正常打开。
"""
import json
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BASE = 'E:/Cursor Project/2-Scientific-Skills-for-Clinical_Trial/review_materials'
OUT = BASE + '/TVAX-009项目3期临床试验启动前沟通交流ppt-糖尿病亚组6页-20260903.pptx'

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

with open(BASE + '/diabetes_immuno_result.json', encoding='utf-8') as f:
    DATA = json.load(f)

FAS = DATA['fas']
PPS = DATA['pps']

# 方案 A：7 组全列（0,1月 A1/A2；0,2月 B1/B2；0,1,6月 C1/C3/C2）
GROUPS = [
    ('A1', '低剂量组(A1)', '0,1月程序', None),
    ('A2', '高剂量组(A2)', '0,1月程序', None),
    ('B1', '低剂量组(B1)', '0,2月程序', None),
    ('B2', '高剂量组(B2)', '0,2月程序', None),
    ('C1', '阳性对照组(C1)', '0,1,6月程序', '18-59岁'),
    ('C3', '0,1,6月高剂量组(C3)', '0,1,6月程序', '≥60岁'),
    ('C2', '阳性对照组(C2)', '0,1,6月程序', '≥60岁'),
]
PROGRAMS = [
    ('0,1月程序', ['A1', 'A2']),
    ('0,2月程序', ['B1', 'B2']),
    ('0,1,6月程序', ['C1', 'C3', 'C2']),
]
PROG_OF = {code: prog for code, _, prog, _ in GROUPS}
LABEL = {code: lab for code, lab, _, _ in GROUPS}
ANNO = {code: a for code, _, _, a in GROUPS}

PROG_TIMEPOINTS = {
    '0,1月程序': ['M1', 'M2', 'M3', 'M6', 'M7', 'M8'],
    '0,2月程序': ['M1', 'M2', 'M3', 'M4', 'M7', 'M8'],
    '0,1,6月程序': ['M1', 'M2', 'M6', 'M7', 'M8'],
}
ALL_TP = ['M1', 'M2', 'M3', 'M4', 'M6', 'M7', 'M8']

REL_MAP = {
    '全免后1个月': {'0,1月程序': 'M2', '0,2月程序': 'M3', '0,1,6月程序': 'M7'},
    '全免后2个月': {'0,1月程序': 'M3', '0,2月程序': 'M4', '0,1,6月程序': 'M8'},
}

# 表头填充（浅蓝灰），数据格保持白色
HDR_FILL = 'D9E2F3'


# ---------------- XML 片段 ----------------
def _esc(s):
    return s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')


LN_T = ('<a:lnT w="12700" cap="flat" cmpd="sng" algn="ctr">'
        '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
        '<a:prstDash val="solid"/><a:round/>'
        '<a:headEnd type="none" w="med" len="med"/>'
        '<a:tailEnd type="none" w="med" len="med"/></a:lnT>')


def _ln_b(solid):
    if solid:
        return ('<a:lnB w="12700" cap="flat" cmpd="sng" algn="ctr">'
                '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
                '<a:prstDash val="solid"/><a:round/>'
                '<a:headEnd type="none" w="med" len="med"/>'
                '<a:tailEnd type="none" w="med" len="med"/></a:lnB>')
    return '<a:lnB><a:noFill/></a:lnB>'


def make_tc(text, bold=False, grid_span=None, row_span=None, lnB_solid=False, lnT=True, fill=None, font_size=800):
    attrs = ''
    if grid_span:
        attrs += ' gridSpan="%d"' % grid_span
    if row_span:
        attrs += ' rowSpan="%d"' % row_span
    b = ' b="1"' if bold else ''
    s = '<a:tc xmlns:a="%s"%s><a:txBody><a:bodyPr/><a:lstStyle/>' % (A_NS, attrs)
    for line in text.split('\n'):
        s += ('<a:p><a:pPr indent="0" algn="ctr">'
              '<a:lnSpc><a:spcPct val="100000"/></a:lnSpc><a:buNone/></a:pPr>'
              '<a:r><a:rPr lang="zh-CN" altLang="en-US" sz="%d" kern="0"%s>'
              '<a:solidFill><a:schemeClr val="tx1"/></a:solidFill><a:effectLst/>'
              '<a:latin typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
              '<a:ea typeface="微软雅黑" panose="020B0503020204020204" pitchFamily="34" charset="-122"/>'
              '<a:cs typeface="Times New Roman" panose="02020603050405020304" charset="0"/>'
              '<a:sym typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
              '</a:rPr><a:t>%s</a:t></a:r></a:p>' % (font_size, b, _esc(line)))
    s += '</a:txBody>'
    s += '<a:tcPr marL="36195" marR="0" marT="0" marB="0" anchor="ctr">'
    s += '<a:lnL><a:noFill/></a:lnL><a:lnR><a:noFill/></a:lnR>'
    if lnT:
        s += LN_T
    s += _ln_b(lnB_solid)
    if fill:
        s += '<a:solidFill><a:srgbClr val="%s"/></a:solidFill>' % fill
    else:
        s += '<a:noFill/>'
    s += '</a:tcPr></a:tc>'
    return s


def table_str(col_widths, rows):
    s = '<a:tbl xmlns:a="%s"><a:tblPr firstRow="0" firstCol="0" bandRow="0"/><a:tblGrid>' % A_NS
    for w in col_widths:
        s += '<a:gridCol w="%d"/>' % w
    s += '</a:tblGrid>'
    total_h = 0
    for h, cells in rows:
        s += '<a:tr h="%d">%s</a:tr>' % (h, ''.join(cells))
        total_h += h
    s += '</a:tbl>'
    return s, total_h


# ---------------- 取值 ----------------
def fmt_num(v):
    return '/' if v is None else '%.2f' % v


def get_vals(data, code, tag, tp, mode):
    rec = data[code][tag][tp]
    if mode == 'sero':
        return [str(rec['N_阳转']), str(rec['n_阳转']), fmt_num(rec['pct'])]
    return [str(rec['N_gmc']), fmt_num(rec['gmc'])]


# ---------------- 表格构建 ----------------
def build_absolute(data, mode):
    """绝对时间点表：19 行（有/无 × M1-M8，含空行分隔）。返回 (xml, width, height)。"""
    if mode == 'sero':
        lead_w = [504000, 612000]
        sub_w = [310000, 310000, 410000]
        sub_headers = ['N', 'n', '%']
        header_h = 190500
        you_h = 175260
        wu_h = 165100
        blank_h = [37669, 64303]
        fz = 800
    else:
        lead_w = [512551, 612000]
        sub_w = [450000, 650000]
        sub_headers = ['例数', 'GMC']
        header_h = 214250
        you_h = 169564
        wu_h = 169258
        blank_h = [36000, 36000]
        fz = 700

    lead_headers = ['糖尿病病史', '时间点\n(首剂接种后）']
    col_widths = list(lead_w)
    for _, codes in PROGRAMS:
        for _ in codes:
            col_widths.extend(sub_w)
    n_sub = sum(len(codes) * len(sub_w) for _, codes in PROGRAMS)

    rows = []
    # R0：前导列(rowSpan 3) + 程序列头(gridSpan)
    cells = [make_tc(lead_headers[0], bold=True, row_span=3, lnB_solid=True, fill=HDR_FILL, font_size=fz),
             make_tc(lead_headers[1], bold=True, row_span=3, lnB_solid=True, fill=HDR_FILL, font_size=fz)]
    for pname, codes in PROGRAMS:
        cells.append(make_tc(pname, bold=True, grid_span=len(codes) * len(sub_w), lnB_solid=True, fill=HDR_FILL, font_size=fz))
    rows.append((header_h, cells))

    # R1：组别列头
    cells = []
    for _, codes in PROGRAMS:
        for code in codes:
            txt = LABEL[code]
            if ANNO[code]:
                txt += '\n(' + ANNO[code] + ')'
            cells.append(make_tc(txt, bold=True, grid_span=len(sub_w), lnB_solid=True, fill=HDR_FILL, font_size=fz))
    rows.append((header_h, cells))

    # R2：子列头
    cells = []
    for _, codes in PROGRAMS:
        for _ in codes:
            for sh in sub_headers:
                cells.append(make_tc(sh, bold=True, lnB_solid=True, fill=HDR_FILL, font_size=fz))
    rows.append((header_h, cells))

    # 数据段：有 / 空行 / 无
    for tag, tag_label, h in [('有', '有', you_h), ('无', '无', wu_h)]:
        for i, tp in enumerate(ALL_TP):
            last_row = (tag == '无' and tp == 'M8')
            cells = []
            if i == 0:
                cells.append(make_tc(tag_label, bold=False, row_span=7, lnB_solid=False, fill='F2F2F2', font_size=fz))
            cells.append(make_tc(tp, bold=False, lnB_solid=last_row, font_size=fz))
            for _, codes in PROGRAMS:
                for code in codes:
                    if tp in PROG_TIMEPOINTS[PROG_OF[code]]:
                        for v in get_vals(data, code, tag, tp, mode):
                            cells.append(make_tc(v, bold=False, lnB_solid=last_row, font_size=fz))
                    else:
                        for _ in sub_headers:
                            cells.append(make_tc('/', bold=False, lnB_solid=last_row, font_size=fz))
            rows.append((h, cells))
        if tag == '有':
            for bh in blank_h:
                cells = [make_tc('', bold=False, lnB_solid=False, lnT=False),
                         make_tc('', bold=False, lnB_solid=False, lnT=False)]
                for _ in range(n_sub):
                    cells.append(make_tc('', bold=False, lnB_solid=False, lnT=False))
                rows.append((bh, cells))

    xml, height = table_str(col_widths, rows)
    return xml, sum(col_widths), height


def build_relative(data, mode, analysis_set):
    """相对时间点表：7 行。返回 (xml, width, height)。"""
    if mode == 'sero':
        lead_w = [400000, 420000, 420000]
        sub_w = [260000, 260000, 320000]
        sub_headers = ['N', 'n', '%']
        header_h = [180000, 180000, 180000]
        data_h = [240000, 240000, 240000, 240000]
    else:
        lead_w = [400000, 420000, 420000]
        sub_w = [280000, 680000]
        sub_headers = ['例数', 'GMC']
        header_h = [190000, 190000, 190000]
        data_h = [220000, 220000, 220000, 220000]

    tp_w = 320000
    lead_headers = ['分析集', '糖尿病病史', '时间点\n(全免后）']
    fz = 700

    col_widths = list(lead_w)
    for _, codes in PROGRAMS:
        col_widths.append(tp_w)
        for _ in codes:
            col_widths.extend(sub_w)

    rows = []
    # R0
    cells = [make_tc(lead_headers[0], bold=True, row_span=3, lnB_solid=True, fill=HDR_FILL, font_size=fz),
             make_tc(lead_headers[1], bold=True, row_span=3, lnB_solid=True, fill=HDR_FILL, font_size=fz),
             make_tc(lead_headers[2], bold=True, row_span=3, lnB_solid=True, fill=HDR_FILL, font_size=fz)]
    for pname, codes in PROGRAMS:
        cells.append(make_tc(pname, bold=True, grid_span=1 + len(codes) * len(sub_w), lnB_solid=True, fill=HDR_FILL, font_size=fz))
    rows.append((header_h[0], cells))

    # R1
    cells = []
    for _, codes in PROGRAMS:
        cells.append(make_tc('对应的时间点', bold=True, row_span=2, lnB_solid=True, fill=HDR_FILL, font_size=fz))
        for code in codes:
            txt = LABEL[code]
            if ANNO[code]:
                txt += '\n(' + ANNO[code] + ')'
            cells.append(make_tc(txt, bold=True, grid_span=len(sub_w), lnB_solid=True, fill=HDR_FILL, font_size=fz))
    rows.append((header_h[1], cells))

    # R2
    cells = []
    for _, codes in PROGRAMS:
        for _ in codes:
            for sh in sub_headers:
                cells.append(make_tc(sh, bold=True, lnB_solid=True, fill=HDR_FILL, font_size=fz))
    rows.append((header_h[2], cells))

    # 数据 4 行
    ri = 0
    for di, (tag, tag_label) in enumerate([('有', '有'), ('无', '无')]):
        for li, tp_label in enumerate(['全免后1个月', '全免后2个月']):
            last_row = (di == 1 and li == 1)
            cells = []
            if li == 0:
                if di == 0:
                    cells.append(make_tc(analysis_set, bold=False, row_span=4, lnB_solid=last_row, fill='F2F2F2', font_size=fz))
                cells.append(make_tc(tag_label, bold=False, row_span=2, lnB_solid=last_row, fill='F2F2F2', font_size=fz))
            cells.append(make_tc(tp_label, bold=False, lnB_solid=last_row, font_size=fz))
            for pname, codes in PROGRAMS:
                m = REL_MAP[tp_label][pname]
                cells.append(make_tc(m, bold=False, lnB_solid=last_row, font_size=fz))
                for code in codes:
                    for v in get_vals(data, code, tag, m, mode):
                        cells.append(make_tc(v, bold=False, lnB_solid=last_row, font_size=fz))
            rows.append((data_h[ri], cells))
            ri += 1

    xml, height = table_str(col_widths, rows)
    return xml, sum(col_widths), height


# ---------------- 幻灯片构建 ----------------
def add_textbox(slide, text, left, top, width, height, size_pt, bold, color, fill=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for line in text.split('\n'):
        p = tf.paragraphs[0] if line == text.split('\n')[0] else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.name = '微软雅黑'
        run.font.color.rgb = RGBColor.from_string(color)
    if fill:
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        tb.fill.background()
    tb.line.fill.background()
    return tb


def add_table_frame(slide, xml_str, left, top, width, height):
    # 先用 1x1 占位表，再替换为自定义 <a:tbl>，避免 deepcopy 既有形状
    placeholder = slide.shapes.add_table(1, 1, Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)))
    gf = placeholder._element  # p:graphicFrame
    tbl_el = etree.fromstring(xml_str)
    old_tbl = gf.find('.//' + qn('a:tbl'))
    old_tbl.getparent().replace(old_tbl, tbl_el)
    return gf


def build_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    kind = spec['kind']
    sh = prs.slide_height

    # 侧栏标签
    sb = spec['sidebar']
    add_textbox(slide, '2期免疫原性结果', sb[0], sb[1], sb[2], sb[3],
                size_pt=20, bold=True, color='FFFFFF', fill='A11C1D',
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 标题
    t = spec['title_pos']
    add_textbox(slide, spec['title'], t[0], t[1], t[2], t[3],
                size_pt=14, bold=True, color='1F1F1F')

    # 表格：根据实际高宽动态排布，避免与描述重叠或超出右边界
    TABLE_LEFT = int(0.10 * 914400)
    GAP_TABLE = int(0.35 * 914400)
    GAP_DESC = int(0.20 * 914400)
    BOTTOM_MARGIN = int(0.20 * 914400)

    if kind.endswith('_abs'):
        xml_str, tw, th = spec['tables'][0]
        top = int(1.0 * 914400)
        add_table_frame(slide, xml_str, TABLE_LEFT, top, tw, th)
        desc_top = top + th + GAP_DESC
        desc_h = sh - desc_top - BOTTOM_MARGIN
        add_textbox(slide, spec['desc'], TABLE_LEFT, desc_top, int(9.5 * 914400), desc_h,
                    size_pt=10, bold=False, color='333333')
    else:
        xml1, tw1, th1 = spec['tables'][0]
        xml2, tw2, th2 = spec['tables'][1]
        top1 = int(1.0 * 914400)
        add_table_frame(slide, xml1, TABLE_LEFT, top1, tw1, th1)
        top2 = top1 + th1 + GAP_TABLE
        add_table_frame(slide, xml2, TABLE_LEFT, top2, tw2, th2)
        desc_top = top2 + th2 + GAP_DESC
        desc_h = sh - desc_top - BOTTOM_MARGIN
        add_textbox(slide, spec['desc'], int(5.5 * 914400), desc_top, int(4.25 * 914400), desc_h,
                    size_pt=10, bold=False, color='333333')


def main():
    desc_dist = ('本分析纳入基线合并糖尿病病史的受试者共 30 例，其中 18-59 岁 18 例'
                 '（0,1月低剂量组 5 例、0,1月高剂量组 3 例、0,2月低剂量组 2 例、'
                 '0,2月高剂量组 3 例、阳性对照组 5 例）、≥60 岁 12 例'
                 '（0,1,6月高剂量组 8 例、阳性对照组 4 例）。')

    specs = [
        # 1. 阳转率 PPS 绝对时间点
        {'kind': 'sero_abs',
         'title': '有、无糖尿病病史人群接种后抗-HBs阳转率（PPS）_绝对时间点',
         'desc': desc_dist + '结果显示：有糖尿病病史人群完成全程 2 剂/3 剂接种后的阳转率总体达到较高水平，'
                             '0,2月程序两组在 M3 均达 100.00%，0,1,6月高剂量组及两个阳性对照组在 M7 均达 100.00%。',
         'tables': [build_absolute(PPS, 'sero')]},
        # 2. 阳转率 FAS 绝对时间点
        {'kind': 'sero_abs',
         'title': '有、无糖尿病病史人群接种后抗-HBs阳转率（FAS）_绝对时间点',
         'desc': '基于 FAS 分析的趋势与 PPS 一致，各组有、无糖尿病病史人群完成全程 2 剂/3 剂接种后的阳转率均达到较高水平。',
         'tables': [build_absolute(FAS, 'sero')]},
        # 3. 阳转率 相对时间点（PPS + FAS）
        {'kind': 'sero_rel',
         'title': '有、无糖尿病病史人群接种后抗-HBs阳转率（相对时间点）',
         'desc': '全免后 1 个月、2 个月，各组有、无糖尿病病史人群的阳转率均维持在较高水平，FAS 与 PPS 趋势一致。',
         'tables': [build_relative(PPS, 'sero', 'PPS'), build_relative(FAS, 'sero', 'FAS')]},
        # 4. GMC PPS 绝对时间点
        {'kind': 'gmc_abs',
         'title': '有、无糖尿病病史人群接种后抗-HBs GMC（PPS）_绝对时间点',
         'desc': '本分析纳入基线合并糖尿病病史的受试者共 30 例，其中 18-59 岁 18 例、≥60 岁 12 例。'
                 '结果显示：各组有、无糖尿病病史人群完成全程 2 剂/3 剂接种后各时间点的抗体 GMC 均呈现上升趋势；'
                 '抗-HBs 浓度低于检测下限（<2.00 mIU/mL）的样本按 LLOQ/2（1.00 mIU/mL）计入 GMC 计算。',
         'tables': [build_absolute(PPS, 'gmc')]},
        # 5. GMC FAS 绝对时间点
        {'kind': 'gmc_abs',
         'title': '有、无糖尿病病史人群接种后抗-HBs GMC（FAS）_绝对时间点',
         'desc': '基于 FAS 分析的趋势与 PPS 一致，各组有、无糖尿病病史人群各时间点抗体 GMC 的变化趋势相似。',
         'tables': [build_absolute(FAS, 'gmc')]},
        # 6. GMC 相对时间点（PPS + FAS）
        {'kind': 'gmc_rel',
         'title': '有、无糖尿病病史人群接种后抗-HBs GMC（相对时间点）',
         'desc': '全免后 1 个月、2 个月，各组有、无糖尿病病史人群的抗体 GMC 均维持在一定水平，FAS 与 PPS 趋势一致。',
         'tables': [build_relative(PPS, 'gmc', 'PPS'), build_relative(FAS, 'gmc', 'FAS')]},
    ]

    # 版式坐标（EMU）：仅保留侧栏与标题，表格/描述位置由 build_slide 根据实际尺寸动态计算
    LAYOUT = {
        'sero_abs': {
            'sidebar': (0, 0.259 * 914400, 2.926 * 914400, 0.374 * 914400),
            'title_pos': (2.926 * 914400, 0.312 * 914400, 7.106 * 914400, 0.242 * 914400),
        },
        'sero_rel': {
            'sidebar': (0, 0.259 * 914400, 2.926 * 914400, 0.374 * 914400),
            'title_pos': (3.098 * 914400, 0.300 * 914400, 7.106 * 914400, 0.242 * 914400),
        },
        'gmc_abs': {
            'sidebar': (0, 0.227 * 914400, 2.984 * 914400, 0.435 * 914400),
            'title_pos': (2.984 * 914400, 0.286 * 914400, 6.904 * 914400, 0.337 * 914400),
        },
        'gmc_rel': {
            'sidebar': (0, 0.227 * 914400, 2.984 * 914400, 0.435 * 914400),
            'title_pos': (3.340 * 914400, 0.258 * 914400, 6.904 * 914400, 0.337 * 914400),
        },
    }

    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 5143500

    for spec in specs:
        spec.update(LAYOUT[spec['kind']])
        build_slide(prs, spec)

    prs.save(OUT)
    print('saved:', OUT)
    print('total slides:', len(prs.slides))


if __name__ == '__main__':
    main()
