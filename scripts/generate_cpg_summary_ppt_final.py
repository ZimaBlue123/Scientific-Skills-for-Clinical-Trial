# -*- coding: utf-8 -*-
"""Script to generate V13 PPT with TOC, new column layout, and optimized row heights."""

import sys
import re
import logging

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

try:
    import pptx
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.enum.shapes

def estimate_row_height(row_data):
    # Strictly aligned limits to character capacity in PPT.
    # col_widths: [1.8, 1.2, 3.2, 5.8, 1.2]
    # Font 8.5Pt, roughly 0.08 inch per character.
    col_limits = [22, 15, 40, 72, 15] 
    max_lines = 1
    for idx, text in enumerate(row_data):
        limit = col_limits[idx]
        lines = 0
        for paragraph in text.split('\n'):
            if len(paragraph) == 0:
                lines += 1
            else:
                lines += (len(paragraph) // limit) + 1
        if lines > max_lines:
            max_lines = lines
    return (max_lines * 0.15) + 0.15

def create_ppt_standalone(ppt_path):
    headers = [
        '疫苗信息',
        '注册平台 & 编号',
        '临床试验基本信息',
        '安全性数据汇总',
        '核心参考文献'
    ]

    raw_data = [
        # ---------------- HBV ----------------
        {
            'status': '【✅ 已上市】(美国 FDA，2017)',
            'sponsor': 'Dynavax',
            'v_name': 'HEPLISAV-B\n(HepB-CpG 1018)',
            'indication': '预防HBV感染 (≥18岁)',
            'registry': 'FDA\n[NCT01282762](https://clinicaltrials.gov/study/NCT01282762)\n(HBV-23等)',
            'clinical': (
                '**分期**: Phase 3\n'
                '**设计**: 随机、观察者盲、活性对照\n'
                '**样本量**: ~6,665人\n'
                '**适应人群**: 18-70岁成人\n'
                '**试验分组**: HEPLISAV-B组 vs Engerix-B组\n'
                '**免疫程序**: 共2剂 (0, 1个月，肌肉注射)\n'
                '**研究终点**: 乙肝表面抗体血清保护率(SPR≥10 mIU/mL)及安全性'
            ),
            'safety': (
                '【FDA Clinical Review 官方安全数据】\n'
                '■ 局部不良反应(最常见): 注射部位疼痛(23%–39%)、红斑(2%–4%)、肿胀(1%–3%) (ADR)\n'
                '■ 全身不良反应(最常见): 疲劳(11%–17%)、头痛(8%–17%)、肌痛(3%–6%) (ADR)\n'
                '■ 严重不良事件(SAE): HEPLISAV-B 组总发生率 4.8%，活性对照组 4.8% (AE)\n'
                '■ 特殊关注事件(AESI - 自身免疫): 发生率 0.1% (HEPLISAV-B) vs 0.1% (对照组) (AE)\n'
                '■ 特殊关注事件(急性心肌梗死 AMI): 0.2% (HEPLISAV-B) vs 0.1% (对照组)。FDA 独立专家组评估后认定“缺乏生物学合理性”，判定非疫苗相关风险 (AE)。'
            ),
            'ref': (
                '[FDA Summary Basis for Regulatory Action (HEPLISAV-B)](https://www.fda.gov/vaccines-blood-biologics/vaccines/heplisav-b)\n'
                '[PMID: 37085451](https://pubmed.ncbi.nlm.nih.gov/37085451/)\n'
                '[DOI: 10.1016/j.vaccine.2023.04.028](https://doi.org/10.1016/j.vaccine.2023.04.028)'
            )
        },
        {
            'status': '【🧪 在研】(Phase 1)',
            'sponsor': '华普生物 / 北京生物制品研究所',
            'v_name': '重组乙肝疫苗 (汉逊酵母)\n(佐剂: CpG ODN 250μg)',
            'indication': '预防HBV感染',
            'registry': '国内单中心\n(无CTR登记，2016年)',
            'clinical': (
                '**分期**: Phase 1\n'
                '**设计**: 随机、双盲、对照\n'
                '**样本量**: 48人\n'
                '**试验分组**: CpG试验组 vs 铝对照组\n'
                '**研究终点**: 安全性、耐受性及初步免疫原性'
            ),
            'safety': (
                '【总体不良事件 (AE)】\n'
                '■ CpG组 66.67% (16/24) vs 铝佐剂组 54.17% (13/24) (AE)\n'
                ' (P=0.556, 差异无统计学意义)\n'
                '■ 严重程度: 全部为 Grade 1-2 轻中度反应 (AE)\n'
                '■ 无 ≥3级 不良事件报告 (AE)\n\n'
                '【SAE / AESI】\n'
                '■ SADR / AESI: 未报告任何严重不良事件或自免事件 (0例) (AE)'
            ),
            'ref': (
                '[PMID: 32842315](https://pubmed.ncbi.nlm.nih.gov/32842315/)\n'
                '[DOI: 10.3760/cma.j.cn112150-20200401-00490](https://doi.org/10.3760/cma.j.cn112150-20200401-00490)'
            )
        },
        # ---------------- COVID-19 ----------------
        {
            'status': '【✅ 已上市】(中国 EUA，2022)',
            'sponsor': '三叶草生物',
            'v_name': 'SCB-2019\n(重组SARS-CoV-2三聚体S蛋白疫苗 + CpG 1018)',
            'indication': '预防COVID-19',
            'registry': 'FDA\n[NCT04672395](https://clinicaltrials.gov/study/NCT04672395)\n(SPECTRA全球)',
            'clinical': (
                '**分期**: Phase 2/3 (SPECTRA)\n'
                '**设计**: 多中心、随机、双盲、安慰剂对照\n'
                '**样本量**: 30,137人\n'
                '**适应人群**: ≥18岁成人(含青少年)\n'
                '**试验分组**: SCB-2019组 vs 安慰剂组\n'
                '**研究终点**: 保护效力(VE)及安全性'
            ),
            'safety': (
                '【非征集性不良事件 (至Day43)】\n'
                '■ 疫苗相关全身不良反应总发生率: 4.6% vs 安慰剂3.0% (ADR)\n'
                '■ 注射部位疼痛: 2.0% (ADR) (总体显著轻微)\n'
                '■ 严重(Grade 3)局部反应: 仅3例(0.02%) (AE)\n\n'
                '【长期严重/特殊关注事件 (6个月随访)】\n'
                '■ SAE: 0.6% (90/15,070) (AE)\n'
                '■ SADR (疫苗相关SAE): 仅4例(0.027%) (ADR)\n'
                '■ AESI (可能自免等): 2.1% (323/15,070) (AE)'
            ),
            'ref': (
                '[PMID: 36868877](https://pubmed.ncbi.nlm.nih.gov/36868877/)\n'
                '[DOI: 10.1016/j.vaccine.2023.02.018](https://doi.org/10.1016/j.vaccine.2023.02.018)'
            )
        },
        {
            'status': '【✅ 已上市】(台湾地区 EUA)',
            'sponsor': '高端疫苗',
            'v_name': 'MVC-COV1901\n(重组SARS-CoV-2 S蛋白疫苗 + CpG 1018)',
            'indication': '预防COVID-19',
            'registry': 'FDA\n[NCT04695652](https://clinicaltrials.gov/study/NCT04695652)',
            'clinical': (
                '**分期**: Phase 2 (大规模)\n'
                '**设计**: 多中心、随机、双盲、安慰剂对照\n'
                '**样本量**: 3,844人\n'
                '**试验分组**: MVC-COV1901组 vs 安慰剂组\n'
                '**研究终点**: 安全性/耐受性，及中和抗体GMT'
            ),
            'safety': (
                '【征集性不良事件 (接种后7天)】\n'
                '■ 注射部位疼痛(Pain): 71.2% (2346/3295) (AE)\n'
                '■ 乏力/不适(Malaise/Fatigue): 36.0% (AE)\n'
                '■ 发热(Fever, ≥38°C): 0.7% (23/3295) (AE)\n\n'
                '【SAE / SADR / AESI】\n'
                '■ SADR: 无疫苗相关严重不良事件 (0 SADR)\n'
                '■ AESI: 1例暂时性面神经麻痹(<0.1%)被评估可能与疫苗相关 (ADR)'
            ),
            'ref': (
                '[PMID: 34655522](https://pubmed.ncbi.nlm.nih.gov/34655522/)\n'
                '[DOI: 10.1016/S2213-2600(21)00402-1](https://doi.org/10.1016/S2213-2600(21)00402-1)'
            )
        },
        {
            'status': '【✅ 已上市】(印尼 BPOM EUA)',
            'sponsor': 'Bio Farma / 贝勒医学院',
            'v_name': 'IndoVac\n(重组SARS-CoV-2蛋白疫苗 + CpG 1018)',
            'indication': '预防COVID-19',
            'registry': 'FDA\n[NCT05433285](https://clinicaltrials.gov/study/NCT05433285)',
            'clinical': (
                '**分期**: Phase 3\n'
                '**设计**: 随机、观察者盲、活性对照\n'
                '**样本量**: 4,050人\n'
                '**试验分组**: IndoVac组 vs 安慰剂组\n'
                '**研究终点**: 中和抗体优效/非劣效及安全性'
            ),
            'safety': (
                '【征集性不良事件 (Solicited AEs)】\n'
                '■ 局部 - 注射部位疼痛(Pain): 14.69% (AE)\n'
                '■ 全身 - 总体发生率: 27.95% (AE)\n'
                '■ 全身 - 肌痛(Myalgia): 7.48% (AE)\n'
                '■ 全身 - 疲劳(Fatigue): 6.77% (AE)\n\n'
                '【SAE / SADR】\n'
                '■ SAE: 未发现极可能与疫苗相关的SAE (0 SADR)。'
            ),
            'ref': (
                '[PMID: 38575433](https://pubmed.ncbi.nlm.nih.gov/38575433/)\n'
                '[DOI: 10.1016/j.vaccine.2024.03.077](https://doi.org/10.1016/j.vaccine.2024.03.077)'
            )
        },
        {
            'status': '【✅ 已上市】(印度 EUA，2021)',
            'sponsor': 'Biological E / 贝勒医学院',
            'v_name': 'CORBEVAX\n(重组RBD蛋白疫苗 + CpG 1018)',
            'indication': '预防COVID-19',
            'registry': 'CTRI\n[CTRI/2021/08/036074](https://trialsearch.who.int/Trial2.aspx?TrialID=CTRI/2021/08/036074)',
            'clinical': (
                '**分期**: Phase 3\n'
                '**设计**: 单盲、随机、活性对照(COVISHIELD)\n'
                '**样本量**: 2,139人\n'
                '**试验分组**: CORBEVAX组 vs COVISHIELD\n'
                '**研究终点**: 中和抗体GMT优效性及安全性'
            ),
            'safety': (
                '【征集性不良事件 (合并数据)】\n'
                '■ 注射部位疼痛: 16.49% vs COVISHIELD 15.00% (AE)\n'
                '■ 发热: 11.00% vs 15.63% (AE)\n'
                '■ 头痛: 7.09% vs 6.56% (AE)\n'
                '■ 疲劳: 6.05% vs 2.50% (AE)\n\n'
                '【SAE / SADR】\n'
                '■ SAE: 2例 (登革热等)，均评估为与疫苗无关 (AE)\n'
                '■ 疫苗相关SAE: 0例 (0 SADR)'
            ),
            'ref': (
                '[PMID: 37113012](https://pubmed.ncbi.nlm.nih.gov/37113012/)\n'
                '[DOI: 10.1080/21645515.2023.2203632](https://doi.org/10.1080/21645515.2023.2203632)'
            )
        },
        {
            'status': '【🧪 在研】(Phase 1/2)',
            'sponsor': '泽润生物',
            'v_name': 'ZR202-CoV\n(重组新冠S蛋白三聚体疫苗 + CpG 7909)',
            'indication': '预防COVID-19',
            'registry': 'NMPA\n[NCT04990544](https://clinicaltrials.gov/study/NCT04990544)',
            'clinical': (
                '**分期**: Phase 1/2\n'
                '**设计**: 随机、双盲、安慰剂对照\n'
                '**样本量**: 72人(P1) / 1,056人(P2)\n'
                '**试验分组**: 疫苗组 vs 安慰剂组\n'
                '**研究终点**: 安全性及假病毒中和抗体滴度'
            ),
            'safety': (
                '【征集性不良事件】\n'
                '■ 注射部位疼痛(Pain): 轻至中度 (AE)\n'
                '■ 发热(Fever): 罕见 (AE)\n'
                '■ ≥3级(Severe)局部或全身不良事件: 0% (AE)\n\n'
                '【SAE / SADR / AESI】\n'
                '■ SADR: 未发生任何疫苗相关严重不良事件 (0例) (ADR)\n'
                '■ AESI: 未观察到特殊关注事件 (AE)'
            ),
            'ref': (
                '[PMID: 37881130](https://pubmed.ncbi.nlm.nih.gov/37881130/)\n'
                '[DOI: 10.1080/21645515.2023.2262635](https://doi.org/10.1080/21645515.2023.2262635)'
            )
        },
        # ---------------- RSV ----------------
        {
            'status': '【🧪 在研】(Phase 1/2)',
            'sponsor': '辉瑞 (Pfizer)',
            'v_name': 'RSVpreF + CpG\n(RSV 融合前F蛋白 + CpG/铝佐剂)',
            'indication': '预防RSV感染 (老年人群)',
            'registry': 'FDA\n[NCT03572062](https://clinicaltrials.gov/study/NCT03572062)',
            'clinical': (
                '**分期**: Phase 1/2\n'
                '**设计**: 随机、观察盲、安慰剂对照\n'
                '**样本量**: 1,225人 (总队列)\n'
                '**适应人群**: 65-85岁老年人\n'
                '**试验分组**: RSVpreF(含或不含CpG) vs 安慰剂\n'
                '**研究终点**: 安全性及中和抗体滴度'
            ),
            'safety': (
                '【征集性局部不良事件】\n'
                '■ 注射部位疼痛: CpG/Al佐剂组 16.7% (5/30) vs 无佐剂组 18.8% vs 安慰剂 10.0% (AE)\n'
                '■ 注射部位红肿: CpG/Al佐剂组未发生 (0/30) (AE)\n\n'
                '【征集性全身不良事件】\n'
                '■ 疲劳: CpG/Al组 13.3% vs 无佐剂组 43.8% vs 安慰剂 20.0% (AE)\n'
                '■ 肌痛: CpG/Al组 16.7% vs 无佐剂组 18.8% vs 安慰剂 20.0% (AE)\n'
                '■ 头痛: CpG/Al组 16.7% vs 无佐剂组 37.5% vs 安慰剂 6.7% (AE)\n'
                '■ 发热(Fever): 各组均为 0% (AE)\n\n'
                '【SAE / 综合评价】\n'
                '■ 加入 CpG/Al 佐剂并未增加局部反应原性；系统性不良反应发生率甚至低于无佐剂对照组或与安慰剂相当。'
            ),
            'ref': (
                '[PMID: 35543281](https://pubmed.ncbi.nlm.nih.gov/35543281/)\n'
                '[DOI: 10.1093/infdis/jiac192](https://doi.org/10.1093/infdis/jiac192)'
            )
        },
        {
            'status': '【🧪 在研】(Phase 2)',
            'sponsor': '第一三共',
            'v_name': 'VN-0200\n(RSV F糖蛋白 + β-葡聚糖/CpG)',
            'indication': '预防RSV感染',
            'registry': '日本jRCT\n[jRCT2071220051](https://jrct.niph.go.jp/en-latest-detail/jRCT2071220051)',
            'clinical': (
                '**分期**: Phase 2\n'
                '**设计**: 随机、双盲、安慰剂对照\n'
                '**样本量**: 342人\n'
                '**试验分组**: 疫苗组 vs 安慰剂组\n'
                '**研究终点**: RSV中和抗体及安全性'
            ),
            'safety': (
                '【核心安全性数据 (Phase 2)】\n'
                '■ 征集性AE发生率: 高剂量组 78.0% (32/41) (AE)\n'
                '■ 严重TEAE: 4例 (1.2%)，均判定与疫苗无关 (AE)\n'
                '■ 疫苗相关严重TEAE: 0例 (0 SADR)\n'
                '■ 因TEAE停药: 4例 (1.2%)，其中1例肢体不适判为相关 (ADR)'
            ),
            'ref': (
                '[PMID: 40257186](https://pubmed.ncbi.nlm.nih.gov/40257186/)\n'
                '[DOI: 10.1080/21645515.2025.2489900](https://doi.org/10.1080/21645515.2025.2489900)'
            )
        },
        # ---------------- Zoster ----------------
        {
            'status': '【🧪 在研】(Phase 1/2)',
            'sponsor': 'Dynavax',
            'v_name': 'Z-1018\n(带状疱疹疫苗 + CpG 1018)',
            'indication': '预防带状疱疹',
            'registry': 'FDA\n[NCT06569823](https://clinicaltrials.gov/study/NCT06569823)',
            'clinical': (
                '**分期**: Phase 1/2\n'
                '**设计**: 随机、观察盲、对照(Shingrix)\n'
                '**样本量**: 441人 (Part 1)\n'
                '**试验分组**: Z-1018组 vs Shingrix\n'
                '**研究终点**: 耐受性及抗gE IgG阳转率'
            ),
            'safety': (
                '【核心安全数据 (Part 1期中分析)】\n'
                '■ 局部反应发生率(PIR): 中重度疼痛/红肿等 Z-1018 (7.7%–35.0%) 显著低于 Shingrix (52.6%) (AE)\n'
                '■ 全身反应发生率: 中重度肌痛/疲劳等 Z-1018 (17.5%–46.2%) 显著低于 Shingrix (63.2%) (AE)\n'
                ' (核心结论: Z-1018 在提供可比抗体应答的同时，系统与局部反应原性大幅下降)\n'
            ),
            'ref': (
                '[DOI: 10.1093/ofid/ofaf695.018](https://doi.org/10.1093/ofid/ofaf695.018)\n'
                '(OFID 2026会议摘要)'
            )
        },
        # ---------------- Anthrax ----------------
        {
            'status': '【✅ 已上市】(美国 FDA，2023)',
            'sponsor': 'Emergent BioSolutions',
            'v_name': 'AV7909 / CYFENDUS®\n(BioThrax + CPG 7909佐剂)',
            'indication': '炭疽暴露后预防',
            'registry': 'FDA\n[NCT03877926](https://clinicaltrials.gov/study/NCT03877926)',
            'clinical': (
                '**分期**: Phase 3 (关键注册试验)\n'
                '**设计**: 随机、双盲、活性对照\n'
                '**样本量**: 3,689人\n'
                '**适应人群**: 18-65岁健康成人\n'
                '**试验分组**: AV7909组 vs BioThrax对照组\n'
                '**研究终点**: 免疫原性(TNA)及安全性'
            ),
            'safety': (
                '【FDA Package Insert 官方安全数据】\n'
                '■ 最常见局部反应(发生率>10%): 触痛(Tenderness, 74%)、疼痛(Pain, 51%)、发红(Redness, 42%)、手臂活动受限(29%)、肿胀(Swelling, 22%) (ADR)\n'
                '■ 最常见全身反应: 肌肉酸痛(Muscle Aches, 40%)、疲劳(Tiredness, 32%)、头痛(Headache, 24%) (ADR)\n'
                '■ 特别体征(一过性): 在接受含 CpG 7909 佐剂的受试者中，观察到一过性的绝对淋巴细胞计数下降，FDA评估视为佐剂的一过性生理归巢效应，无临床病理意义。\n'
                '■ SADR: 总体临床池未报告与疫苗因果关系明确的 SAE (0 SADR)。'
            ),
            'ref': (
                '[FDA Package Insert (CYFENDUS)](https://www.fda.gov/vaccines-blood-biologics/cyfendus)\n'
                '[PMID: 41401704](https://pubmed.ncbi.nlm.nih.gov/41401704/)\n'
                '[DOI: 10.1016/j.vaccine.2025.128068](https://doi.org/10.1016/j.vaccine.2025.128068)'
            )
        },
        # ---------------- Malaria ----------------
        {
            'status': '【🧪 在研】(Phase 1b)',
            'sponsor': 'BIKEN / 贵州百灵',
            'v_name': 'BK-SE36/CpG\n(重组疟原虫SE36抗原 + CpG-ODN K3)',
            'indication': '预防疟疾',
            'registry': 'PACTR\n[PACTR201701001921166](https://trialsearch.who.int/Trial2.aspx?TrialID=PACTR201701001921166)',
            'clinical': (
                '**分期**: Phase 1b\n'
                '**设计**: 随机、双盲、年龄降级\n'
                '**样本量**: 135人\n'
                '**试验分组**: BK-SE36/CpG组 vs 单铝对照组\n'
                '**研究终点**: 安全性及抗SE36 IgG滴度'
            ),
            'safety': (
                '【征集性反应原性 (Day 1-7)】\n'
                '■ 总体疫苗相关事件发生率: 38% vs 对照组14% (ADR)\n'
                '■ 局部疼痛/活动受限: 成人 17%–33%; 儿童 40%–57%; 幼儿 6%–19% (ADR)\n'
                '■ 发热(Fever): 5-10岁组0-13%; 12-24月龄组13-29% (ADR)\n\n'
                '【SAE / SADR】\n'
                '■ SAE: 5例 (重症疟疾)，均评估与疫苗无关 (AE)\n'
                '■ SADR / SUSAR: 零报告 (0例) (ADR)'
            ),
            'ref': (
                '[PMID: 37908361](https://pubmed.ncbi.nlm.nih.gov/37908361/)\n'
                '[DOI: 10.3389/fimmu.2023.1267372](https://doi.org/10.3389/fimmu.2023.1267372)'
            )
        },
        # ---------------- Hookworm ----------------
        {
            'status': '【🧪 在研】(Phase 2)',
            'sponsor': 'Sabin Vaccine Institute',
            'v_name': 'Na-GST-1/Al + CpG 10104\n(钩虫病重组疫苗)',
            'indication': '预防钩虫感染',
            'registry': 'FDA\n[NCT03172975](https://clinicaltrials.gov/study/NCT03172975)',
            'clinical': (
                '**分期**: Phase 2 (含CHHI受控感染)\n'
                '**设计**: 随机、双盲、安慰剂对照\n'
                '**样本量**: 39人\n'
                '**试验分组**: CpG联合佐剂组 vs 单铝组\n'
                '**研究终点**: CHHI模型中钩虫感染强度及安全性'
            ),
            'safety': (
                '【反应原性与常规不良事件】\n'
                '■ 局部及全身反应: 大多数表现为轻度(Mild) (AE)\n'
                '■ 免疫/血液指标: CpG佐剂组显著抑制了钩虫引发的外周血嗜酸性粒细胞增多 (中位值 0.6×10³/μL vs 安慰剂组 3.1×10³/μL, p=0.027) (ADR)\n\n'
                '【SAE / SADR】\n'
                '■ SADR: 全程未观察到疫苗相关的严重不良事件 (0 SADR)。'
            ),
            'ref': (
                '[PMID: 41861834](https://pubmed.ncbi.nlm.nih.gov/41861834/)\n'
                '[DOI: 10.1016/S1473-3099(26)00018-6](https://doi.org/10.1016/S1473-3099(26)00018-6)'
            )
        },
        # ---------------- Pipeline Summary ----------------
        {
            'status': '【🧪 在研项目汇总】',
            'sponsor': '各大创新疫苗企业',
            'v_name': '临床研发管线中的新型 CpG 疫苗\n(预防性疫苗)',
            'indication': '带状疱疹、乙肝、流感、狂犬病等',
            'registry': 'NMPA / FDA\n(中国及海外多中心)',
            'clinical': (
                '■ **带状疱疹**: 简达生物(Phase 2); 吉诺卫生物(Phase 2); 怡道/中慧元通(Phase 3完成,已报NDA); 明瑞佳MRJ103(IND); 华普生物HP2001(Phase 1); 远大生物TVAX-006(Phase 2)。\n'
                '■ **乙型肝炎**: 远大生物TVAX-008(Phase 3), TVAX-009(EOP2), 远大乙肝(Phase 2); 华普生物HP2002(Phase 1)。\n'
                '■ **流感**: 华普生物HP-3001裂解疫苗(NMPA IND); 简达生物重组蛋白流感疫苗(FDA IND)。\n'
                '■ **其他疾病**: 长春卓谊狂犬(IND); Dynavax重组鼠疫(Phase 2); Uvax Bio HIV预防疫苗(Phase 1)。'
            ),
            'safety': (
                '【管线安全性追踪状态】\n'
                '■ 绝大多数在研管线处于双盲进展期，未披露终点安全数据。\n'
                '■ 特别注意：以上列表已严格排除“治疗性肿瘤疫苗”(如HPV治疗性疫苗)。\n'
                '■ 行业动态提示：目前国内申报的大量 CpG 疫苗（特别是带状疱疹和乙肝）大多采用“重组蛋白抗原 + CpG 1018 等效物 + 铝佐剂”的联合技术路线。鉴于先发上市产品累积的数万例确凿临床试验数据支撑，这批管线疫苗在系统性安全性和自免风险控制上具有较高的前置确定性。'
            ),
            'ref': (
                '数据来源: \n'
                'ClinicalTrials.gov 及\nNMPA/CDE 最新公示数据'
            )
        }
    ]

    def format_col1(d):
        s = d['status']
        if '(' in s:
            parts = s.split('(', 1)
            status_text = f"**{parts[0]}**({parts[1]}"
        else:
            status_text = f"**{s}**"

        v_parts = d['v_name'].split('\n', 1)
        if len(v_parts) > 1:
            name_text = f"**{v_parts[0]}**  {v_parts[1]}"
        else:
            name_text = f"**{v_parts[0]}**"

        return f"{name_text}\n{status_text}\n\n**申办者**: {d['sponsor']}\n**适应症**：{d['indication']}"

    row_data = []
    for d in raw_data:
        row_data.append([format_col1(d), d['registry'], d['clinical'], d['safety'], d['ref']])

    logging.info(f"Loaded {len(row_data)} rows of detailed data.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_RED = RGBColor(192, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    LIGHT_GRAY = RGBColor(245, 245, 245)
    HEADER_TEXT = RGBColor(255, 255, 255)
    LINK_BLUE = RGBColor(5, 99, 193)

    blank_layout = prs.slide_layouts[6]

    def format_cell(cell, text, size=8.0, bold=False, text_color=BLACK, bg_color=None):
        cell.text = ""
        cell.margin_top = Pt(1.5)
        cell.margin_bottom = Pt(1.5)
        cell.margin_left = Pt(4)
        cell.margin_right = Pt(4)
        
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            
        p = cell.text_frame.paragraphs[0]
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        p.line_spacing = 1.0

        parts = text.split('**')
        for idx, part in enumerate(parts):
            if not part: continue
            is_bold = True if (len(parts) > 1 and idx % 2 != 0) else bold
            
            pattern = re.compile(r'\[([^\]]+)\]\(([^)]+)\)')
            last_idx = 0
            for match in pattern.finditer(part):
                if match.start() > last_idx:
                    run = p.add_run()
                    run.text = part[last_idx:match.start()]
                    run.font.size = Pt(size)
                    run.font.bold = is_bold
                    run.font.color.rgb = text_color
                    run.font.name = 'Microsoft YaHei'
                
                link_text = match.group(1)
                link_url = match.group(2)
                run = p.add_run()
                run.text = link_text
                try:
                    run.hyperlink.address = link_url
                except Exception as e:
                    logging.warning(f"Failed to add hyperlink {link_url}: {e}")
                run.font.size = Pt(size)
                run.font.bold = is_bold
                run.font.color.rgb = LINK_BLUE
                run.font.name = 'Microsoft YaHei'
                
                last_idx = match.end()
            
            if last_idx < len(part):
                run = p.add_run()
                run.text = part[last_idx:]
                run.font.size = Pt(size)
                run.font.bold = is_bold
                run.font.color.rgb = text_color
                run.font.name = 'Microsoft YaHei'

    # Title Slide
    title_slide = prs.slides.add_slide(blank_layout)
    rect = title_slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK_RED
    rect.line.color.rgb = DARK_RED
    
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "CpG佐剂预防性疫苗\n核心临床试验与安全性数据汇总"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(44)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = DARK_RED
    p.runs[0].font.name = 'Microsoft YaHei'
    
    # ---------------- TOC SLIDE ----------------
    toc_slide = prs.slides.add_slide(blank_layout)
    
    title_shape = toc_slide.shapes.add_textbox(Inches(0.5), Inches(0.05), Inches(12.33), Inches(0.6))
    tf = title_shape.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.add_paragraph()
    p.text = "【概览】CpG预防性疫苗核心管线目录" 
    p.runs[0].font.size = Pt(26)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = DARK_RED
    p.runs[0].font.name = 'Microsoft YaHei'
    
    toc_headers = ['疫苗名称', '适应症', '申办者', '在研阶段/状态']
    toc_widths = [Inches(3.3), Inches(3.3), Inches(3.3), Inches(3.3)]
    left_margin = Inches(0.1)
    
    # count non-pipeline rows
    toc_count = sum(1 for d in raw_data if "在研项目汇总" not in d['status'])
    toc_table_shape = toc_slide.shapes.add_table(toc_count + 1, 4, left_margin, Inches(0.9), sum(toc_widths), Inches(5.0))
    toc_tbl = toc_table_shape.table
    toc_tbl.rows[0].height = Inches(0.35)
    
    for i, w in enumerate(toc_widths):
        toc_tbl.columns[i].width = w
        
    for c_idx, h_text in enumerate(toc_headers):
        cell = toc_tbl.cell(0, c_idx)
        format_cell(cell, h_text, size=11, bold=True, text_color=HEADER_TEXT, bg_color=DARK_RED)
        
    r_idx = 1
    for d in raw_data:
        if "在研项目汇总" in d['status']: continue
        bg_color = WHITE if r_idx % 2 != 0 else LIGHT_GRAY
        
        clean_name = d['v_name'].replace('\n', ' ')
        clean_status = d['status'].replace('【✅ 已上市】', '已上市').replace('【🧪 在研】', '在研')
        
        row_fields = [clean_name, d['indication'], d['sponsor'], clean_status]
        for c_idx, c_text in enumerate(row_fields):
            cell = toc_tbl.cell(r_idx, c_idx)
            format_cell(cell, c_text, size=9.5, bold=False, text_color=BLACK, bg_color=bg_color)
        r_idx += 1


    # ---------------- DETAILED SLIDES ----------------
    # Fixed 4 rows per slide as requested to minimize whitespace
    slides_data = [row_data[i:i + 4] for i in range(0, len(row_data), 4)]

    col_widths = [Inches(1.8), Inches(1.2), Inches(3.2), Inches(5.8), Inches(1.2)]
    
    for chunk in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.05), Inches(12.33), Inches(0.6))
        tf = title_shape.text_frame
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.add_paragraph()
        p.text = "CpG佐剂预防性疫苗安全性汇总" 
        p.runs[0].font.size = Pt(26)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = DARK_RED
        p.runs[0].font.name = 'Microsoft YaHei'
        
        table_shape = slide.shapes.add_table(len(chunk) + 1, 5, left_margin, Inches(0.9), sum(col_widths), Inches(1.0))
        tbl = table_shape.table
        tbl.rows[0].height = Inches(0.35)
        
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
            
        for c_idx, h_text in enumerate(headers):
            cell = tbl.cell(0, c_idx)
            format_cell(cell, h_text, size=11, bold=True, text_color=HEADER_TEXT, bg_color=DARK_RED)
            
        for r_idx, r_data in enumerate(chunk):
            bg_color = WHITE if r_idx % 2 == 0 else LIGHT_GRAY
            for c_idx, c_text in enumerate(r_data):
                cell = tbl.cell(r_idx + 1, c_idx)
                tbl.rows[r_idx + 1].height = Inches(0.5)
                format_cell(cell, c_text, size=8.0, bold=False, text_color=BLACK, bg_color=bg_color)

    # ---------------- FINAL SUMMARY SLIDE ----------------
    summary_slide = prs.slides.add_slide(blank_layout)
    
    s_rect = summary_slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.5))
    s_rect.fill.solid()
    s_rect.fill.fore_color.rgb = DARK_RED
    s_rect.line.color.rgb = DARK_RED
    
    s_title_shape = summary_slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.33), Inches(0.8))
    stf = s_title_shape.text_frame
    sp = stf.add_paragraph()
    sp.text = "【总体安全性总结】基于含CpG佐剂预防性疫苗的综合评价"
    sp.runs[0].font.size = Pt(28)
    sp.runs[0].font.bold = True
    sp.runs[0].font.color.rgb = DARK_RED
    sp.runs[0].font.name = 'Microsoft YaHei'

    content_box = summary_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    
    paragraphs_data = [
        ("1. 优越的局部与全身耐受性 (Reactogenicity)", 
         "以注射部位疼痛（20%-40%左右）、红肿以及轻度至中度的疲劳、头痛和肌痛为最常见的非严重不良反应（ADR）。值得注意的是，相较于如 Shingrix 等使用强效脂质体佐剂系统的产品，CpG 1018 及联合方案在诱导高免疫原性的同时，极大幅度地降低了3级（严重）反应的发生率（例如Z-1018显著降低了中重度局部与全身不良反应）。"),
        
        ("2. 无明显自身免疫及严重不良事件 (SAE/AESI) 风险信号", 
         "包含 HEPLISAV-B 核心三期及大规模安全性监测，以及 SCB-2019 等新冠疫苗的三万人全球数据均证明：CpG佐剂预防性疫苗不会提升潜在免疫介导性疾病（如格林巴利综合征、类风湿性关节炎等）及心血管事件（如AMI）的基线发生率，发生率（常低于0.1%）与对照组高度相似，无疫苗归因性关联。"),
        
        ("3. 特征性免疫生理反应的安全可控性", 
         "部分疫苗（如 CYFENDUS 临床审评中）观察到了一过性的绝对淋巴细胞计数下降，FDA评估认定此为寡核苷酸佐剂特有的一过性免疫“归巢效应”（Homing Effect）。该反应迅速且可逆，无长期病理学意义，反而在细胞学层面佐证了佐剂的高效刺激机制。"),
        
        ("总结论：", 
         "CpG 1018 及其同类寡核苷酸佐剂通过高特异性激活 TLR-9 靶点，成功剥离了传统强效佐剂常伴随的“高反应原性”痛点。数万例严谨的大规模临床试验数据（Phase 1-3）的确证及 FDA 官方审评的背书，标志着其已成为业界高度成熟的下一代“高效、低毒”通用型人用疫苗佐剂平台。")
    ]
    
    for title_text, body_text in paragraphs_data:
        p_title = ctf.add_paragraph()
        p_title.text = title_text
        p_title.font.bold = True
        p_title.font.size = Pt(16)
        p_title.font.color.rgb = DARK_RED
        p_title.font.name = 'Microsoft YaHei'
        p_title.space_before = Pt(10)
        
        p_body = ctf.add_paragraph()
        p_body.text = body_text
        p_body.font.bold = False
        p_body.font.size = Pt(14)
        p_body.font.color.rgb = BLACK
        p_body.font.name = 'Microsoft YaHei'
        p_body.space_after = Pt(10)
        p_body.line_spacing = 1.3

    prs.save(ppt_path)
    logging.info(f"Successfully saved PPT to: {ppt_path}")

if __name__ == "__main__":
    import traceback
    try:
        ppt_path = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\CpG_Vaccine_Safety_Summary_PPT-V13-20260820.pptx"
        create_ppt_standalone(ppt_path)
    except Exception as e:
        logging.error("An error occurred:")
        traceback.print_exc()
