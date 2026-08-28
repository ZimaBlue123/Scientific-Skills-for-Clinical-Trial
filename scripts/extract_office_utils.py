from __future__ import annotations

"""
extract_office_utils.py
Unified extraction utilities for DOCX, PPTX, and XLSX files.
"""

import argparse
import logging
import os
import sys
from pathlib import Path

# ==============================================================================
# DOCX Extraction
# ==============================================================================
#!/usr/bin/env python3
"""
Extract text from .docx / .doc files.

Supports a single file, a list of files, or every supported file in a
folder. Modern .docx is handled via python-docx; legacy .doc is handled
via the Windows-only Word COM interface.

Output formats
--------------
* --format text  (default) plain paragraphs + tables-as-pipes
* --format md    Markdown-flavoured rendering. Replaces the now-archived
                 scripts/extract_docx_to_md.py.

Usage
-----
    py -3 scripts/extract_docx_full.py document.docx [output.txt]
    py -3 scripts/extract_docx_full.py document.docx --format md output.md
    py -3 scripts/extract_docx_full.py folder/ [output.txt]
"""


from collections.abc import Sequence

LOG_FORMAT = "%(asctime)s [%(levelname)s] extract_docx_full: %(message)s"
logger = logging.getLogger("extract_docx_full")

ENCODING = "utf-8"


def extract_docx(filepath: Path) -> str:
    """Extract all text from a .docx using python-docx."""
    from docx import Document  # local import keeps module importable when missing

    doc = Document(str(filepath))
    full_text: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                full_text.append(" | ".join(row_text))

    return "\n\n".join(full_text)


def extract_docx_to_markdown(filepath: Path) -> str:
    """Render a .docx to Markdown.

    Paragraph styles named "Heading 1..3" / "Title" map to ``#``-style
    headings; tables render as GFM pipe tables. Replaces the archived
    ``extract_docx_to_md.py``.
    """
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(str(filepath))
    out: list[str] = []

    def _heading_level(para: Paragraph) -> int | None:
        name = (para.style.name or "").strip().lower()
        if name == "title":
            return 0
        if name.startswith("heading "):
            try:
                return int(name.split()[1])
            except (IndexError, ValueError):
                return None
        return None

    def _esc(s: str) -> str:
        return s.replace("|", "\\|").replace("\n", " ")

    def _render_table(tbl: Table) -> list[str]:
        rows = tbl.rows
        if not rows:
            return []
        md: list[str] = []
        header = [_esc(c.text.strip()) for c in rows[0].cells]
        md.append("| " + " | ".join(header) + " |")
        md.append("|" + "|".join(["---"] * len(header)) + "|")
        for row in rows[1:]:
            cells = [_esc(c.text.strip()) for c in row.cells]
            md.append("| " + " | ".join(cells) + " |")
        return md

    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            para = Paragraph(child, doc)
            text = para.text.rstrip()
            if not text:
                continue
            level = _heading_level(para)
            if level == 0:
                out.append(f"# {text}")
            elif level is not None and 1 <= level <= 6:
                out.append("#" * (level + 1) + " " + text)
            else:
                out.append(text)
        elif isinstance(child, CT_Tbl):
            tbl = Table(child, doc)
            out.extend(_render_table(tbl))
    return "\n\n".join(out) + "\n"


def extract_doc_legacy(filepath: Path) -> str | None:
    """Extract text from a legacy .doc file using Word COM (Windows only)."""
    try:
        import pythoncom  # type: ignore[import-not-found]
        import win32com.client  # type: ignore[import-not-found]
    except ImportError:
        logger.error("pywin32 not available - cannot extract .doc files")
        return None

    pythoncom.CoInitialize()
    word = None
    doc = None
    try:
        try:
            word = win32com.client.Dispatch("Word.Application")
        except Exception as exc:  # noqa: BLE001
            logger.error("failed to launch Word.Application: %s", exc)
            return None
        word.Visible = False
        try:
            doc = word.Documents.Open(os.path.abspath(filepath))
            return str(doc.Content.Text or "")
        except Exception as exc:  # noqa: BLE001
            logger.error("COM error while reading %s: %s", filepath, exc)
            return None
        finally:
            if doc is not None:
                try:
                    doc.Close(False)
                except Exception:  # noqa: BLE001
                    logger.debug("doc.Close raised", exc_info=True)
    finally:
        if word is not None:
            try:
                word.Quit()
            except Exception:  # noqa: BLE001
                logger.debug("word.Quit raised", exc_info=True)
        try:
            pythoncom.CoUninitialize()
        except Exception:  # noqa: BLE001
            logger.debug("pythoncom.CoUninitialize raised", exc_info=True)


def extract_file(filepath: Path) -> str | None:
    """Dispatch to the right backend based on the file extension."""
    suffix = filepath.suffix.lower()
    if suffix == ".docx":
        return extract_docx(filepath)
    if suffix == ".doc":
        return extract_doc_legacy(filepath)
    logger.error("unsupported format: %s", suffix)
    return None


def extract_folder(
    folder_path: Path,
    output_path: Path | None = None,
) -> str:
    """Extract text from all docx/doc files in ``folder_path``."""
    files = [f for f in os.listdir(folder_path) if f.lower().endswith((".docx", ".doc"))]

    combined: list[str] = [f"Found {len(files)} files", ""]

    for fname in files:
        path = Path(folder_path) / fname
        combined.append(f"\n========== {fname} ==========\n")
        try:
            text = extract_file(path)
            combined.append(text if text else "[FAILED to extract content]")
        except Exception as exc:  # noqa: BLE001
            combined.append(f"Error: {exc}")
        combined.append("")

    result = "\n".join(combined)

    if output_path is not None:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(result, encoding=ENCODING)
        print(f"Extracted {len(files)} files to: {output_path}")

    return result


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="extract_docx_full",
        description="Extract text from .docx and legacy .doc files.",
    )
    parser.add_argument("input", help="Input .docx/.doc file or a folder.")
    parser.add_argument("output", nargs="?", default=None, help="Optional output .txt path.")
    parser.add_argument(
        "--log-level",
        default="INFO",
        choices=("DEBUG", "INFO", "WARNING", "ERROR"),
        help="Logging verbosity (default: %(default)s).",
    )
    parser.add_argument(
        "--format",
        default="text",
        choices=("text", "md"),
        dest="fmt",
        help="Output format (default: %(default)s).",
    )
    return parser


def main(argv: Sequence[str] | None = None) -> int:
    parser = _build_parser()
    args = parser.parse_args(argv)
    logging.basicConfig(level=getattr(logging, args.log_level), format=LOG_FORMAT)

    input_p = Path(args.input)
    output_p = Path(args.output) if args.output else None

    if not input_p.exists():
        logger.error("input does not exist: %s", input_p)
        return 2

    if input_p.is_dir():
        extract_folder(input_p, output_p)
        return 0

    if args.fmt == "md":
        try:
            text = extract_docx_to_markdown(input_p)
        except Exception as exc:  # noqa: BLE001
            logger.exception("failed to extract md from %s: %s", input_p, exc)
            return 1
    else:
        try:
            text = extract_file(input_p)
        except Exception as exc:  # noqa: BLE001
            logger.exception("failed to extract %s: %s", input_p, exc)
            return 1

    if not text:
        logger.error("failed to extract content: %s", input_p)
        return 1

    if output_p is not None:
        output_p.parent.mkdir(parents=True, exist_ok=True)
        output_p.write_text(text, encoding=ENCODING)
        print(f"Extracted to: {output_p}")
    else:
        print(text)
    return 0


# ==============================================================================
# PPTX Extraction
# ==============================================================================
"""全量 PPTX 内容提取器（表格 + 文本框 + 组合形状 + 备注）.

用法:
    python extract_pptx_full.py <input.pptx> [output.txt]

相比 scripts/extract_pptx.py（仅提取纯文本），本脚本完整提取:
- 每个 shape 的类型/位置/尺寸（用于排版分析）
- 表格逐行逐列单元格内容（含合并单元格标注）
- 组合形状递归展开
- 幻灯片备注
"""


from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _iter_shape_text(shape, depth=0):
    """递归提取 shape 内容，返回文本片段列表."""
    fragments = []
    indent = "  " * depth

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        fragments.append(f"{indent}[GROUP]")
        for sub in shape.shapes:
            fragments.extend(_iter_shape_text(sub, depth + 1))
        return fragments

    if shape.has_table:
        tbl = shape.table
        fragments.append(f"{indent}[TABLE {len(tbl.rows)}行 x {len(tbl.columns)}列]")
        for r_idx, row in enumerate(tbl.rows):
            cells = []
            for c_idx, cell in enumerate(row.cells):
                text = cell.text.replace("\n", "⏎").strip()
                cells.append(text if text else "∅")
            fragments.append(f"{indent}  R{r_idx}: {' | '.join(cells)}")
        return fragments

    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            top = int(shape.top) if shape.top is not None else -1
            left = int(shape.left) if shape.left is not None else -1
            fragments.append(f"{indent}[TEXT @{top},{left}] {text}")
        return fragments

    return fragments


def extract_full(pptx_path, output_path):
    prs = Presentation(pptx_path)
    lines = []
    lines.append(f"# 文件: {Path(pptx_path).name}")
    lines.append(f"# 幻灯片尺寸: {prs.slide_width} x {prs.slide_height} EMU")
    lines.append(f"# 幻灯片数: {len(prs.slides)}")
    lines.append("")

    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"===== Slide {i} =====")
        for shape in slide.shapes:
            lines.extend(_iter_shape_text(shape))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"[NOTES] {notes}")
        lines.append("")

    content = "\n".join(lines)
    Path(output_path).write_text(content, encoding="utf-8")
    print(f"OK: {len(prs.slides)} slides -> {output_path} ({len(content)} chars)")


# ==============================================================================
# XLSX Extraction
# ==============================================================================
#!/usr/bin/env python3
"""
Extract text/content from .xlsx files (UTF-8, robust).

Why this exists
---------------
``openpyxl.load_workbook`` raises ``ValueError`` when an .xlsx contains
non-conforming XML such as ``<autoFilter ref="...">`` with cell references
outside the canonical pattern. Several real-world EDC exports (e.g. some
Taimei/Taibo/Tongxin exports used by Chinese vaccine clinical trials) ship
with this kind of legacy artifact.

This script parses the .xlsx zip + workbook.xml + sharedStrings.xml directly,
which is robust to that failure mode and still preserves all cell text and
numeric content.

Usage
-----
    # Single file -> stdout
    py -3 scripts/extract_xlsx_full.py workbook.xlsx

    # Single file -> output file
    py -3 scripts/extract_xlsx_full.py workbook.xlsx -o dump.txt

    # Batch a folder
    py -3 scripts/extract_xlsx_full.py some_folder/ -o combined_dump.txt

Output format
-------------
One section per sheet, separated by 78 '#' chars. Each non-empty row is
written as ``R{nnnn}: cell1 | cell2 | ...`` with newlines inside cells
collapsed to `` ⏎ ``.

Dependencies
------------
stdlib only (zipfile + xml.etree.ElementTree).
"""


import contextlib
import logging
import re
import xml.etree.ElementTree as ET
import zipfile
from collections.abc import Iterator

LOG_FORMAT = "%(asctime)s [%(levelname)s] extract_xlsx_full: %(message)s"
logger = logging.getLogger("extract_xlsx_full")
ENCODING = "utf-8"

NS = {
    "n": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

# Pre-compiled regexes for cell reference parsing.
_REF_RE = re.compile(r"([A-Za-z]+)(\d+)(?::([A-Za-z]+)(\d+))?$")
_COL_LETTERS_RE = re.compile(r"([A-Za-z]+)")


def _col_letters_to_idx(s: str) -> int:
    """A -> 0, B -> 1, ..., AA -> 26, etc.

    Non-letter characters are ignored to stay robust against malformed
    cell references (e.g. ``A1$`` from legacy EDC exports).
    """
    n = 0
    for ch in s:
        u = ch.upper()
        if not ("A" <= u <= "Z"):
            continue
        n = n * 26 + (ord(u) - 64)
    return n - 1


def _split_ref(ref: str) -> tuple[int, int]:
    """Split ``A1`` or ``A1:B3`` into ((c1, r1), (c2, r2))."""
    m = _REF_RE.match(ref)
    if not m:
        return ((-1, -1), (-1, -1))
    c1 = _col_letters_to_idx(m.group(1))
    r1 = int(m.group(2)) - 1
    if m.group(3):
        c2 = _col_letters_to_idx(m.group(3))
        r2 = int(m.group(4)) - 1
    else:
        c2, r2 = c1, r1
    return ((c1, r1), (c2, r2))


def _load_shared_strings(zf: zipfile.ZipFile) -> list[str]:
    """Return the workbook's shared string table (decoded as text)."""
    try:
        data = zf.read("xl/sharedStrings.xml")
    except KeyError:
        return []
    try:
        root = ET.fromstring(data)
    except ET.ParseError as e:
        logger.warning("sharedStrings.xml parse error: %s", e)
        return []
    out: list[str] = []
    for si in root.findall("n:si", NS):
        # concatenate all <t> nodes (handles rich text via <r>)
        text = "".join(
            (t.text or "")
            for t in si.iter("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t")
        )
        out.append(text)
    return out


def _load_sheet_targets(zf: zipfile.ZipFile) -> list[tuple[str, str]]:
    """Return [(sheet_name, sheet_target_path), ...] in workbook order."""
    try:
        wb_xml = ET.fromstring(zf.read("xl/workbook.xml"))
    except ET.ParseError as e:
        logger.warning("workbook.xml parse error: %s", e)
        return []
    sheets_meta: list[tuple[str, str]] = []
    sheets = wb_xml.find("n:sheets", NS)
    if sheets is None:
        return []
    for s in sheets.findall("n:sheet", NS):
        name = s.attrib.get("name", "")
        rid = s.attrib.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id",
            "",
        )
        sheets_meta.append((name, rid))

    rels: dict[str, str] = {}
    try:
        rx = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))
    except (KeyError, ET.ParseError) as e:
        if isinstance(e, ET.ParseError):
            logger.warning("workbook.xml.rels parse error: %s", e)
        rx = None
    if rx is not None:
        for r in rx:
            rels[r.attrib.get("Id", "")] = r.attrib.get("Target", "")
    return [(name, _resolve_target(rels.get(rid, ""))) for name, rid in sheets_meta]


def _resolve_target(target: str) -> str:
    if not target:
        return ""
    if target.startswith("xl/"):
        return target
    return "xl/" + target


def _iter_sheet_rows(
    zf: zipfile.ZipFile,
    path: str,
    sst: list[str],
) -> Iterator[tuple[int, list[str]]]:
    """Yield (row_idx_1based, [cell_text_in_column_order]) for a sheet."""
    try:
        raw = zf.read(path)
    except KeyError:
        return
    try:
        ws = ET.fromstring(raw)
    except ET.ParseError as e:
        logger.warning("ParseError in %s: %s", path, e)
        return

    # Collect merged cell ranges so empty trailing cells inherit values.
    merged: list[tuple[int, int, int, int]] = []
    for mc in ws.iter("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}mergeCell"):
        ref = mc.attrib.get("ref", "")
        if ":" not in ref:
            continue
        (c1, r1), (c2, r2) = _split_ref(ref.split(":", 1)[0]), _split_ref(ref.split(":", 1)[1])
        if c1 < 0 or r1 < 0 or c2 < 0 or r2 < 0:
            continue
        merged.append((c1, r1, c2, r2))

    sheet_data = ws.find("n:sheetData", NS)
    if sheet_data is None:
        return

    for row in sheet_data.findall("n:row", NS):
        try:
            r_attr = int(row.attrib.get("r", "0"))
        except ValueError:
            continue
        row_cells: dict[int, str] = {}
        for c in row.findall("n:c", NS):
            ref = c.attrib.get("r", "")
            t = c.attrib.get("t", "n")
            v_el = c.find("n:v", NS)
            is_el = c.find("n:is", NS)
            val = ""
            if t == "s" and v_el is not None and v_el.text is not None:
                try:
                    val = sst[int(v_el.text)]
                except (ValueError, IndexError):
                    val = f"<sst!{v_el.text}>"
            elif t == "inlineStr" and is_el is not None:
                val = "".join(
                    (tt.text or "")
                    for tt in is_el.iter(
                        "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t"
                    )
                )
            elif t == "b" and v_el is not None:
                val = "TRUE" if v_el.text == "1" else "FALSE"
            elif v_el is not None:
                val = v_el.text or ""
            # extract column index
            col_m = _COL_LETTERS_RE.match(ref)
            if col_m:
                row_cells[_col_letters_to_idx(col_m.group(1))] = val

        # Apply merged ranges for this row.
        for c1, r1, c2, r2 in merged:
            if r1 <= r_attr - 1 <= r2 and c1 in row_cells:
                for cc in range(c1, c2 + 1):
                    row_cells.setdefault(cc, row_cells[c1])

        if any(v.strip() for v in row_cells.values()):
            # Sort by column index for stable column order across sheets.
            ordered = [row_cells.get(i, "") for i in sorted(row_cells)]
            yield r_attr, ordered


def extract_xlsx(xlsx_path: Path) -> str:
    """Dump every sheet of *xlsx_path* as UTF-8 text."""
    out: list[str] = [f"# FILE: {xlsx_path.name}", f"# Path: {xlsx_path}"]
    try:
        with zipfile.ZipFile(str(xlsx_path)) as zf:
            names = zf.namelist()
            sst = _load_shared_strings(zf)
            sheets = _load_sheet_targets(zf)
            out.append(f"# Sheets ({len(sheets)}): {[s[0] for s in sheets]}")
            for sheet_name, target in sheets:
                out.append("")
                out.append("#" * 78)
                out.append(f"# Sheet: {sheet_name}  path={target}")
                out.append("#" * 78)
                if not target or target not in names:
                    out.append(f"# (missing sheet path: {target})")
                    continue
                for r_idx, cells in _iter_sheet_rows(zf, target, sst):
                    cells_out = [c.replace("\n", " ⏎ ") for c in cells]
                    out.append(f"R{r_idx:04d}: " + " | ".join(cells_out))
    except zipfile.BadZipFile as e:
        logger.error("not a valid .xlsx (bad zip): %s: %s", xlsx_path, e)
        return f"# FILE: {xlsx_path.name}\n# ERROR: bad zip: {e}"
    except OSError as e:
        logger.error("I/O error reading %s: %s", xlsx_path, e)
        return f"# FILE: {xlsx_path.name}\n# ERROR: I/O: {e}"
    return "\n".join(out)


def _iter_xlsx_files(target: Path) -> Iterator[Path]:
    if target.is_file():
        yield target
        return
    if not target.is_dir():
        return
    for entry in sorted(target.iterdir()):
        if entry.suffix.lower() == ".xlsx" and not entry.name.startswith("~"):
            yield entry


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Robust .xlsx -> UTF-8 text dumper (zip+xml).")
    parser.add_argument("target", help=".xlsx file or folder containing .xlsx files")
    parser.add_argument("-o", "--output", default=None, help="output file (default: stdout)")
    parser.add_argument("--log-level", default="WARNING")
    args = parser.parse_args(argv)

    logging.basicConfig(
        level=getattr(logging, args.log_level.upper(), logging.WARNING),
        format=LOG_FORMAT,
    )

    target = Path(args.target)
    if not target.exists():
        logger.error("Path not found: %s", target)
        return 2

    chunks: list[str] = []
    for xlsx in _iter_xlsx_files(target):
        logger.info("processing %s", xlsx)
        chunks.append(extract_xlsx(xlsx))

    text = "\n".join(chunks)
    if args.output:
        out_path = Path(args.output)
        try:
            out_path.parent.mkdir(parents=True, exist_ok=True)
            out_path.write_text(text, encoding=ENCODING)
        except OSError as e:
            logger.error("failed to write %s: %s", out_path, e)
            return 1
        print(f"OK: wrote {args.output} ({len(text):,} chars)")
    else:
        with contextlib.suppress(Exception):  # some streams are not reconfigurable
            sys.stdout.reconfigure(encoding=ENCODING, errors="replace")
        sys.stdout.write(text)
    return 0


def extract_office_text(file_path):
    ext = Path(file_path).suffix.lower()
    if ext == ".docx":
        return extract_docx(Path(file_path))
    elif ext == ".doc":
        return extract_doc_legacy(Path(file_path))
    elif ext == ".pptx":
        out_path = Path(file_path).with_suffix(".txt")
        extract_full(file_path, str(out_path))
        return "PPTX extracted to " + str(out_path)
    elif ext == ".xlsx":
        return extract_xlsx(Path(file_path))
    else:
        raise ValueError("Unsupported format: " + ext)
