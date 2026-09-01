"""V7 自检脚本：核验 Slide 8 新行内容、链接、格式、整体链接计数"""

from pptx import Presentation
from pptx.util import Emu

P = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v7.pptx"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

prs = Presentation(P)
print("=== 全局 ===")
print(f"slides: {len(prs.slides._sldIdLst)}")

# 1. Slide 8 表格结构
slide = prs.slides[7]
print("\n=== Slide 8 ===")
for sh in slide.shapes:
    if sh.has_table:
        tb = sh.table
        print(f"rows={len(tb.rows)} cols={len(tb.columns)}")
        print(f"row heights: {[round(Emu(r.height).inches, 2) for r in tb.rows]}")
        # 新行 row 3
        print("\n--- row 3 (新行) 内容 ---")
        for ci in range(5):
            cell = tb.cell(3, ci)
            paras = []
            for para in cell.text_frame.paragraphs:
                runs = []
                for run in para.runs:
                    hlink = run._r.find(".//a:hlinkClick", {"a": A_NS, "r": R_NS})
                    link_txt = ""
                    if hlink is not None:
                        rid = hlink.get(f"{{{R_NS}}}id")
                        link_txt = f"[->{rid}]"
                    runs.append(run.text + link_txt)
                paras.append("".join(runs))
            print(f"c{ci}:")
            for p in paras:
                print(f"   {p}")
        # 检查新行每个单元格字体 sz
        print("\n--- row 3 字体检查 (应全部 sz=1000) ---")
        ok = True
        for ci in range(5):
            for run in tb.cell(3, ci).text_frame.paragraphs:
                for r in run.runs:
                    rPr = r._r.rPr
                    if rPr is not None and rPr.get("sz") != "1000":
                        ok = False
                        print(f"  BAD sz in c{ci}: {rPr.get('sz')} text={r.text[:30]}")
        print("  ALL sz=1000 OK" if ok else "  FONT MISMATCH FOUND")

# 2. Slide 8 rels 校验
print("\n=== Slide 8 hyperlink rels ===")
targets = {}
for rid, rel in slide.part.rels.items():
    if "hyperlink" in rel.reltype:
        targets[rid] = str(rel.target_ref)
for rid in sorted(targets):
    print(f"  {rid}: {targets[rid]}")
expect = {
    "rId13": "https://clinicaltrials.gov/study/NCT04193189",
    "rId14": "https://clinicaltrials.gov/study/NCT01195246",
    "rId15": "https://pubmed.ncbi.nlm.nih.gov/39616603/",
    "rId16": "https://pubmed.ncbi.nlm.nih.gov/36269938/",
}
print("\n--- 新链接期望核对 ---")
allok = True
for rid, url in expect.items():
    got = targets.get(rid)
    status = "OK" if got == url else f"MISMATCH got={got}"
    if got != url:
        allok = False
    print(f"  {rid} {status}")
print("  ALL NEW LINKS OK" if allok else "  LINK ERROR!")

# 3. 全文件链接计数（跨 slide 统计）
print("\n=== 全文件超链接统计 ===")
total = 0
for si, s in enumerate(prs.slides):
    n = 0
    for rel in s.part.rels.values():
        if "hyperlink" in rel.reltype:
            n += 1
    if n:
        print(f"  Slide {si + 1}: {n} links")
    total += n
print(f"  TOTAL hyperlink rels: {total}")

# 4. 其他 slide 内容未受影响（抽查 slide 7/9 行数）
for idx in (6, 8):
    s = prs.slides[idx]
    for sh in s.shapes:
        if sh.has_table:
            print(f"\nSlide {idx + 1} table rows={len(sh.table.rows)} (应为 3)")
