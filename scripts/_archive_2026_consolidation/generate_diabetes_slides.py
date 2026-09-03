# -*- coding: utf-8 -*-
"""
TVAX-009 新增 6 页「有/无糖尿病病史」免疫原性分析 PPT

- 读取 diabetes_immuno_result.json（阳转率 + GMC 的 FAS/PPS 结果）
- 参照 32-34 / 38-40 页版式，构建「方案 A」7 组表格（无 P 值列）
- 克隆参考页，替换表格 + 标题 + 描述文字
- 插入原第 42 页与第 43 页之间（0-based 索引 42..47）
- 另存为新文件（不动原始 PPT）
"""
import json, copy
from lxml import etree
from pptx import Presentation

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def qn(t):
    return '{%s}%s' % (A_NS, t)


def pn(t):
    return '{%s}%s' % (P_NS, t)


BASE = 'E:/Cursor Project/2-Scientific-Skills-for-Clinical_Trial/review_materials'
SRC = BASE + '/TVAX-009项目3期临床试验启动前沟通交流ppt-20260902（临床部分）.pptx'
OUT = BASE + '/TVAX-009项目3期临床试验启动前沟通交流ppt-20260903（临床部分-新增糖尿病亚组）.pptx'

with open(BASE + '/diabetes_immuno_result.json', encoding='utf-8') as f:
    DATA = json.load(f)

FAS = DATA['fas']
PPS = DATA['pps']

# 方案 A：7 组全列（0,1月 A1/A2；0,2月 B1/B2；0,1,6月 C1/C3/C2）
GROUPS = [
    ('A1', '低剂量组(A1)', '0,1月程序', None),
    ('A2', '高剂量组(A2)', '0,1月程序', None),
    ('B1', '低剂量组(B1)', '0,2月程序', None),
    ('B2', '高剂量组(B2)', '0,2月程序', None),
    ('C1', '阳性对照组(C1)', '0,1,6月程序', '18-59岁'),
    ('C3', '0,1,6月高剂量组(C3)', '0,1,6月程序', '≥60岁'),
    ('C2', '阳性对照组(C2)', '0,1,6月程序', '≥60岁'),
]
PROGRAMS = [
    ('0,1月程序', ['A1', 'A2']),
    ('0,2月程序', ['B1', 'B2']),
    ('0,1,6月程序', ['C1', 'C3', 'C2']),
]
PROG_OF = {code: prog for code, _, prog, _ in GROUPS}
LABEL = {code: lab for code, lab, _, _ in GROUPS}
ANNO = {code: a for code, _, _, a in GROUPS}

PROG_TIMEPOINTS = {
    '0,1月程序': ['M1', 'M2', 'M3', 'M6', 'M7', 'M8'],
    '0,2月程序': ['M1', 'M2', 'M3', 'M4', 'M7', 'M8'],
    '0,1,6月程序': ['M1', 'M2', 'M6', 'M7', 'M8'],
}
ALL_TP = ['M1', 'M2', 'M3', 'M4', 'M6', 'M7', 'M8']

REL_MAP = {
    '全免后1个月': {'0,1月程序': 'M2', '0,2月程序': 'M3', '0,1,6月程序': 'M7'},
    '全免后2个月': {'0,1月程序': 'M3', '0,2月程序': 'M4', '0,1,6月程序': 'M8'},
}


# ---------------- XML 片段 ----------------
def _esc(s):
    return s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')


LN_T = ('<a:lnT w="12700" cap="flat" cmpd="sng" algn="ctr">'
        '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
        '<a:prstDash val="solid"/><a:round/>'
        '<a:headEnd type="none" w="med" len="med"/>'
        '<a:tailEnd type="none" w="med" len="med"/></a:lnT>')


def _ln_b(solid):
    if solid:
        return ('<a:lnB w="12700" cap="flat" cmpd="sng" algn="ctr">'
                '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
                '<a:prstDash val="solid"/><a:round/>'
                '<a:headEnd type="none" w="med" len="med"/>'
                '<a:tailEnd type="none" w="med" len="med"/></a:lnB>')
    return '<a:lnB><a:noFill/></a:lnB>'


def make_tc(text, bold=False, grid_span=None, row_span=None, lnB_solid=False, lnT=True):
    attrs = ''
    if grid_span:
        attrs += ' gridSpan="%d"' % grid_span
    if row_span:
        attrs += ' rowSpan="%d"' % row_span
    b = ' b="1"' if bold else ''
    s = '<a:tc xmlns:a="%s"%s><a:txBody><a:bodyPr/><a:lstStyle/>' % (A_NS, attrs)
    for line in text.split('\n'):
        s += ('<a:p><a:pPr indent="0" algn="ctr">'
              '<a:lnSpc><a:spcPct val="100000"/></a:lnSpc><a:buNone/></a:pPr>'
              '<a:r><a:rPr lang="zh-CN" altLang="en-US" sz="800" kern="0"%s>'
              '<a:solidFill><a:schemeClr val="tx1"/></a:solidFill><a:effectLst/>'
              '<a:latin typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
              '<a:ea typeface="微软雅黑" panose="020B0503020204020204" pitchFamily="34" charset="-122"/>'
              '<a:cs typeface="Times New Roman" panose="02020603050405020304" charset="0"/>'
              '<a:sym typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
              '</a:rPr><a:t>%s</a:t></a:r></a:p>' % (b, _esc(line)))
    s += '</a:txBody>'
    s += '<a:tcPr marL="36195" marR="0" marT="0" marB="0" anchor="ctr">'
    s += '<a:lnL><a:noFill/></a:lnL><a:lnR><a:noFill/></a:lnR>'
    if lnT:
        s += LN_T
    s += _ln_b(lnB_solid)
    s += '<a:noFill/></a:tcPr></a:tc>'
    return s


def table_str(col_widths, rows):
    s = '<a:tbl xmlns:a="%s"><a:tblPr firstRow="1" firstCol="1" bandRow="1"/><a:tblGrid>' % A_NS
    for w in col_widths:
        s += '<a:gridCol w="%d"/>' % w
    s += '</a:tblGrid>'
    for h, cells in rows:
        s += '<a:tr h="%d">%s</a:tr>' % (h, ''.join(cells))
    s += '</a:tbl>'
    return s


# ---------------- 取值 ----------------
def fmt_num(v):
    return '/' if v is None else '%.2f' % v


def get_vals(data, code, tag, tp, mode):
    rec = data[code][tag][tp]
    if mode == 'sero':
        return [str(rec['N_阳转']), str(rec['n_阳转']), fmt_num(rec['pct'])]
    return [str(rec['N_gmc']), fmt_num(rec['gmc'])]


# ---------------- 表格构建 ----------------
def build_absolute(data, mode):
    """绝对时间点表：19 行（有/无 × M1-M8，含空行分隔）"""
    if mode == 'sero':
        lead_w = [504000, 612000]
        sub_w = [310000, 310000, 410000]
        sub_headers = ['N', 'n', '%']
        header_h = 190500
        you_h = 175260
        wu_h = 165100
        blank_h = [37669, 64303]
    else:
        lead_w = [512551, 612000]
        sub_w = [420000, 580000]
        sub_headers = ['例数', 'GMC']
        header_h = 214250
        you_h = 169564
        wu_h = 169258
        blank_h = [36000, 36000]

    lead_headers = ['糖尿病病史', '时间点\n(首剂接种后）']
    col_widths = list(lead_w)
    for _, codes in PROGRAMS:
        for _ in codes:
            col_widths.extend(sub_w)
    n_sub = sum(len(codes) * len(sub_w) for _, codes in PROGRAMS)

    rows = []
    # R0：前导列(rowSpan 3) + 程序列头(gridSpan)
    cells = [make_tc(lead_headers[0], bold=True, row_span=3, lnB_solid=True),
             make_tc(lead_headers[1], bold=True, row_span=3, lnB_solid=True)]
    for pname, codes in PROGRAMS:
        cells.append(make_tc(pname, bold=True, grid_span=len(codes) * len(sub_w), lnB_solid=True))
    rows.append((header_h, cells))

    # R1：组别列头
    cells = []
    for _, codes in PROGRAMS:
        for code in codes:
            txt = LABEL[code]
            if ANNO[code]:
                txt += '\n(' + ANNO[code] + ')'
            cells.append(make_tc(txt, bold=True, grid_span=len(sub_w), lnB_solid=True))
    rows.append((header_h, cells))

    # R2：子列头
    cells = []
    for _, codes in PROGRAMS:
        for _ in codes:
            for sh in sub_headers:
                cells.append(make_tc(sh, bold=True, lnB_solid=True))
    rows.append((header_h, cells))

    # 数据段：有 / 空行 / 无
    for tag, tag_label, h in [('有', '有', you_h), ('无', '无', wu_h)]:
        for i, tp in enumerate(ALL_TP):
            last_row = (tag == '无' and tp == 'M8')
            cells = []
            if i == 0:
                cells.append(make_tc(tag_label, bold=False, row_span=7, lnB_solid=False))
            cells.append(make_tc(tp, bold=False, lnB_solid=last_row))
            for _, codes in PROGRAMS:
                for code in codes:
                    if tp in PROG_TIMEPOINTS[PROG_OF[code]]:
                        for v in get_vals(data, code, tag, tp, mode):
                            cells.append(make_tc(v, bold=False, lnB_solid=last_row))
                    else:
                        for _ in sub_headers:
                            cells.append(make_tc('/', bold=False, lnB_solid=last_row))
            rows.append((h, cells))
        if tag == '有':
            for bh in blank_h:
                cells = [make_tc('', bold=False, lnB_solid=False, lnT=False),
                         make_tc('', bold=False, lnB_solid=False, lnT=False)]
                for _ in range(n_sub):
                    cells.append(make_tc('', bold=False, lnB_solid=False, lnT=False))
                rows.append((bh, cells))

    return table_str(col_widths, rows), sum(col_widths)


def build_relative(data, mode, analysis_set):
    """相对时间点表：7 行（全免后 1/2 个月 × 有/无）"""
    if mode == 'sero':
        lead_w = [430000, 430000, 380000]
        sub_w = [300000, 300000, 350000]
        sub_headers = ['N', 'n', '%']
        header_h = [192405, 191770, 193040]
        data_h = [268605, 269240, 269240, 269240]
    else:
        lead_w = [430000, 430000, 380000]
        sub_w = [350000, 580000]
        sub_headers = ['例数', 'GMC']
        header_h = [206459, 206459, 206459]
        data_h = [231234, 231234, 231234, 231234]

    tp_w = 360000
    lead_headers = ['分析集', '糖尿病病史', '时间点\n(全免后）']

    col_widths = list(lead_w)
    for _, codes in PROGRAMS:
        col_widths.append(tp_w)
        for _ in codes:
            col_widths.extend(sub_w)

    rows = []
    # R0
    cells = [make_tc(lead_headers[0], bold=True, row_span=3, lnB_solid=True),
             make_tc(lead_headers[1], bold=True, row_span=3, lnB_solid=True),
             make_tc(lead_headers[2], bold=True, row_span=3, lnB_solid=True)]
    for pname, codes in PROGRAMS:
        cells.append(make_tc(pname, bold=True, grid_span=1 + len(codes) * len(sub_w), lnB_solid=True))
    rows.append((header_h[0], cells))

    # R1
    cells = []
    for _, codes in PROGRAMS:
        cells.append(make_tc('对应的时间点', bold=True, row_span=2, lnB_solid=True))
        for code in codes:
            txt = LABEL[code]
            if ANNO[code]:
                txt += '\n(' + ANNO[code] + ')'
            cells.append(make_tc(txt, bold=True, grid_span=len(sub_w), lnB_solid=True))
    rows.append((header_h[1], cells))

    # R2
    cells = []
    for _, codes in PROGRAMS:
        for _ in codes:
            for sh in sub_headers:
                cells.append(make_tc(sh, bold=True, lnB_solid=True))
    rows.append((header_h[2], cells))

    # 数据 4 行
    ri = 0
    for di, (tag, tag_label) in enumerate([('有', '有'), ('无', '无')]):
        for li, tp_label in enumerate(['全免后1个月', '全免后2个月']):
            last_row = (di == 1 and li == 1)
            cells = []
            if li == 0:
                if di == 0:
                    cells.append(make_tc(analysis_set, bold=False, row_span=4, lnB_solid=last_row))
                cells.append(make_tc(tag_label, bold=False, row_span=2, lnB_solid=last_row))
            cells.append(make_tc(tp_label, bold=False, lnB_solid=last_row))
            for pname, codes in PROGRAMS:
                m = REL_MAP[tp_label][pname]
                cells.append(make_tc(m, bold=False, lnB_solid=last_row))
                for code in codes:
                    for v in get_vals(data, code, tag, m, mode):
                        cells.append(make_tc(v, bold=False, lnB_solid=last_row))
            rows.append((data_h[ri], cells))
            ri += 1

    return table_str(col_widths, rows), sum(col_widths)


# ---------------- 幻灯片处理 ----------------
def duplicate_slide(prs, src_index):
    src = prs.slides[src_index]
    new = prs.slides.add_slide(src.slide_layout)
    for shape in list(new.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(shape._element))
    return new


def replace_text(shape, new_text, bold=True):
    tf = shape.text_frame
    first_rpr = None
    for p in tf.paragraphs:
        if p.runs:
            rpr = p.runs[0]._r.find(qn('rPr'))
            if rpr is not None:
                first_rpr = copy.deepcopy(rpr)
            break
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    run = p0.add_run()
    run.text = new_text
    if first_rpr is not None:
        if not bold:
            first_rpr.attrib.pop('b', None)
        existing = run._r.find(qn('rPr'))
        if existing is not None:
            run._r.remove(existing)
        run._r.insert(0, first_rpr)


def replace_table(shape, tbl_el, new_width):
    gf = shape._element
    for old in gf.findall('.//' + qn('tbl')):
        old.getparent().remove(old)
    gd = gf.find('.//' + qn('graphicData'))
    gd.append(tbl_el)
    ext = gf.find(pn('xfrm')).find(qn('ext'))
    ext.set('cx', str(int(new_width)))


def modify_cloned_slide(new_slide, title, desc, tbl_specs):
    ti = 0
    for shape in new_slide.shapes:
        if shape.has_table:
            if ti < len(tbl_specs):
                tbl_el, w = tbl_specs[ti]
                replace_table(shape, tbl_el, w)
                ti += 1
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text
            if text.strip() == '2期免疫原性结果':
                continue
            if '_绝对时间点' in text or '相对时间点' in text:
                replace_text(shape, title, bold=True)
            else:
                replace_text(shape, desc, bold=False)


# ---------------- 主流程 ----------------
def main():
    prs = Presentation(SRC)

    desc_dist = ('本分析纳入基线合并糖尿病病史的受试者共 30 例，其中 18-59 岁 18 例'
                 '（0,1月低剂量组 5 例、0,1月高剂量组 3 例、0,2月低剂量组 2 例、'
                 '0,2月高剂量组 3 例、阳性对照组 5 例）、≥60 岁 12 例'
                 '（0,1,6月高剂量组 8 例、阳性对照组 4 例）。')

    specs = [
        # 1. 阳转率 PPS 绝对时间点
        {'src': 31,
         'title': '有、无糖尿病病史人群接种后抗-HBs阳转率（PPS）_绝对时间点',
         'desc': desc_dist + '结果显示：有糖尿病病史人群完成全程 2 剂/3 剂接种后的阳转率总体达到较高水平，'
                             '0,2月程序两组在 M3 均达 100.00%，0,1,6月高剂量组及两个阳性对照组在 M7 均达 100.00%。',
         'tables': [build_absolute(PPS, 'sero')]},
        # 2. 阳转率 FAS 绝对时间点
        {'src': 32,
         'title': '有、无糖尿病病史人群接种后抗-HBs阳转率（FAS）_绝对时间点',
         'desc': '基于 FAS 分析的趋势与 PPS 一致，各组有、无糖尿病病史人群完成全程 2 剂/3 剂接种后的阳转率均达到较高水平。',
         'tables': [build_absolute(FAS, 'sero')]},
        # 3. 阳转率 相对时间点（PPS + FAS）
        {'src': 33,
         'title': '有、无糖尿病病史人群接种后抗-HBs阳转率（相对时间点）',
         'desc': '全免后 1 个月、2 个月，各组有、无糖尿病病史人群的阳转率均维持在较高水平，FAS 与 PPS 趋势一致。',
         'tables': [build_relative(PPS, 'sero', 'PPS'), build_relative(FAS, 'sero', 'FAS')]},
        # 4. GMC PPS 绝对时间点
        {'src': 37,
         'title': '有、无糖尿病病史人群接种后抗-HBs GMC（PPS）_绝对时间点',
         'desc': '本分析纳入基线合并糖尿病病史的受试者共 30 例，其中 18-59 岁 18 例、≥60 岁 12 例。'
                 '结果显示：各组有、无糖尿病病史人群完成全程 2 剂/3 剂接种后各时间点的抗体 GMC 均呈现上升趋势；'
                 '抗-HBs 浓度低于检测下限（<2.00 mIU/mL）的样本按 LLOQ/2（1.00 mIU/mL）计入 GMC 计算。',
         'tables': [build_absolute(PPS, 'gmc')]},
        # 5. GMC FAS 绝对时间点
        {'src': 38,
         'title': '有、无糖尿病病史人群接种后抗-HBs GMC（FAS）_绝对时间点',
         'desc': '基于 FAS 分析的趋势与 PPS 一致，各组有、无糖尿病病史人群各时间点抗体 GMC 的变化趋势相似。',
         'tables': [build_absolute(FAS, 'gmc')]},
        # 6. GMC 相对时间点（PPS + FAS）
        {'src': 39,
         'title': '有、无糖尿病病史人群接种后抗-HBs GMC（相对时间点）',
         'desc': '全免后 1 个月、2 个月，各组有、无糖尿病病史人群的抗体 GMC 均维持在一定水平，FAS 与 PPS 趋势一致。',
         'tables': [build_relative(PPS, 'gmc', 'PPS'), build_relative(FAS, 'gmc', 'FAS')]},
    ]

    for spec in specs:
        new = duplicate_slide(prs, spec['src'])
        tbl_specs = [(etree.fromstring(ts), w) for ts, w in spec['tables']]
        modify_cloned_slide(new, spec['title'], spec['desc'], tbl_specs)

    # 重排：将末尾新增 6 页移到 0-based 索引 42（原第 43 页之前）
    sldIdLst = prs.slides._sldIdLst
    all_ids = list(sldIdLst)
    new_ids = all_ids[-6:]
    for sid in new_ids:
        sldIdLst.remove(sid)
    for i, sid in enumerate(new_ids):
        sldIdLst.insert(42 + i, sid)

    prs.save(OUT)
    print('saved:', OUT)
    print('total slides:', len(prs.slides))


if __name__ == '__main__':
    main()
