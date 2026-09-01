#!/usr/bin/env python3
"""Add small clinical-semantic icons to the 4 phase-card title bars on slide 2
of the user-optimized TVAX-006 overview PPTX (in-place, preserving user edits).

Icons (white, on the dark-red title bar, left of each phase title):
  Phase I  Australia -> globe   (white circle + red cross = overseas/global)
  Phase I  China     -> shield  (freeform shield = safety)
  Phase II China     -> molecule(hexagon = immunogenicity)
  Phase III China    -> up arrow(efficacy)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

DARK_RED = RGBColor(0xC0, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PPTX = r"E:/Cursor Project/2-Scientific-Skills-for-Clinical_Trial/review_materials/006 PPT/TVAX-006_Clinical_Development_Overview_20260831.pptx"

ICON = 0.3      # icon square side, inches
BAR_Y = 1.25    # title-bar top, inches
BAR_H = 0.52    # title-bar height, inches

# (phase title text, icon kind)
TITLES = [
    ("Phase I — Australia", "globe"),
    ("Phase I — China", "shield"),
    ("Phase II — China", "molecule"),
    ("Phase III — China", "arrow"),
]


def _fill(sp, rgb):
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb


def _no_line(sp):
    sp.line.fill.background()
    sp.shadow.inherit = False


def add_shape(slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        _fill(sp, fill)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_icon(slide, x, y, kind):
    s = ICON
    if kind == "globe":
        # white disc + red longitude/latitude cross
        add_shape(slide, x, y, s, s, WHITE, MSO_SHAPE.OVAL)
        add_shape(slide, x + 0.04, y + 0.138, 0.22, 0.024, DARK_RED)   # equator
        add_shape(slide, x + 0.138, y + 0.04, 0.024, 0.22, DARK_RED)   # meridian
    elif kind == "shield":
        fb = slide.shapes.build_freeform(Inches(x), Inches(y))
        fb.add_line_segments([
            (Inches(x + s), Inches(y)),           # top-right
            (Inches(x + s), Inches(y + 0.20)),    # right shoulder
            (Inches(x + s / 2), Inches(y + s)),   # bottom point
            (Inches(x), Inches(y + 0.20)),        # left shoulder
        ], close=True)
        sp = fb.convert_to_shape()
        _fill(sp, WHITE)
        _no_line(sp)
    elif kind == "molecule":
        add_shape(slide, x, y, s, s, WHITE, MSO_SHAPE.HEXAGON)
    elif kind == "arrow":
        add_shape(slide, x, y, s, s, WHITE, MSO_SHAPE.UP_ARROW)


def main():
    prs = Presentation(PPTX)
    slide = prs.slides[1]

    by_text = {}
    wanted = {t for t, _ in TITLES}
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() in wanted:
            by_text[sh.text_frame.text.strip()] = sh

    for title_text, kind in TITLES:
        tb = by_text[title_text]
        x_in = tb.left.inches
        icon_x = x_in - 0.02               # just inside the bar's left padding
        icon_y = BAR_Y + (BAR_H - ICON) / 2
        add_icon(slide, icon_x, icon_y, kind)
        # shift the title text right to clear the icon
        tb.left = Inches(x_in + 0.40)
        tb.width = Inches(tb.width.inches - 0.40)

    prs.save(PPTX)
    print("SAVED:", PPTX)
    print("icons added:", len(TITLES))


if __name__ == "__main__":
    main()
