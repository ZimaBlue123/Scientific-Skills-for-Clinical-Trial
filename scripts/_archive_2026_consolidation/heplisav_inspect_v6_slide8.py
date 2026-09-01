"""检查 V6 Slide 7/8/9 表格结构，确定特定人群行现状"""

import sys

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.oxml.ns import qn

V6 = r"E:/Cursor Project/2-Scientific-Skills-for-Clinical_Trial/review_materials/60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v6.pptx"
prs = Presentation(V6)


def cell_links(cell):
    out = []
    for r_elem in cell._tc.iter(qn("a:r")):
        rPr = r_elem.find(qn("a:rPr"))
        if rPr is None:
            continue
        hl = rPr.find(qn("a:hlinkClick"))
        if hl is None:
            continue
        rid = hl.get(qn("r:id"))
        rtext = "".join(t.text or "" for t in r_elem.findall(qn("a:t")))
        url = ""
        if rid:
            try:
                url = cell.part.rels[rid].target_ref
            except KeyError:
                pass
        if url.startswith("http") and rtext.strip():
            out.append((rtext.strip(), url))
    return out


for idx in (6, 7, 8):
    slide = prs.slides[idx]
    print("=" * 100)
    print("SLIDE %d" % (idx + 1))
    print("=" * 100)
    for shp in slide.shapes:
        if shp.has_text_frame and not shp.has_table and shp.text_frame.text.strip():
            print("  [标题] %s" % shp.text_frame.text.replace("\n", " ")[:60])
        if shp.has_table:
            tbl = shp.table
            print("  表格 %d 行 x %d 列" % (len(tbl.rows), len(tbl.columns)))
            for ri in range(len(tbl.rows)):
                cell0 = tbl.cell(ri, 0).text_frame.text.replace("\n", " ")[:35]
                cell4 = (
                    tbl.cell(ri, 4).text_frame.text.replace("\n", " | ")[:80]
                    if len(tbl.columns) > 4
                    else ""
                )
                links = cell_links(tbl.cell(ri, 4)) if len(tbl.columns) > 4 else []
                print("  行%d: [%s] | 参考: %s" % (ri, cell0, cell4))
                for rt, u in links:
                    print("        ↳ [%r] -> %s" % (rt[:35], u))
