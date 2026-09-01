"""
V8 生成脚本：基于 v7，按用户要求重组 Slide 8 + 新增 Slide 9
  1. Slide 8：保留表头 + 综合数据行 + 特定人群行(HBV-17/19)，补充样本量/分组/研究终点；
     删除原整合行(BEe-HIVe+HBV-18)
  2. Slide 9（新页，插在 Slide 8 之后）：表头 + BEe-HIVe 一行 + HBV-18 一行
     - 两个试验拆成两行
     - 补充研究终点
     - 分组单独列、样本量独立描述
  3. 链接全部嵌入，格式与 v7 行2 对齐（Arial/微软雅黑 10pt、白底、边距、lnSpc 100%）
"""

import copy
import sys
from xml.sax.saxutils import escape

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches

SRC = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v7.pptx"
DST = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v8.pptx"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
RT_HYPERLINK = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"

RPR_TMPL = (
    '<a:rPr xmlns:a="{a}" xmlns:r="{r}" sz="1000" dirty="0">'
    '<a:latin typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
    '<a:ea typeface="微软雅黑" panose="020B0503020204020204" pitchFamily="34" charset="-122"/>'
    '<a:sym typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
    "{hlink}"
    "</a:rPr>"
)


def make_rpr(rid=None):
    hlink = f'<a:hlinkClick r:id="{rid}"/>' if rid else ""
    return RPR_TMPL.format(a=A_NS, r=R_NS, hlink=hlink)


def make_p(segments):
    """segments: list of (text, rid_or_None)"""
    runs = ""
    for text, rid in segments:
        runs += f'<a:r>{make_rpr(rid)}<a:t xml:space="preserve">{escape(text)}</a:t></a:r>'
    xml = (
        f'<a:p xmlns:a="{A_NS}" xmlns:r="{R_NS}">'
        '<a:pPr><a:lnSpc><a:spcPct val="100000"/></a:lnSpc>'
        '<a:spcBef><a:spcPts val="100"/></a:spcBef>'
        '<a:spcAft><a:spcPts val="100"/></a:spcAft></a:pPr>'
        f"{runs}</a:p>"
    )
    return etree.fromstring(xml)


def set_cell(cell, paras, part):
    """paras: list of segments-lists; 清空 cell 并写入段落"""
    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    for child in list(txBody):
        if child.tag not in (qn("a:bodyPr"), qn("a:lstStyle")):
            txBody.remove(child)
    for segs in paras:
        txBody.append(make_p(segs))
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))
        tc.insert(0, tcPr)
    tcPr.set("marL", "50800")
    tcPr.set("marR", "50800")
    tcPr.set("marT", "19050")
    tcPr.set("marB", "19050")
    for sf in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(sf)
    fill = etree.SubElement(tcPr, qn("a:solidFill"))
    clr = etree.SubElement(fill, qn("a:srgbClr"))
    clr.set("val", "FFFFFF")


def add_row_xml(table, height_inches):
    """XML 层面追加一行：深拷贝最后一行模板（继承 tcPr 格式），清空文本"""
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    if not trs:
        sys.exit("ERROR: no <a:tr> in table")
    new_tr = copy.deepcopy(trs[-1])
    new_tr.set("h", str(Emu(Inches(height_inches))))
    for tc in new_tr.findall(qn("a:tc")):
        txBody = tc.find(qn("a:txBody"))
        if txBody is not None:
            for child in list(txBody):
                if child.tag not in (qn("a:bodyPr"), qn("a:lstStyle")):
                    txBody.remove(child)
    tbl.append(new_tr)
    return table.rows[len(table.rows) - 1]


def delete_row(table, index):
    """删除表格第 index 行"""
    trs = table._tbl.findall(qn("a:tr"))
    if index >= len(trs):
        sys.exit(f"ERROR: delete_row index {index} out of range ({len(trs)})")
    trs[index].getparent().remove(trs[index])


def set_row_height(table, index, height_inches):
    trs = table._tbl.findall(qn("a:tr"))
    trs[index].set("h", str(Emu(Inches(height_inches))))


def find_table(slide):
    for sh in slide.shapes:
        if sh.has_table:
            return sh.table
    return None


def duplicate_slide(prs, index):
    """深拷贝 index 幻灯片为新 slide（含超链接 rels 与 rId 映射），追加到末尾"""
    source = prs.slides[index]
    slide = prs.slides.add_slide(source.slide_layout)
    # 删除 add_slide 自动生成的占位符
    for shp in list(slide.shapes):
        shp._element.getparent().remove(shp._element)
    # 深拷贝 spTree 子元素（跳过组属性）
    # NOTE: spTree 子元素是 p: 命名空间（presentationml），不能用 qn('a:nvGrpSpPr')（drawingml）
    # 否则跳过判断永远为 False，导致源 slide 的 nvGrpSpPr/grpSpPr 被重复拷贝，
    # 违反 CT_GroupShape schema（nvGrpSpPr/grpSpPr 只能各出现一次），PowerPoint 会拒绝打开。
    spTree = source.shapes._spTree
    for child in list(spTree):
        if child.tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            continue
        slide.shapes._spTree.append(copy.deepcopy(child))
    # 复制超链接 rels 并建立 rId 映射
    rmap = {}
    for rid, rel in source.part.rels.items():
        if rel.reltype == RT_HYPERLINK:
            rmap[rid] = slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
    # 更新 hlinkClick 的 rId 指向
    for hlink in slide.shapes._spTree.iter(qn("a:hlinkClick")):
        old = hlink.get(qn("r:id"))
        if old in rmap:
            hlink.set(qn("r:id"), rmap[old])
    return slide


def move_slide(prs, old_index, new_index):
    """移动幻灯片 old_index -> new_index"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[old_index]
    xml_slides.remove(el)
    if new_index > old_index:
        new_index -= 1
    xml_slides.insert(new_index, el)


def cleanup_orphan_hyperlink_rels(slide):
    """删除 slide 中不再被任何 hlinkClick 引用的超链接 rel"""
    used = set()
    for hlink in slide.shapes._spTree.iter(qn("a:hlinkClick")):
        rid = hlink.get(qn("r:id"))
        if rid:
            used.add(rid)
    removed = []
    for rid in list(slide.part.rels.keys()):
        rel = slide.part.rels[rid]
        if rel.reltype == RT_HYPERLINK and rid not in used:
            slide.part.drop_rel(rid)
            removed.append(rid)
    return removed


def set_textbox_text(shape, text):
    """改写 TextBox 第一个段落文本（保留格式）"""
    tf = shape.text_frame
    ps = list(tf.paragraphs)
    for p in ps[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    runs = p0.runs
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.add_run().text = text


def main():
    prs = Presentation(SRC)

    # ================= 1. 先复制 Slide 8（4 行原始状态）到新页 =================
    new_slide = duplicate_slide(prs, 7)
    move_slide(prs, len(prs.slides) - 1, 8)  # 插到第 9 页位置
    s9 = prs.slides[8]
    t9 = find_table(s9)
    if t9 is None:
        sys.exit("ERROR: new slide table not found")

    # ---- Slide 9 表格改造：删综合行(1)、特定人群行(2) -> 表头 + 整合行 ----
    delete_row(t9, 2)
    delete_row(t9, 1)

    # ---- Slide 9 标题 ----
    for sh in s9.shapes:
        if sh.has_text_frame and not sh.has_table:
            set_textbox_text(
                sh,
                "【已上市产品】CpG佐剂预防性疫苗安全性汇总（HEPLISAV-B：Phase 3 特定人群：HIV感染者、透析患者）",
            )

    # ---- Slide 9 行1：BEe-HIVe；行2：HBV-18 ----
    part9 = s9.part
    rid9 = {}
    urls9 = {
        "nct_beehive": "https://clinicaltrials.gov/study/NCT04193189",
        "nct_hbv18": "https://clinicaltrials.gov/study/NCT01195246",
        "pmid_jama": "https://pubmed.ncbi.nlm.nih.gov/39616603/",
        "pmid_girndt": "https://pubmed.ncbi.nlm.nih.gov/36269938/",
    }
    for key, url in urls9.items():
        rid9[key] = part9.relate_to(url, RT_HYPERLINK, is_external=True)

    c0_bee = [
        [("HEPLISAV-B", None)],
        [("(Phase 3 特定人群: HIV感染者)", None)],
        [("🏦 Dynavax", None)],
        [("🎯 预防乙型肝炎", None)],
    ]
    c1_bee = [
        [("✅【上市】", None)],
        [("(2017年FDA获批)", None)],
    ]
    c2_bee = [
        [
            ("▪ 试验: BEe-HIVe (", None),
            ("NCT04193189", rid9["nct_beehive"]),
            ("), Phase 3, 开放标签、随机, 10国41中心 (2020-2024)", None),
        ],
        [("▪ 适应人群: HIV感染者 (A组=既往乙肝疫苗无应答者; B组=乙肝初免者)", None)],
        [
            (
                "▪ 试验分组: A组 2-CpG(HEPLISAV-B 2剂) / 3-CpG(HEPLISAV-B 3剂) / 3-alum(Engerix-B 3剂); B组 HEPLISAV-B 3剂",
                None,
            )
        ],
        [("▪ 免疫程序: HEPLISAV-B 2剂(0,4周)或3剂(0,4,24周) vs Engerix-B 3剂(0,4,24周)", None)],
        [("▪ 样本量: 638例 (A组 187 / 188 / 186; B组 74)", None)],
        [
            (
                "▪ 研究终点: 主要=①血清保护应答(抗-HBs≥10mIU/mL; 2剂组第12周/3剂组第28周) ②AE发生率(入组至72周); 次要=各访视SPR、抗-HBs滴度分层、接种后4周内Grade≥2 AE",
                None,
            )
        ],
    ]
    c3_bee = [
        [("▪ 免疫原性: A组SPR 93.1% / 99.4% vs 80.6% (组别顺序 2-CpG/3-CpG/3-alum)", None)],
        [("▪ 任何AE(全研究期至72周, CT.gov 2025-07发布): 73.3% / 77.7% / 73.7% (B组 87.8%)", None)],
        [("▪ 接种后4周内 Grade≥2 AE: 33.7% / 45.7% / 43.5%", None)],
        [
            (
                "▪ 常见AE(按组别顺序): 注射部位疼痛 25.1% / 39.9% / 22.0%; 头痛 16.0% / 20.2% / 17.2%",
                None,
            )
        ],
        [
            (
                "▪ 疲劳 14.4% / 17.6% / 17.7%; 不适 11.8% / 18.1% / 14.0%; 肌痛 9.1% / 17.6% / 13.4%",
                None,
            )
        ],
        [
            (
                "▪ SAE: 6.4%(12例) / 8.0%(15例) / 5.9%(11例), 事件分散(HIV背景)且研究者判定均与疫苗无关; 死亡 0.5%(1例) / 0.5%(1例) / 0",
                None,
            )
        ],
    ]
    c4_bee = [
        [("NCT04193189", rid9["nct_beehive"]), (" (BEe-HIVe)", None)],
        [("PMID: 39616603", rid9["pmid_jama"]), (" (Marks 2025, JAMA)", None)],
    ]

    set_cell(t9.rows[1].cells[0], c0_bee, part9)
    set_cell(t9.rows[1].cells[1], c1_bee, part9)
    set_cell(t9.rows[1].cells[2], c2_bee, part9)
    set_cell(t9.rows[1].cells[3], c3_bee, part9)
    set_cell(t9.rows[1].cells[4], c4_bee, part9)
    set_row_height(t9, 1, 2.5)

    # ---- 复制行1 模板 -> 行2 (HBV-18) ----
    new_row = add_row_xml(t9, 2.3)

    c0_hbv = [
        [("HEPLISAV-B", None)],
        [("(Phase 3 特定人群: 透析患者)", None)],
        [("🏦 Dynavax", None)],
        [("🎯 预防乙型肝炎", None)],
    ]
    c1_hbv = [
        [("✅【上市】", None)],
        [("(2017年FDA获批)", None)],
    ]
    c2_hbv = [
        [
            ("▪ 试验: HBV-18 (", None),
            ("NCT01195246", rid9["nct_hbv18"]),
            ("), Phase 3, 随机、开放标签, 德国20中心 (2010-2012)", None),
        ],
        [("▪ 适应人群: 血液透析无应答成人", None)],
        [("▪ 试验分组: HEPLISAV 0.5mL / Engerix-B 2.0mL(加倍剂量) / Fendrix 0.5mL", None)],
        [("▪ 免疫程序: 单剂加强免疫", None)],
        [("▪ 样本量: 155例 (54 / 50 / 51)", None)],
        [
            (
                "▪ 研究终点: 主要=第4周SPR(抗-HBs≥10mIU/mL); 次要=各组注射后反应与AE总发生率(至第12周)",
                None,
            )
        ],
    ]
    c3_hbv = [
        [("▪ 免疫原性: 第4周SPR 52.8% vs 32.6% vs 43.1% (组别顺序 CpG/Engerix-B/Fendrix)", None)],
        [
            (
                "▪ 局部反应: 9.3% / 8.2% / 31.4% (疼痛为主 9.3% / 8.2% / 29.4%; Fendrix组4例中度余均轻度; 局部反应Fendrix显著高 p=.007)",
                None,
            )
        ],
        [
            (
                "▪ 全身反应: 18.5% / 12.2% / 19.6% (不适 9.3% / 4.2% / 2.0%; 头痛 9.3% / 2.1% / 2.0%; 疲劳 7.4% / 0 / 9.8%; 肌痛 7.4% / 4.1% / 7.8%); 每组各1例重度",
                None,
            )
        ],
        [("▪ 任意AE: 44.4% / 44.0% / 43.1%; SAE: 18.5% / 18.0% / 13.7% (均判与疫苗无关)", None)],
        [("▪ 相关AE: 3.7% / 4.0% / 3.9%; 相关SAE: 0; 因AE退出: 0", None)],
    ]
    c4_hbv = [
        [("NCT01195246", rid9["nct_hbv18"]), (" (HBV-18)", None)],
        [("PMID: 36269938", rid9["pmid_girndt"]), (" (Girndt 2022)", None)],
    ]

    set_cell(new_row.cells[0], c0_hbv, part9)
    set_cell(new_row.cells[1], c1_hbv, part9)
    set_cell(new_row.cells[2], c2_hbv, part9)
    set_cell(new_row.cells[3], c3_hbv, part9)
    set_cell(new_row.cells[4], c4_hbv, part9)

    # ---- Slide 9 清理孤儿 rel（原综合/特定人群行的 11 条链接）----
    removed9 = cleanup_orphan_hyperlink_rels(s9)

    # ================= 2. 修改 Slide 8 =================
    s8 = prs.slides[7]
    t8 = find_table(s8)
    if t8 is None:
        sys.exit("ERROR: Slide 8 table not found")
    part8 = s8.part

    # 行2（特定人群）重写 c2/c3；c0/c1/c4 保留
    rid8 = {
        "nct17": "rId8",  # NCT00985426 (HBV-17)  已存在
        "nct19": "rId9",  # NCT01282762 (HBV-19)  已存在
    }
    # 保险：确认 rId8/rId9 存在，否则 relate_to
    rels8 = s8.part.rels
    if rid8["nct17"] not in rels8:
        rid8["nct17"] = s8.part.relate_to(
            "https://clinicaltrials.gov/study/NCT00985426", RT_HYPERLINK, is_external=True
        )
    if rid8["nct19"] not in rels8:
        rid8["nct19"] = s8.part.relate_to(
            "https://clinicaltrials.gov/study/NCT01282762", RT_HYPERLINK, is_external=True
        )

    c2_spec = [
        [("▪ 分期: Phase 3/3b (HBV-17 CKD主研究 + HBV-19 34个月随访)", None)],
        [("▪ 设计: HBV-17 随机、观察者盲、多中心; HBV-19 开放标签、多中心", None)],
        [("▪ 适应人群: 慢性肾病(CKD)成人", None)],
        [("▪ 试验分组: HEPLISAV-B(CpG) 3剂 vs Engerix-B 加倍剂量4剂", None)],
        [
            (
                "▪ 免疫程序: HEPLISAV-B 0/4/24周(第8周安慰剂模拟) vs Engerix-B 双倍剂量(2×1.0mL) 0/4/8/24周",
                None,
            )
        ],
        [
            (
                "▪ 样本量: HBV-17 521例 (HEPLISAV-B 258 / Engerix-B 263); HBV-19 147例 (HepB-CpG 72 / HepB-Eng 75)",
                None,
            )
        ],
        [
            (
                "▪ 研究终点: HBV-17 主要=第28周SPR(≥10mIU/mL)非劣效/优效; HBV-19 主要=SPR持久性(6-48个月)",
                None,
            )
        ],
    ]
    c3_spec = [
        [
            ("▪ HBV-17 (", None),
            ("NCT00985426", rid8["nct17"]),
            (
                ") CT.gov结果: 第28周SPR 89.9%(204/227) vs 81.8%(198/242), 组间差8%(95%CI 1.7-14.3), 非劣效且优效",
                None,
            ),
        ],
        [
            (
                "▪ HBV-17 安全性: 因AE退出 0 vs 1例; SAE为CKD人群背景事件(肾衰/心衰/感染等), 两组无失衡; 研究期间死亡 7 vs 3例(均判与疫苗无关)",
                None,
            )
        ],
        [
            ("▪ HBV-19 (", None),
            ("NCT01282762", rid8["nct19"]),
            (
                ") Girndt 2023: 随访期末维持血清保护 77.4% vs 70.3% (P=0.4723); 维持抗-HBs≥100mIU/mL中位时间 37.6 vs 17.4个月 (P=0.0076)",
                None,
            ),
        ],
        [
            (
                "▪ HBV-19 安全性: 任意AE 16.7% vs 18.8%; 相关AE 4.2% vs 0; SAE 4.2% vs 3.1%(均判无关); 因AE中止 0; 死亡 2.7% vs 1.4%(判无关)",
                None,
            )
        ],
        [("▪ 结论: 长期耐受性良好, 无新增安全信号; CpG组抗体水平更高、保护更持久", None)],
    ]

    set_cell(t8.rows[2].cells[2], c2_spec, part8)
    set_cell(t8.rows[2].cells[3], c3_spec, part8)
    set_row_height(t8, 2, 2.4)

    # 删除行3（原 BEe-HIVe+HBV-18 整合行）
    delete_row(t8, 3)

    # ---- Slide 8 清理孤儿 rel（原行3 的 4 条链接）----
    removed8 = cleanup_orphan_hyperlink_rels(s8)

    prs.save(DST)
    print("SAVED:", DST)
    print("Slide 9 link rIds:", {k: v for k, v in rid9.items()})
    print("Slide 9 removed orphan rels:", removed9)
    print("Slide 8 removed orphan rels:", removed8)
    print("Slide 8 rId8/rId9 used:", rid8)


if __name__ == "__main__":
    main()
