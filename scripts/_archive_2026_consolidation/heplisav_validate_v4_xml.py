"""XML 完整性验证：V4 文件所有部件可解析、关系引用有效"""

import sys
import zipfile

sys.stdout.reconfigure(encoding="utf-8")
from lxml import etree

V4 = r"E:/Cursor Project/2-Scientific-Skills-for-Clinical_Trial/review_materials/60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v4.pptx"
NS = {"r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}

z = zipfile.ZipFile(V4)
names = z.namelist()
print("zip 部件数:", len(names))

# 1) 所有 xml 可解析
bad = 0
for n in names:
    if n.endswith(".xml") or n.endswith(".rels"):
        try:
            etree.fromstring(z.read(n))
        except Exception as e:
            bad += 1
            print("  解析失败:", n, e)
print("XML 解析失败数:", bad)

# 2) 检查每个 slide 的 rels 引用是否存在于包中
rels_bad = 0
for n in names:
    if n.endswith(".rels"):
        root = etree.fromstring(z.read(n))
        rels_dir = n.replace("_rels/", "").rsplit("/", 1)[0] + "/" if "_rels/" in n else ""
        for rel in root:
            target = rel.get("Target", "")
            if (
                target.startswith("http")
                or target.startswith("mailto")
                or target.startswith("https")
            ):
                continue
            # 解析相对路径
            full = rels_dir + target
            # 规范化
            parts = full.split("/")
            stack = []
            for p in parts:
                if p == "..":
                    if stack:
                        stack.pop()
                elif p in ("", "."):
                    continue
                else:
                    stack.append(p)
            full = "/".join(stack)
            if full not in names:
                # 有些 target 以 / 开头（绝对包路径）
                alt = target.lstrip("/")
                if alt not in names:
                    rels_bad += 1
                    print("  悬空引用:", n, "->", target)
print("悬空关系引用数:", rels_bad)

# 3) 检查 slide7/9 的 sldId 顺序与文本
from pptx import Presentation

prs = Presentation(V4)
print("\n幻灯片总数:", len(prs.slides))
for i in (6, 8):
    slide = prs.slides[i]
    texts = []
    for shp in slide.shapes:
        if shp.has_table:
            texts.append("TBL(%dx%d)" % (len(shp.table.rows), len(shp.table.columns)))
        elif shp.has_text_frame and shp.text_frame.text.strip():
            texts.append(shp.text_frame.text.replace("\n", "")[:30])
    print("Slide %d: %s" % (i + 1, " | ".join(texts)))
print("\n文件大小: %.1f KB" % (len(open(V4, "rb").read()) / 1024))
print("验证完成")
