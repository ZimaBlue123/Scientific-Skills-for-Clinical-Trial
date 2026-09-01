"""
count_pptx_hyperlinks.py
========================
只读统计 .pptx 文件中所有超链接（hlinkClick）：
- 按段落统计（同一段落内拆成多个 run 的链接计 1 个）
- 解析 rId → 实际 URL（依赖 slide part 的 rels）
- 输出：总链接数、去重 URL 数、按页分布、按域名分类
用法：python count_pptx_hyperlinks.py <file.pptx>
"""

import sys
from collections import Counter, OrderedDict

from pptx import Presentation
from pptx.oxml.ns import qn


def iter_link_paragraphs(shape, part):
    """遍历 shape 内所有带 hlinkClick 的段落，yield (文本, url)。"""
    # 表格
    if shape.has_table:
        tbl = shape.table
        for ri in range(len(tbl.rows)):
            for ci in range(len(tbl.columns)):
                cell = tbl.cell(ri, ci)
                for para in cell.text_frame.paragraphs:
                    url = _para_url(para, part)
                    if url:
                        yield (para.text, url, "表格[%d,%d]" % (ri, ci))
        return
    # 文本框 / 图形
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            url = _para_url(para, part)
            if url:
                yield (para.text, url, "文本框")


def _para_url(para, part):
    """段落中第一个带 hlinkClick 的 run → 解析 URL。"""
    for run in para.runs:
        rPr = run._r.find(qn("a:rPr"))
        if rPr is None:
            continue
        hlink = rPr.find(qn("a:hlinkClick"))
        if hlink is None:
            continue
        rId = hlink.get(qn("r:id"))
        if not rId:
            continue
        try:
            rel = part.rels[rId]
            return rel.target_ref
        except KeyError:
            return "?rId=%s" % rId
    return None


def main(path):
    prs = Presentation(path)
    total = 0
    unique_urls = Counter()
    by_slide = OrderedDict()
    domain_count = Counter()
    details = []

    for idx, slide in enumerate(prs.slides):
        slide_links = []
        for shp in slide.shapes:
            part = shp.part
            for text, url, where in iter_link_paragraphs(shp, part):
                total += 1
                unique_urls[url] += 1
                domain = url.split("/")[2] if url.startswith(("http://", "https://")) else url
                domain_count[domain] += 1
                slide_links.append((url, text.replace("\n", " ")[:40], where))
                details.append((idx + 1, url, text.replace("\n", " ")[:40], where))
        if slide_links:
            by_slide[idx + 1] = slide_links

    print("=" * 80)
    print("文件:", path)
    print("总幻灯片数:", len(prs.slides))
    print("链接段落总数:", total)
    print("去重 URL 数:", len(unique_urls))
    print()
    print("----- 按页分布 -----")
    for sidx, links in by_slide.items():
        urls = [u for u, _, _ in links]
        print("Slide %d: %d 个链接" % (sidx, len(links)))
        for u in set(urls):
            print("    %s (x%d)" % (u, urls.count(u)))
    print()
    print("----- 按域名分类 -----")
    for d, c in domain_count.most_common():
        print("  %-50s %d" % (d, c))
    print()
    print("----- 链接明细 -----")
    for sidx, url, text, where in details:
        print("S%02d [%s] %s -> %s" % (sidx, where, text, url))


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python count_pptx_hyperlinks.py <file.pptx>")
        sys.exit(1)
    main(sys.argv[1])
