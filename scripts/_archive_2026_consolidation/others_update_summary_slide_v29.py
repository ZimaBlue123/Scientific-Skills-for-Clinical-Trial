"""
update_summary_slide_v29.py
在 V29 PPT（CpG_Vaccine_Safety_Summary-V29-20260827.pptx）最后一页
【总体安全性总结】上，基于 V29 详表最新数据重写正文：
  - 修正数据口径（Z-1018 中重度反应范围、HEPLISAV-B 自身免疫事件 0.2%/0.1% vs 对照 0%）
  - 补充反应原性谱（CORBEVAX、SCB-2019）、0 疫苗相关 SAE 证据链、特殊人群板块
  - 正文 17pt；板块标题与"结论"加粗；其余继承原样式
用法: python update_summary_slide_v29.py
"""

import re

from pptx import Presentation
from pptx.util import Pt

PPT = r"e:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\CpG_Vaccine_Safety_Summary-V29-20260827.pptx"

# ---- 终页总结（精简克制版，数据与 V29 详表逐条核对） ----
SUMMARY_LINES = [
    "1. 局部与全身反应原性 (Reactogenicity)",
    "CpG 预防性疫苗最常见的 ADR 为轻至中度注射部位疼痛及疲劳、头痛等全身反应，"
    "多为 1-2 级、一过性、自限性。已上市产品中注射部位疼痛发生率因适应症人群及"
    "抗原平台不同差异较大（约 15%-86%：IndoVac 14.69%、HEPLISAV-B 23%-39%、"
    "MVC-COV1901 71.2%、CYFENDUS 86.3%），未观察到随剂次累积加重。Z-1018 与 "
    "Shingrix 对照研究中，中重度局部反应（7.7%-35.0% vs 52.6%）及全身反应"
    "（17.5%-46.2% vs 63.2%）发生率较低。",
    "",
    "2. 严重不良事件与特殊关注事件 (SAE/AESI)",
    "已上市产品及在研管线的 III 期/上市后数据中，多数 SAE 及 AESI 经评估与疫苗"
    "无关，未发现与疫苗有明确因果关系的 SADR 信号或致死性事件（SCB-2019 疫苗"
    "相关 SAE 0.027%、HEPLISAV-B 新发自身免疫事件 0.1%-0.2% 且对照组为 0、"
    "CYFENDUS SAE 1.9% 均判无关）。个别 AESI 被评估为可能与疫苗相关（CYFENDUS "
    "自身免疫病因事件 3 例、MVC-COV1901 面神经麻痹 1 例）。",
    "",
    "3. 特征性生理反应",
    "部分管线观察到一过性实验室指标变化：Na-GST-1 中 CpG 佐剂抑制钩虫感染所致"
    "外周血嗜酸性粒细胞增多（中位 0.6 vs 3.1×10³/μL，p=0.027）；HEPLISAV-B "
    "早期高剂量组见一过性 ALT/AST 轻度升高，补体 C3/C4 无改变，无 ANA/抗 DNA "
    "转阳。",
    "",
    "结论",
    "现有临床及上市后数据支持：明确因果关系的 ADR 主要为一过性注射部位反应及"
    "轻中度全身反应。CpG 佐剂（CpG 1018、CpG 7909、CpG-QCX1 等）在增强免疫"
    "原性的同时，整体安全性特征成熟，未观察到严重相关毒性风险增加。",
]

HEADER_RE = re.compile(r"^\d+\.\s|^结论$")


def main():
    prs = Presentation(PPT)
    slide = prs.slides[-1]
    target = None
    for s in slide.shapes:
        if s.has_text_frame and s.text and "1. " in s.text:
            target = s
            break
    if target is None:
        raise SystemExit("未找到总结页正文文本框，中止。")

    tf = target.text_frame
    # 清空全部现有段落
    for p_elem in list(tf.paragraphs):
        p_elem._p.getparent().remove(p_elem._p)
    # 重建段落（全部重新 add，避免首段不存在的问题）
    for line in SUMMARY_LINES:
        p = tf.add_paragraph()
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = line
        r.font.size = Pt(17)
        if HEADER_RE.match(line):
            r.font.bold = True
    prs.save(PPT)
    print(f"OK: 总结页已更新 ({len(SUMMARY_LINES)} 段)，保存至 {PPT}")


if __name__ == "__main__":
    main()
