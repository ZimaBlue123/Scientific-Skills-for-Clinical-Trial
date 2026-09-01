"""
V7 生成脚本：基于 v6，在 Slide 8（HEPLISAV-B Phase 3 页）表格末尾新增行 3
内容：两个特定人群试验的安全性数据汇总
  - 试验1: BEe-HIVe (NCT04193189)  HIV感染者, JAMA 2025 (PMID 39616603) + CT.gov 结果
  - 试验2: HBV-18 (NCT01195246)    透析无应答者, Girndt 2022 (PMID 36269938)

格式对齐 row 2：Arial/微软雅黑 10pt(sz=1000)，tcPr 白底 + 边距 50800/19050，
段落 lnSpc 100% + spcBef/spcAft 1pt，链接以 a:hlinkClick 嵌入 run rPr。
"""

import copy
import sys
from xml.sax.saxutils import escape

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches

SRC = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v6.pptx"
DST = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v7.pptx"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
RT_HYPERLINK = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"

RPR_TMPL = (
    '<a:rPr xmlns:a="{a}" xmlns:r="{r}" sz="1000" dirty="0">'
    '<a:latin typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
    '<a:ea typeface="微软雅黑" panose="020B0503020204020204" pitchFamily="34" charset="-122"/>'
    '<a:sym typeface="Arial" panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
    "{hlink}"
    "</a:rPr>"
)


def make_rpr(rid=None):
    hlink = f'<a:hlinkClick r:id="{rid}"/>' if rid else ""
    return RPR_TMPL.format(a=A_NS, r=R_NS, hlink=hlink)


def make_p(segments):
    """segments: list of (text, rid_or_None)"""
    runs = ""
    for text, rid in segments:
        runs += f'<a:r>{make_rpr(rid)}<a:t xml:space="preserve">{escape(text)}</a:t></a:r>'
    xml = (
        f'<a:p xmlns:a="{A_NS}" xmlns:r="{R_NS}">'
        '<a:pPr><a:lnSpc><a:spcPct val="100000"/></a:lnSpc>'
        '<a:spcBef><a:spcPts val="100"/></a:spcBef>'
        '<a:spcAft><a:spcPts val="100"/></a:spcAft></a:pPr>'
        f"{runs}</a:p>"
    )
    return etree.fromstring(xml)


def set_cell(cell, paras, part):
    """paras: list of segments-lists; 清空 cell 并写入段落"""
    tc = cell._tc
    # 清空已有 txBody 内容（保留 bodyPr/lstStyle）
    txBody = tc.find(qn("a:txBody"))
    for child in list(txBody):
        if child.tag not in (qn("a:bodyPr"), qn("a:lstStyle")):
            txBody.remove(child)
    for segs in paras:
        txBody.append(make_p(segs))
    # 设置单元格属性（白底 + 边距，对齐 row 2）
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))
        tc.insert(0, tcPr)
    tcPr.set("marL", "50800")
    tcPr.set("marR", "50800")
    tcPr.set("marT", "19050")
    tcPr.set("marB", "19050")
    # 移除旧填充后加白底
    for sf in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(sf)
    fill = etree.SubElement(tcPr, qn("a:solidFill"))
    clr = etree.SubElement(fill, qn("a:srgbClr"))
    clr.set("val", "FFFFFF")


def add_row_xml(table, height_inches):
    """XML 层面追加一行：深拷贝最后一行模板（继承 tcPr 格式），清空文本"""
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    if not trs:
        sys.exit("ERROR: no <a:tr> in table")
    new_tr = copy.deepcopy(trs[-1])
    new_tr.set("h", str(Emu(Inches(height_inches))))
    for tc in new_tr.findall(qn("a:tc")):
        txBody = tc.find(qn("a:txBody"))
        if txBody is not None:
            for child in list(txBody):
                if child.tag not in (qn("a:bodyPr"), qn("a:lstStyle")):
                    txBody.remove(child)
    tbl.append(new_tr)
    return table.rows[len(table.rows) - 1]


def main():
    prs = Presentation(SRC)
    slide = prs.slides[7]  # Slide 8
    table = None
    for sh in slide.shapes:
        if sh.has_table:
            table = sh.table
            break
    if table is None:
        sys.exit("ERROR: Slide 8 no table found")

    # ---- 新增行 ----
    new_row = add_row_xml(table, 2.2)

    # ---- 建立超链接关系，返回 rId ----
    part = slide.part
    rid = {}
    urls = {
        "nct_beehive": "https://clinicaltrials.gov/study/NCT04193189",
        "nct_hbv18": "https://clinicaltrials.gov/study/NCT01195246",
        "pmid_jama": "https://pubmed.ncbi.nlm.nih.gov/39616603/",
        "pmid_girndt": "https://pubmed.ncbi.nlm.nih.gov/36269938/",
    }
    for key, url in urls.items():
        rid[key] = part.relate_to(url, RT_HYPERLINK, is_external=True)

    # ---- 5 列内容 ----
    c0 = [
        [("HEPLISAV-B", None)],
        [("(Phase 3 特定人群: HIV感染者、透析患者)", None)],
        [("🏦 Dynavax", None)],
        [("🎯 预防乙型肝炎", None)],
    ]
    c1 = [
        [("✅【上市】", None)],
        [("(2017年FDA获批)", None)],
    ]
    c2 = [
        [("▪ 试验1 BEe-HIVe (NCT04193189): Phase 3, HIV感染者", None)],
        [("  - 设计: 开放标签、随机, 10国41中心", None)],
        [
            (
                "  - 样本量: 638例 [A组无应答者: 2-CpG 187 / 3-CpG 188 / 3-alum(Engerix-B) 186; B组初免者 74]",
                None,
            )
        ],
        [("  - 免疫程序: HEPLISAV-B 2剂(0,4周)或3剂(0,4,24周) vs Engerix-B 3剂(0,4,24周)", None)],
        [("  - 血清学转换率(SPR): 93.1% / 99.4% vs 80.6%", None)],
        [("▪ 试验2 HBV-18 (NCT01195246): Phase 3, 透析无应答者", None)],
        [("  - 设计: 随机、开放标签, 德国20中心", None)],
        [("  - 样本量: 155例 (CpG 54 / Engerix-B 50 / Fendrix 51)", None)],
        [("  - 免疫程序: 单剂加强免疫", None)],
        [("  - 血清学转换率(SPR): 52.8% vs 32.6% vs 43.1%", None)],
    ]
    c3 = [
        [
            ("▪ 试验1 BEe-HIVe (", None),
            ("NCT04193189", rid["nct_beehive"]),
            (")｜HIV感染者｜全研究期至72周 (CT.gov 2025-07发布)", None),
        ],
        [
            (
                "  分组(组别顺序: 2-CpG / 3-CpG / 3-alum): A组无应答者 187/188/186例; B组初免者 74例",
                None,
            )
        ],
        [
            (
                "  任何AE: 73.3% / 77.7% / 73.7% (B组 87.8%); 末次接种后4周内 Grade≥2 AE: 33.7% / 45.7% / 43.5%",
                None,
            )
        ],
        [
            (
                "  常见AE(按组别顺序): 注射部位疼痛 25.1% / 39.9% / 22.0%; 头痛 16.0% / 20.2% / 17.2%",
                None,
            )
        ],
        [
            (
                "  疲劳 14.4% / 17.6% / 17.7%; 不适 11.8% / 18.1% / 14.0%; 肌痛 9.1% / 17.6% / 13.4%",
                None,
            )
        ],
        [
            (
                "  SAE: 6.4%(12例) / 8.0%(15例) / 5.9%(11例), 事件分散且研究者判定均与疫苗无关; 死亡 0.5%(1例) / 0.5%(1例) / 0",
                None,
            )
        ],
        [
            ("▪ 试验2 HBV-18 (", None),
            ("NCT01195246", rid["nct_hbv18"]),
            (")｜透析无应答者 (Girndt 2022)", None),
        ],
        [("  分组(组别顺序: CpG / Engerix-B / Fendrix): 54 / 50 / 51例 (单剂加强)", None)],
        [
            (
                "  局部反应: 9.3% / 8.2% / 31.4% (疼痛为主: 9.3% / 8.2% / 29.4%; Fendrix组4例中度, 余均轻度)",
                None,
            )
        ],
        [
            (
                "  全身反应: 18.5% / 12.2% / 19.6% (不适 9.3% / 4.2% / 2.0%; 头痛 9.3% / 2.1% / 2.0%; 疲劳 7.4% / 0 / 9.8%; 肌痛 7.4% / 4.1% / 7.8%)",
                None,
            )
        ],
        [
            (
                "  任意AE: 44.4% / 44.0% / 43.1%; SAE: 18.5% / 18.0% / 13.7% (均与疫苗无关); 相关AE: 3.7% / 4.0% / 3.9%; 相关SAE: 0; 因AE退出: 0",
                None,
            )
        ],
    ]
    c4 = [
        [("NCT04193189", rid["nct_beehive"]), (" (BEe-HIVe)", None)],
        [("NCT01195246", rid["nct_hbv18"]), (" (HBV-18)", None)],
        [("PMID: 39616603", rid["pmid_jama"]), (" (Marks 2025, JAMA)", None)],
        [("PMID: 36269938", rid["pmid_girndt"]), (" (Girndt 2022)", None)],
    ]

    set_cell(new_row.cells[0], c0, part)
    set_cell(new_row.cells[1], c1, part)
    set_cell(new_row.cells[2], c2, part)
    set_cell(new_row.cells[3], c3, part)
    set_cell(new_row.cells[4], c4, part)

    prs.save(DST)
    print("SAVED:", DST)
    print("links:")
    for k, v in rid.items():
        print(f"  {k}: rId={v} -> {urls[k]}")


if __name__ == "__main__":
    main()
