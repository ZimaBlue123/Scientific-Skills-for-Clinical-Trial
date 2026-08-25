"""
PPT Template Injector

A utility to separate presentation design from data logic.
Instead of hardcoding colors, fonts, and positions in Python, you create a
.pptx template in PowerPoint with placeholders like {{title}} and {{summary}}.
This module will inject the data into those placeholders while preserving
the original styling.
"""

from __future__ import annotations

import logging
from typing import Any

from pptx import Presentation

LOGGER = logging.getLogger(__name__)


def _replace_text_in_paragraph(paragraph: Any, replacements: dict[str, str]) -> None:
    """
    Replaces text in a paragraph while attempting to preserve formatting.
    PowerPoint splits text into multiple 'runs' arbitrarily.
    For a robust replacement, we combine run text, replace it, and put it
    back in the first run, preserving the original run's font/size/color.
    """
    p_text = paragraph.text
    if not any(key in p_text for key in replacements):
        return

    # Store first run's formatting to re-apply it after replacement
    if not paragraph.runs:
        return

    first_run = paragraph.runs[0]
    font_name = first_run.font.name
    font_size = first_run.font.size
    font_bold = first_run.font.bold
    font_color = None
    if first_run.font.color and hasattr(first_run.font.color, "rgb"):
        font_color = first_run.font.color.rgb

    # Do the replacement on the full paragraph text
    new_text = p_text
    for key, value in replacements.items():
        new_text = new_text.replace(key, str(value))

    # Clear old runs
    for run in paragraph.runs:
        run.text = ""

    # Assign the new text to the first run
    first_run.text = new_text
    if font_name:
        first_run.font.name = font_name
    if font_size:
        first_run.font.size = font_size
    if font_bold is not None:
        first_run.font.bold = font_bold
    if font_color:
        first_run.font.color.rgb = font_color


def inject_data_into_template(
    template_path: str, output_path: str, slide_data: list[dict[str, str]]
) -> None:
    """
    Takes a path to a template PPTX.
    slide_data is a list of dictionaries. Each dictionary represents a slide's replacements.
    E.g.: [{'{{title}}': 'Slide 1 Title'}, {'{{title}}': 'Slide 2 Title'}]
    """
    prs = Presentation(template_path)

    for i, slide_dict in enumerate(slide_data):
        if i < len(prs.slides):
            slide = prs.slides[i]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        _replace_text_in_paragraph(paragraph, slide_dict)
                # Also check tables
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                _replace_text_in_paragraph(paragraph, slide_dict)
        else:
            LOGGER.warning(f"Not enough slides in template. Skipping data index {i}")

    prs.save(output_path)
    LOGGER.info(f"Template injection complete. Saved to {output_path}")
