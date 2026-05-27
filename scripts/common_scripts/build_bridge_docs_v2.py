from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Pt as PPTPt
from pathlib import Path

out_dir = Path('deliverables')
out_dir.mkdir(exist_ok=True)

# =========================
# Word (full templates)
# =========================
word_path = out_dir / '海外桥接-完整模板包.docx'
doc = Document()
doc.styles['Normal'].font.name = 'Calibri'
doc.styles['Normal'].font.size = Pt(11)

doc.add_heading('海外桥接完整模板包（可直接复制）', 0)
doc.add_paragraph('版本：V2.0\n适用项目：重组带状疱疹疫苗（CHO细胞）\n说明：本模板按“海外桥接总方案大纲 + 四国监管问题清单”完整展开，结合国内I/II期、澳洲I期CSR证据与ICH/WHO/EMA通用桥接原则。')

doc.add_page_break()
doc.add_heading('模板1：海外桥接总方案大纲（Synopsis/Bridge Study）', 1)

blocks = [
('1. 研究标题', [
'一项在[国家/多国家]开展的、随机、盲法、[阳性对照/安慰剂+阳性对照]的临床桥接研究，用于评估重组带状疱疹疫苗（CHO细胞）在[40岁及以上]人群中的免疫原性与安全性，并支持既有临床证据的区域外推。'
]),
('2. 背景与立题依据', [
'产品：重组带状疱疹疫苗（CHO细胞），抗原为gE，佐剂系统为TVA01。',
'已有证据：国内I期、国内II期阶段性、澳洲I期CSR。',
'监管原则：ICH E5(R1)、ICH E17、WHO TRS 1004 Annex 9、EMA疫苗临床评价指南。',
'桥接必要性：目标国对本地人群可外推性、本地免疫应答一致性与风险可控性提出监管要求。'
]),
('3. 研究目的', [
'主要目的：验证本品在目标区域人群中的关键免疫原性指标与既有证据链的一致性/非劣性，支撑区域外推。',
'次要目的：评估短期及中期安全性，描述年龄分层（40-49、50-59、≥60、可选≥70）免疫应答差异，评估体液-细胞免疫一致性。'
]),
('4. 研究设计', [
'类型：[单国桥接/多区域桥接]。',
'设计：随机、盲法、[阳性对照/安慰剂+阳性对照]。',
'人群：40岁及以上健康成年人；按年龄层分层随机。',
'样本量：总计[ ]例；每年龄层[ ]例；试验组:对照组=[ ]。',
'程序：D0、D60两剂；肌内注射。',
'细胞免疫亚组：每年龄层前[ ]例（可选）。'
]),
('5. 终点设置', [
'主要终点（建议固定在D90）：第2剂后30天抗gE抗体SCR、GMC（或GMT）；抗VZV抗体SCR、GMC（或GMT）。',
'次要终点：GMI（gE/VZV）；细胞免疫（gE特异CD4+T细胞应答率/频率）；征集AE（0-14天）、非征集AE（0-30天）、SAE/AESI（全程）；第2剂后12/24/36月持久性。'
]),
('6. 统计分析框架', [
'分析集：ITT/FAS、SS、PPS（体液/细胞分别定义）。',
'主要比较：分层（年龄）+总体模型；给出组间差值/比值及95%CI。',
'多重性控制：[层级检验/校正方法]。',
'敏感性分析：ITT与PPS一致性、年龄亚组（≥60/≥70）、方法学一致性（ELISA/FAMA）。'
]),
('7. 外推与桥接论证路径', [
'论证四要素：人群可比性、检测方法可比性、工艺/批次可比性、终点与统计可解释性。',
'针对≥60岁潜在差异，预设协变量校正、分层结果一致性解释与风险-获益综合判断路径。'
]),
('8. 安全风险管理', [
'重点风险：局部反应、发热、3级反应、潜在AESI。',
'DSMB触发规则：[阈值]；暂停/终止标准：[阈值]。',
'SAE/AESI快速报告机制：[24h/72h]。'
]),
('9. 关键操作与质量控制', [
'中央实验室与方法学标准化（ELISA/FAMA）。',
'关键时间窗采血依从性管理。',
'盲态维护与揭盲规则。',
'数据核查、稽查与审计追踪。'
]),
('10. 里程碑计划', [
'FPFV：[ ]',
'LPLV（D90主要分析）：[ ]',
'阶段性CSR：[ ]',
'最终CSR：[ ]'
]),
('11. 附件清单', [
'国内I/II期CSR摘要与对照映射表',
'澳洲I期CSR摘要',
'检测方法学验证与一致性文件',
'SAP摘要',
'CMC桥接说明（临床批-商业批）',
'RMP/PV摘要文件'
])
]

for title, items in blocks:
    doc.add_heading(title, 2)
    for it in items:
        doc.add_paragraph(it, style='List Bullet')

doc.add_page_break()
doc.add_heading('模板2：四国监管问题清单（Russia / Egypt / Brazil / Indonesia）', 1)

doc.add_heading('A. 通用问题（四国共用）', 2)
common = [
'1) 请确认接受以既有境外/境内临床数据为基础的免疫桥接路径，并说明可接受的主要免疫终点（SCR/GMC/GMT/GMI）。',
'2) 请确认桥接研究统计目标应采用非劣、等效或优效，以及可接受的统计界值/置信区间规则。',
'3) 对年龄分层（40-49、50-59、≥60/≥70）是否有强制要求？',
'4) 是否必须设置本地阳性对照（如Shingrix或本国已上市同类）？',
'5) 对免疫检测方法（ELISA/FAMA/中和）是否有指定要求？',
'6) 对CMC桥接（临床批与拟上市商业批一致性）需要哪些补充资料？',
'7) 阶段性提交是否可接受（D90先报、长期随访后补）？'
]
for q in common:
    doc.add_paragraph(q, style='List Bullet')

doc.add_heading('B. 俄罗斯（含EAEU路径）特异问题', 2)
for q in [
'1) 该项目应优先走EAEU统一路径还是俄罗斯本国路径？是否允许并行技术沟通？',
'2) 是否接受基于境外数据+本地桥接队列的注册策略？本地样本量最低建议是多少？',
'3) 对阳性对照疫苗来源、可及性和标签适应症有何要求？',
'4) 对高龄（≥60/≥70）亚组结果若出现“部分指标不占优”时，监管可接受的解释边界是什么？',
'5) 是否接受先上市后补充长期免疫持久性数据？'
]:
    doc.add_paragraph(q, style='List Bullet')

doc.add_heading('C. 埃及（EDA）特异问题', 2)
for q in [
'1) EDA对生物制品采用外部参考/依赖路径时，本项目桥接研究最小本地数据要求是什么？',
'2) 是否接受“免疫原性主要终点 + 安全性支持”的注册证据组合？',
'3) 对样本民族构成与外推论证（ethnic factors）有无格式化要求？',
'4) 对本地PV系统和RMP提交有何最低要件？'
]:
    doc.add_paragraph(q, style='List Bullet')

doc.add_heading('D. 巴西（ANVISA）特异问题', 2)
for q in [
'1) 在RDC 945/2024框架下，本项目桥接临床归入DDCM/DEEC哪一路径更优？',
'2) 对阳性对照与本地受试者比例有无明确门槛？',
'3) 是否接受阶段性分析（主要终点）先行审阅，长期终点后补？',
'4) 对免疫终点统计判据及多重性控制是否有官方偏好？'
]:
    doc.add_paragraph(q, style='List Bullet')

doc.add_heading('E. 印尼（BPOM）特异问题', 2)
for q in [
'1) 在ACTD提交中，桥接证据对Module 5的最低要求是什么？',
'2) 是否接受“参考国批准 + 本地免疫桥接”组合路径？',
'3) 对本地样本量、年龄分层、阳性对照有无明确监管期望？',
'4) 对上市后有效性/安全性真实世界补充数据的要求是什么？'
]:
    doc.add_paragraph(q, style='List Bullet')

doc.add_heading('F. 申办方预置答复材料清单（建议随Q&A同步）', 2)
for q in [
'国内I期CSR摘要（安全+免疫）',
'国内II期阶段性CSR摘要（分年龄层）',
'澳洲I期CSR摘要（英文）',
'免疫检测方法学与跨实验室一致性文件',
'SAP摘要（主要/敏感性分析）',
'CMC桥接说明（临床批-商业批）',
'风险管理计划（RMP）与PV主文件摘要'
]:
    doc.add_paragraph(q, style='List Bullet')

doc.add_heading('附录：≥60岁亚组风险标准话术（Discussion/回函可用）', 1)
doc.add_paragraph('在≥60岁亚组观察到部分免疫指标与阳性对照存在差异。申办方已预设分层统计、协变量校正及敏感性分析，并通过SCR、GMC/GMI、细胞免疫与安全性进行综合获益-风险评估；同时将补充长期持久性及上市后真实世界证据，以降低外推不确定性。')

doc.save(word_path)

# =========================
# PPT (detailed)
# =========================
ppt_path = out_dir / '海外桥接-细化版汇报.pptx'
prs = Presentation()

def add_title(title, subtitle=''):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle

def add_bullets(title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    tf = s.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = PPTPt(18)

add_title('海外桥接策略（细化版）', '重组带状疱疹疫苗（CHO细胞）\n依据：国内I/II期 + 澳洲I期 + ICH/WHO/EMA桥接原则')
add_bullets('1. 项目目标与输出', [
'目标：形成可被俄/埃及/巴西/印尼接受的桥接注册路径',
'输出：1份全球核心桥接方案 + 4份国家附录问题清单',
'原则：最大化复用现有CSR，最小化新增试验负担'
])
add_bullets('2. 现有证据底盘', [
'国内I期：随机盲法安慰剂对照，100例，安全/免疫信号明确',
'国内II期：随机盲法阳性对照，420例，D90分层结果可直接用于桥接论证',
'澳洲I期CSR：英文审评材料基础，结构与分析口径可迁移'
])
add_bullets('3. 国内II期分层信号（桥接关注点）', [
'40-49岁：部分指标相对阳性对照更优',
'50-59岁：多指标与阳性对照相当',
'≥60岁：存在“部分指标相当/部分GMC偏低”情形，需要预置解释框架'
])
add_bullets('4. 桥接研究主设计', [
'设计：随机、盲法、阳性对照（必要时加安慰剂/历史支持）',
'人群：40-49/50-59/≥60（可选≥70）分层随机',
'给药：D0/D60两剂；细胞免疫亚组前置抽样'
])
add_bullets('5. 终点与统计框架', [
'主要终点（D90）：gE与VZV SCR + GMC/GMT',
'次要终点：GMI、CD4+T应答率、AE/SAE/AESI、12/24/36月持久性',
'统计：FAS/SS/PPS并行，分层+总体模型，95%CI与敏感性分析'
])
add_bullets('6. 关键风险与应对', [
'风险1：≥60岁亚组外推不确定性',
'应对：预设分层分析、协变量校正、一致性敏感性分析',
'风险2：方法学差异',
'应对：ELISA/FAMA跨实验室一致性与质控链路'
])
add_bullets('7. 法规依据（桥接逻辑）', [
'ICH E5(R1)：民族因素与境外数据可接受性',
'ICH E17：多区域试验规划与一致性评价',
'WHO TRS 1004 Annex 9：bridging to efficacy框架',
'EMA疫苗临床评价指南：bridging studies原则'
])
add_bullets('8. 四国预沟通问题（摘要）', [
'俄罗斯：EAEU路径优先级、本地样本量、≥60解释边界',
'埃及：依赖路径+本地桥接最小数据包',
'巴西：RDC 945下DDCM/DEEC路径与阶段性提交',
'印尼：ACTD Module 5桥接最低要求'
])
add_bullets('9. 递交资料包结构', [
'核心包：Synopsis + SAP + 核心CSR证据摘要 + 方法学文件',
'国家附录：法规映射、对照可及性说明、本地执行与PV方案',
'补充包：CMC桥接、批次一致性、RMP'
])
add_bullets('10. 里程碑与行动', [
'第1周：锁定终点与统计界值（含≥60策略）',
'第2周：完成四国Q-list并发起预沟通',
'第3周：形成“核心包+国家附录”v1',
'第4周：根据反馈确定是否启动本地桥接队列'
])

prs.save(ppt_path)
print(str(word_path))
print(str(ppt_path))
