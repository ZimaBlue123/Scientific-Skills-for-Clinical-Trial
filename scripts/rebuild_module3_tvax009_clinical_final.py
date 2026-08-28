"""
重建《60岁以上乙肝流调-CpG佐剂安全性-接种程序循证-20260827-updated.pptx》第3模块（V3）。

V3 相对 V2 的变更（用户 2026-08-27 晚要求）：
1. 模块标题改为「03 新佐剂乙肝疫苗与减剂次的临床意义」
2. P16/P17 增加上标文献序号 + 页内 AMA 格式参考文献列表
3. Ⅲ期表述修正：18-59岁试验疫苗已确定2剂程序，Ⅲ期确认第2剂接种时间点（0,1月 vs 0,2月，
   理想状态第2剂于首剂后1-2个月内接种）；不再称"程序探索"
4. 统计学表述均引用原数据并给出具体 P 值（取自 TVAX-009 沟通会 PPT 表格原文）
5. 全文删除"左右"，改用"约"或精确数值
6. 正文纯黑色（参考文献列表可用灰色）；正文字号提升至 11pt 级
7. 输出独立 V3 文件，不覆盖 V2

用法：
    python rebuild_module3_tvax009_clinical_v3.py
输出：
    review_materials/60岁以上乙肝流调-CpG佐剂安全性-接种程序循证-20260827-updated-v3.pptx
"""

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------- 常量 ----------
SRC = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-接种程序循证-20260827-updated.pptx"
DST = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\60岁以上乙肝流调-CpG佐剂安全性-接种程序循证-20260827-updated-v3.pptx"

RED = RGBColor(0xC0, 0x00, 0x00)  # 主强调色 RGB(192,0,0)
BLACK = RGBColor(0x00, 0x00, 0x00)  # 正文纯黑
GRAY = RGBColor(0x80, 0x80, 0x80)  # 仅参考文献列表
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)  # 表格浅灰行
FONT = "微软雅黑"

SLIDE_W = 13.333
SLIDE_H = 7.5
DELETE_FROM = 13  # 0-based：删除第14页起的全部旧模块3内容

# AMA 参考文献（P16: 1-4；P17: 4-5，编号连续）
REFS = {
    1: "李津, 沈永才, 庄辉. 乙型肝炎疫苗2针免疫程序的研究进展. 中国病毒病杂志. 2012;2(4):241-244.",
    2: "Bai X, Chen L, Liu X, et al. Adult hepatitis B virus vaccination coverage in China from 2011 to 2021: a systematic review. Vaccines. 2022;10(6):900. doi:10.3390/vaccines10060900",
    3: "林云, 曹家穗, 何奔, 等. 成人乙肝疫苗全程接种影响因素及免疫效果调查. 中国农村卫生事业管理. 2013;33(6):648-650.",
    4: "Bruxvoort KJ, Slezak J, Huang R, et al. Association of number of doses with hepatitis B vaccine series completion in US adults. JAMA Netw Open. 2020;3(11):e2027577. doi:10.1001/jamanetworkopen.2020.27577",
    5: "LaMori J, Feng X, Pericone CD, et al. Hepatitis vaccination adherence and completion rates and factors associated with low compliance: a claims-based analysis of US adults. PLoS ONE. 2022;17(2):e0264062. doi:10.1371/journal.pone.0264062",
}

# ---------- 基础工具 ----------


def set_run_font(run, size=None, bold=None, color=None, name=FONT, italic=None, sup=False):
    """同时设置 latin 与 east-asian 字体；sup=True 时设为上标。"""
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if color is not None:
        f.color.rgb = color
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)
    if sup:
        rPr.set("baseline", "30000")


def add_textbox(slide, left, top, width, height, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def add_para(tf, first=False):
    if first and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    return p


def para_runs(p, segments, size=11.5, line=1.1, space_after=4, align=PP_ALIGN.LEFT):
    """segments: list of (text, bold, color) 或 (text, bold, color, sup)。"""
    p.alignment = align
    p.line_spacing = line
    p.space_after = Pt(space_after)
    for seg in segments:
        text, bold, color = seg[0], seg[1], seg[2]
        sup = seg[3] if len(seg) > 3 else False
        r = p.add_run()
        r.text = text
        set_run_font(r, size=size, bold=bold, color=color, sup=sup)
    return p


def add_rect(slide, left, top, width, height, color, line=None):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    sh.shadow.inherit = False
    return sh


def add_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)


def base_scaffold(slide, title, subtitle=None):
    """内容页公共元素：顶部色条 + 标题 + 红色竖条。"""
    add_rect(slide, 0, 0, SLIDE_W, 0.14, RED)
    add_rect(slide, 0.62, 0.42, 0.06, 0.40, RED)
    _, tf = add_textbox(slide, 0.78, 0.38, 11.9, 0.48)
    para_runs(tf.paragraphs[0], [(title, True, RED)], size=22, line=1.0, space_after=0)
    if subtitle:
        _, tf2 = add_textbox(slide, 0.78, 0.92, 11.9, 0.30)
        para_runs(tf2.paragraphs[0], [(subtitle, False, BLACK)], size=11, line=1.0, space_after=0)


def add_footnote(slide, text, top=7.05):
    _, tf = add_textbox(slide, 0.62, top, 12.1, 0.34)
    para_runs(tf.paragraphs[0], [(text, False, BLACK)], size=9.5, line=1.0, space_after=0)


def add_refs(slide, nums, top, size=9.5):
    """页内 AMA 参考文献列表（灰色小字）。nums: 序号列表。"""
    _, tf = add_textbox(slide, 0.62, top, 12.1, 1.6)
    # 标题行
    p = add_para(tf, first=True)
    para_runs(p, [("参考文献", True, GRAY)], size=size, line=1.0, space_after=2)
    for n in nums:
        p = add_para(tf)
        para_runs(p, [(f"{n}. {REFS[n]}", False, GRAY)], size=size, line=1.05, space_after=2)


def add_table(
    slide,
    left,
    top,
    width,
    data,
    col_ratios=None,
    header=True,
    font_size=10.5,
    header_size=10.5,
    row_h=0.34,
    header_h=0.36,
    aligns=None,
    bold_col=None,
):
    """data: list of rows；单元格为 str 或 [(text, sup), ...] 富文本段列表。"""
    rows, cols = len(data), len(data[0])
    heights = [Inches(header_h if (header and i == 0) else row_h) for i in range(rows)]
    gf = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(sum(h.inches for h in heights))
    )
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    tblPr.set("firstRow", "1" if header else "0")
    tblPr.set("bandRow", "0")
    if col_ratios:
        total = sum(col_ratios)
        for i, ratio in enumerate(col_ratios):
            tbl.columns[i].width = Emu(int(Inches(width) * ratio / total))
    for i, h in enumerate(heights):
        tbl.rows[i].height = h
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = aligns[ci] if aligns else PP_ALIGN.LEFT
            p.line_spacing = 1.0
            if header and ri == 0:
                r = p.add_run()
                r.text = str(val)
                set_run_font(r, size=header_size, bold=True, color=WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RED
            else:
                if isinstance(val, list):
                    for text, sup in val:
                        r = p.add_run()
                        r.text = text
                        set_run_font(r, size=font_size, bold=False, color=BLACK, sup=sup)
                else:
                    is_bold = bold_col is not None and ci == bold_col
                    r = p.add_run()
                    r.text = str(val)
                    set_run_font(r, size=font_size, bold=is_bold, color=BLACK)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else LIGHT
            _set_cell_border(cell)
    return gf


def _set_cell_border(cell, color="D0D0D0", width_pt="0.5"):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
        ln = parse_xml(
            '<a:{} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{}" cap="flat">'
            '<a:solidFill><a:srgbClr val="{}"/></a:solidFill>'
            '<a:prstDash val="solid"/></a:{}>'.format(
                tag.split(":")[1], int(float(width_pt) * 12700), color, tag.split(":")[1]
            )
        )
        tcPr.append(ln)


def section_cover(prs, num, title, subtitle, meta_lines):
    s = add_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, 0.50, RED)
    _, tf = add_textbox(s, 0.78, 2.05, 12.0, 1.0)
    para_runs(
        tf.paragraphs[0],
        [(f"{num} ", True, RED), (title, True, RED)],
        size=32,
        line=1.05,
        space_after=0,
    )
    _, tf2 = add_textbox(s, 0.78, 3.25, 12.0, 0.45)
    para_runs(tf2.paragraphs[0], [(subtitle, False, BLACK)], size=15, line=1.1, space_after=0)
    _, tf3 = add_textbox(s, 0.78, 4.20, 12.0, 1.4)
    for i, (label, text) in enumerate(meta_lines):
        p = add_para(tf3, first=(i == 0))
        para_runs(
            p,
            [("■ ", True, RED), (label, True, BLACK), ("　" + text, False, BLACK)],
            size=12.5,
            line=1.15,
            space_after=6,
        )
    return s


def delete_old_module(prs, from_idx):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for sldId in slides[from_idx:]:
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


# ---------- 页面构建 ----------


def slide_framework(prs):
    """页2：总体框架——分人群接种程序策略（Ⅲ期设计）。"""
    s = add_slide(prs)
    base_scaffold(
        s,
        "总体框架：分人群接种程序策略（Ⅲ期设计）",
        "数据来源：TVAX-009 Ⅲ期临床试验启动前沟通交流材料（2026.8）",
    )
    data = [
        ["人群", "免疫程序（Ⅲ期）", "试验定位", "临床意义"],
        [
            "18-59岁\n免前阴性人群",
            "试验组确定2剂；第2剂接种时间点Ⅲ期确认：\n0,1月程序或0,2月程序\n（第2剂于首剂后1-2个月内接种）；\n对照组：3剂，0,1,6月程序",
            "2种2剂方案（0,1月、0,2月）与\n3剂对照并行验证（随机、盲法、\n阳性对照、非劣效）",
            "减少1剂、缩短接种周期，\n有望提高全程接种完成率，\n并更早建立血清保护",
        ],
        [
            "60岁及以上\n免前阴性人群",
            "试验组：3剂，0,1,6月程序；\n对照组：3剂，0,1,6月程序",
            "经典3剂程序不变，\n仅比较疫苗本身",
            "免疫衰老人群以3剂提供\n充分免疫刺激，保障应答\n水平与持久性",
        ],
    ]
    add_table(
        s,
        0.62,
        1.42,
        12.1,
        data,
        col_ratios=[1.1, 2.1, 1.6, 1.7],
        font_size=10.5,
        header_size=10.5,
        row_h=1.08,
        header_h=0.38,
        bold_col=0,
    )
    _, tf = add_textbox(s, 0.62, 5.05, 12.1, 1.9)
    paras = [
        [
            ("主要终点：", True, RED),
            (
                "首剂接种后7个月，18-59岁人群中0,1月、0,2月程序试验组与对照3剂组的抗-HBs阳转率；60岁及以上人群中试验组与对照组的阳转率。",
                False,
                BLACK,
            ),
        ],
        [
            ("检验假设：", True, RED),
            ("阳转率非劣效检验，率差双侧95%置信区间下限大于-5%（单侧α=0.025）。", False, BLACK),
        ],
        [
            ("设计逻辑：", True, RED),
            (
                "18-59岁试验疫苗已确定采用2剂程序，Ⅲ期并行验证0,1月与0,2月方案，以确认第2剂接种时间点；≥60岁维持经典3剂程序，与对照在同程序下比较疫苗本身。",
                False,
                BLACK,
            ),
        ],
    ]
    for i, segs in enumerate(paras):
        p = add_para(tf, first=(i == 0))
        para_runs(p, segs, size=11, line=1.12, space_after=5)
    return s


def slide_adherence_gap(prs):
    """页3（P16）：现实动因——多剂次、长周期程序依从性不足 + AMA参考文献。"""
    s = add_slide(prs)
    base_scaffold(
        s, "现实动因：多剂次、长周期程序依从性不足", "循证来源：2-doses/3-doses 依从性支持文献"
    )
    data = [
        ["研究（文献）", "设计与人群", "关键结果"],
        [
            [("李津等, 2012\n（综述）", False)],
            "美国18-49岁成人",
            [("3针全程接种率：高危人群41.8%，\n一般人群31.2%", False), ("1", True)],
        ],
        [
            [("Bai X, et al.\nVaccines 2022\n（系统综述）", False)],
            "21项研究、约34.7万中国成人",
            [
                ("成人乙肝疫苗接种率合并26.27%；\n≥40岁人群仅17.09%（<40岁为36.93%）", False),
                ("2", True),
            ],
        ],
        [
            [("林云等, 2013\n（现场调查）", False)],
            "浙江嘉兴，18岁以上1,159人\n（主动上门、免费接种）",
            [("0-1-6程序实际全程接种率仅69.11%；\n18-24岁组最低（44.65%）", False), ("3", True)],
        ],
        [
            [("Bruxvoort K, et al.\nJAMA Netw Open 2020\n（回顾性队列）", False)],
            "美国成人10,888例\n（Kaiser Permanente）",
            [("3剂铝佐剂疫苗起始者中仅26%\n在推荐窗口内完成全程", False), ("4", True)],
        ],
    ]
    add_table(
        s,
        0.62,
        1.42,
        12.1,
        data,
        col_ratios=[1.25, 1.55, 1.85],
        font_size=10.5,
        header_size=10.5,
        row_h=0.70,
        header_h=0.36,
        bold_col=0,
    )
    _, tf = add_textbox(s, 0.62, 4.60, 12.1, 0.9)
    segs = [
        ("小结：", True, RED),
        (
            "0,1,6月程序跨6个月、需3次就诊，是成人全程接种率低的主要障碍之一；即便免费、主动上门服务，全程接种率仍不足70%。",
            False,
            BLACK,
        ),
        ("1-4", True, BLACK, True),
    ]
    para_runs(tf.paragraphs[0], segs, size=11, line=1.12, space_after=0)
    add_refs(s, [1, 2, 3, 4], top=5.62)
    return s


def slide_2dose_benefit(prs):
    """页4（P17）：2剂程序与完成率改善的循证支持 + AMA参考文献。"""
    s = add_slide(prs)
    base_scaffold(
        s, "循证支持：2剂程序与全程接种完成率改善", "同类CpG佐剂2剂乙肝疫苗的真实世界证据"
    )
    data = [
        ["研究", "设计/样本", "2剂完成率", "3剂完成率", "关键发现"],
        [
            [("Bruxvoort K, et al.\nJAMA Netw Open 2020", False), ("4", True)],
            "美国成人10,888例\n（HepB-CpG 2剂 vs\nHepB-alum 3剂）",
            "45%",
            "26%",
            "2剂起始者完成全程的\n可能性较3剂高约77%",
        ],
        [
            [("LaMori J, et al.\nPLOS ONE 2022", False), ("5", True)],
            "美国商保数据库\n356,828例成人",
            "44.8%\n（24个月）",
            "37.3%\n（24个月）",
            "依从率32.2% vs 14.3%；2剂完成曲线\n约6个月达平台，3剂约需12个月",
        ],
    ]
    add_table(
        s,
        0.62,
        1.42,
        12.1,
        data,
        col_ratios=[1.15, 1.35, 0.75, 0.75, 1.75],
        font_size=10.5,
        header_size=10.5,
        row_h=0.78,
        header_h=0.36,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT],
    )
    _, tf = add_textbox(s, 0.62, 3.80, 12.1, 1.3)
    paras = [
        [
            ("证据解读：", True, RED),
            (
                "在真实世界条件下，同为CpG佐剂的2剂乙肝疫苗（HEPLISAV-B）较3剂铝佐剂疫苗的全程完成率更高（45% vs 26%",
                False,
                BLACK,
            ),
            ("4", False, BLACK, True),
            ("），完成所需时间更短（约6个月 vs 约12个月", False, BLACK),
            ("5", False, BLACK, True),
            ("）。", False, BLACK),
        ],
        [
            ("对TVAX-009的意义：", True, RED),
            (
                "18-59岁人群采用CpG+铝佐剂配方，具备2剂程序的同类证据基础；Ⅲ期并行设置0,1月与0,2月两种2剂方案与3剂对照，将以注册研究证据直接支持第2剂接种时间点的选择。",
                False,
                BLACK,
            ),
        ],
    ]
    for i, segs in enumerate(paras):
        p = add_para(tf, first=(i == 0))
        para_runs(p, segs, size=11, line=1.12, space_after=5)
    add_refs(s, [4, 5], top=5.35)
    return s


def slide_1859_early(prs):
    """页5：18-59岁——早保护（Ⅱ期免疫原性，含精确P值）。"""
    s = add_slide(prs)
    base_scaffold(
        s,
        "18-59岁人群：2剂程序可更早建立血清保护（Ⅱ期）",
        "18-59岁免前阴性人群接种后抗-HBs阳转率（PPS）",
    )
    data = [
        ["时间点", "2剂程序（高剂量组）", "对照（3剂，0,1,6月）", "P值"],
        ["M1（1针后）", "54.86%（0,1月）/ 43.92%（0,2月）", "31.76%", "<0.0001 / 0.0310"],
        ["M2（0,1月程序2针后）", "97.89%", "65.28%", "<0.0001"],
        ["M3（0,2月程序2针后）", "100.00%", "69.44%", "<0.0001"],
        ["M6（对照第3剂前，试验组均2剂）", "98.58%", "78.08%", "<0.0001"],
        [
            "M7（对照3剂全免后1个月）",
            "99.30%（0,1月）/ 100.00%（0,2月）",
            "97.22%",
            "0.3707 / 0.0587",
        ],
    ]
    add_table(
        s,
        0.62,
        1.45,
        12.1,
        data,
        col_ratios=[2.1, 1.75, 1.25, 0.95],
        font_size=10.5,
        header_size=10.5,
        row_h=0.40,
        header_h=0.36,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER],
    )
    _, tf = add_textbox(s, 0.62, 4.20, 12.1, 2.7)
    paras = [
        [
            ("早保护：", True, RED),
            (
                "1针后（M1）高剂量组阳转率即高于对照（54.86%/43.92% vs 31.76%，P<0.0001/P=0.0310）；2针后达97.89%-100.00%，各时间点均显著高于同期对照（P<0.0001）。",
                False,
                BLACK,
            ),
        ],
        [
            ("提前建立全程阳转：", True, RED),
            (
                "对照3剂程序需至M7（第3剂后1个月）阳转率方达峰值（97.22%）；以全程阳转建立保护计，2剂程序较对照提前约5个月。",
                False,
                BLACK,
            ),
        ],
        [
            ("应急场景：", True, RED),
            (
                "0,1月程序M2阳转率即达97.89%，0,2月程序为61.49%（P<0.0001），提示0,1月程序适用于需快速获得保护的接种场景。",
                False,
                BLACK,
            ),
        ],
    ]
    for i, segs in enumerate(paras):
        p = add_para(tf, first=(i == 0))
        para_runs(p, segs, size=11, line=1.12, space_after=5)
    add_footnote(
        s,
        "数据来源：TVAX-009 Ⅱ期临床试验。阳转率定义：免前抗-HBs<10 mIU/mL者，免后相应时间点≥10 mIU/mL的受试者百分比。",
    )
    return s


def slide_1859_durability(prs):
    """页6：18-59岁——免疫持久性（Ⅱ期，含精确P值）。"""
    s = add_slide(prs)
    base_scaffold(
        s, "18-59岁人群：2剂程序的抗体持久性（Ⅱ期）", "阳转率、SPR-100与GMC的动态趋势（PPS）"
    )
    data = [
        ["指标", "数据（2剂程序 vs 对照3剂）", "统计学结果"],
        ["阳转率·M12", "97.86%/98.59% vs 96.53%", "P=0.7229/0.4474，差异无统计学意义"],
        ["阳转率·全免后6个月", "99.30%/99.32% vs 96.53%", "P=0.4473/0.1180，差异无统计学意义"],
        ["SPR-100率·M6", "78.38%-83.69% vs 44.52%", "P<0.0001"],
        ["SPR-100率·M12", "77.86%/81.69% vs 84.03%", "P=0.1852/0.5999，差异无统计学意义"],
        ["GMC·峰值", "374.37/486.11 vs 2311.46 mIU/mL", "P<0.0001，对照峰值更高"],
        ["GMC·全免后6个月", "0,2月 418.26 vs 559.72 mIU/mL", "P=0.1243，差异无统计学意义"],
    ]
    add_table(
        s,
        0.62,
        1.45,
        12.1,
        data,
        col_ratios=[1.35, 1.85, 1.65],
        font_size=10.5,
        header_size=10.5,
        row_h=0.40,
        header_h=0.36,
        bold_col=0,
    )
    _, tf = add_textbox(s, 0.62, 4.55, 12.1, 2.4)
    paras = [
        [
            ("客观呈现：", True, RED),
            (
                "2剂程序的GMC峰值低于对照3剂程序（374.37/486.11 vs 2311.46 mIU/mL，P<0.0001），短期内抗体水平不及对照苗。",
                False,
                BLACK,
            ),
        ],
        [
            ("趋势解读：", True, RED),
            (
                "对照3剂程序GMC自峰值（M7）下降最为明显；至全免后6个月，0,2月程序GMC（418.26）与对照（559.72）差异无统计学意义（P=0.1243）。基于现有12个月数据及下降趋势，2剂程序长期持久性预计不劣于3剂程序。",
                False,
                BLACK,
            ),
        ],
        [
            ("证据状态：", True, RED),
            (
                "M12数据为基于已有数据的分析；免疫持久性随访仍在进行，长期结论有待后续数据支持。",
                False,
                BLACK,
            ),
        ],
    ]
    for i, segs in enumerate(paras):
        p = add_para(tf, first=(i == 0))
        para_runs(p, segs, size=11, line=1.12, space_after=5)
    add_footnote(
        s,
        "数据来源：TVAX-009 Ⅱ期临床试验（PPS）。SPR-100：免前<10 mIU/mL者免后≥100 mIU/mL的百分比；GMC单位为mIU/mL。",
    )
    return s


def slide_60_rationale(prs):
    """页7：≥60岁——经典3剂程序不变的依据。"""
    s = add_slide(prs)
    base_scaffold(s, "60岁及以上人群：维持经典3剂程序（0,1,6月）", "程序选择的免疫学与临床依据")
    _, tf = add_textbox(s, 0.62, 1.42, 12.1, 2.0)
    paras = [
        [
            ("Ⅰ期结论：", True, RED),
            (
                "≥60岁人群免疫功能衰退，2剂应答偏弱，需3剂（0,1,6月）程序提供充分免疫刺激。",
                False,
                BLACK,
            ),
        ],
        [
            ("Ⅱ期设计：", True, RED),
            (
                "≥60岁人群按0,1,6月程序接种3剂（高剂量）；接种2剂后阳转率维持在较高水平，但GMC在第3剂接种前后差异较大，第3剂可显著提升抗体水平（M6 GMC 143.50 vs M7 10389.37 mIU/mL）。",
                False,
                BLACK,
            ),
        ],
        [
            ("Ⅲ期设计：", True, RED),
            (
                "≥60岁人群试验组与对照组均维持3剂、0,1,6月程序，在同程序下以非劣效设计（Δ=-5%）比较疫苗本身，不引入程序变量。",
                False,
                BLACK,
            ),
        ],
    ]
    for i, segs in enumerate(paras):
        p = add_para(tf, first=(i == 0))
        para_runs(p, segs, size=11.5, line=1.15, space_after=7)
    data = [
        ["维度", "依据与设计要点"],
        ["免疫学基础", "免疫衰老导致应答减弱，延长/增加免疫刺激（3剂）为该人群的稳妥选择"],
        [
            "Ⅱ期数据支持",
            "高剂量组第3剂接种前GMC已明显高于对照（M6：143.50 vs 9.35 mIU/mL，P<0.0001），第3剂后阳转率与GMC进一步达到较高水平",
        ],
        ["Ⅲ期程序设置", "试验组3剂（0,1,6月） vs 对照3剂（0,1,6月），非劣效界值-5%"],
    ]
    add_table(
        s,
        0.62,
        3.65,
        12.1,
        data,
        col_ratios=[1.0, 3.0],
        font_size=10.5,
        header_size=10.5,
        row_h=0.60,
        header_h=0.34,
        bold_col=0,
    )
    add_footnote(s, "数据来源：TVAX-009 Ⅰ/Ⅱ期临床试验及Ⅲ期方案设计。")
    return s


def slide_60_data(prs):
    """页8：≥60岁——早保护、高应答与持久性（Ⅱ期，含精确P值）。"""
    s = add_slide(prs)
    base_scaffold(
        s,
        "60岁及以上人群：早保护、高应答与持久性（Ⅱ期）",
        "0,1,6月程序：试验组（高剂量）vs 阳性对照，抗-HBs阳转率（PPS）",
    )
    data = [
        ["时间点", "试验组（3剂，高剂量）", "对照组（3剂）", "P值"],
        ["M1（1针后）", "29.33%", "8.11%", "0.0009"],
        ["M2（2针后1个月）", "86.49%", "31.51%", "<0.0001"],
        ["M3（2针后2个月）", "91.89%", "34.25%", "<0.0001"],
        ["M6（第3剂前）", "92.96%", "44.59%", "<0.0001"],
        ["M7（3剂后1个月）", "100.00%", "91.55%", "0.0280"],
        ["M8（3剂后2个月）", "98.59%", "92.96%", "0.2087"],
        ["M12（FAS）", "98.55%", "85.92%", "0.0055"],
    ]
    add_table(
        s,
        0.62,
        1.45,
        12.1,
        data,
        col_ratios=[1.7, 1.5, 1.3, 0.9],
        font_size=10.5,
        header_size=10.5,
        row_h=0.38,
        header_h=0.36,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER],
    )
    _, tf = add_textbox(s, 0.62, 4.75, 12.1, 2.2)
    paras = [
        [
            ("早保护：", True, RED),
            (
                "首剂后6个月内各时间点阳转率均显著高于对照（P≤0.0009）；对照第3剂前（M6）为92.96% vs 44.59%。",
                False,
                BLACK,
            ),
        ],
        [
            ("高应答：", True, RED),
            (
                "3剂全程后1个月（M7）阳转率达100.00% vs 91.55%（P=0.0280）；GMC同步达峰（10389.37 vs 723.28 mIU/mL，P<0.0001）。",
                False,
                BLACK,
            ),
        ],
        [
            ("持久性：", True, RED),
            (
                "至M12，试验组阳转率维持98.55%，对照降至85.92%（P=0.0055）；M12 GMC为2405.29 vs 177.88 mIU/mL（P<0.0001）。",
                False,
                BLACK,
            ),
        ],
    ]
    for i, segs in enumerate(paras):
        p = add_para(tf, first=(i == 0))
        para_runs(p, segs, size=11, line=1.12, space_after=5)
    add_footnote(s, "数据来源：TVAX-009 Ⅱ期临床试验，60岁及以上免前阴性人群（PPS；M12为FAS）。")
    return s


def slide_summary(prs):
    """页9：模块小结（减剂次表述，含精确P值）。"""
    s = add_slide(prs)
    base_scaffold(s, "模块小结：新佐剂乙肝疫苗与减剂次的临床意义")
    data = [
        ["人群", "程序策略", "早保护", "应答与持久性", "依从性意义"],
        [
            "18-59岁",
            "2剂\n（第2剂时间点Ⅲ期确认：\n0,1月 vs 0,2月）",
            "M1阳转率54.86%/43.92% vs\n31.76%（P<0.0001）；\n全程阳转较对照提前约5个月",
            "M12阳转率97.86%/98.59% vs\n96.53%（P=0.7229/0.4474）；\n全免后6个月0,2月GMC与对照\n差异无统计学意义（P=0.1243）",
            "减少1剂、缩短周期；\n真实世界2剂完成率45% vs\n3剂26%（同类CpG疫苗）",
        ],
        [
            "≥60岁",
            "经典3剂（0,1,6月）\nⅢ期同程序对照",
            "对照第3剂前阳转率\n92.96% vs 44.59%（P<0.0001）",
            "M7达峰100.00% vs 91.55%\n（P=0.0280）；M12维持98.55%\nvs 85.92%（P=0.0055）",
            "以充分免疫刺激保障\n免疫衰老人群应答水平",
        ],
    ]
    add_table(
        s,
        0.62,
        1.42,
        12.1,
        data,
        col_ratios=[0.75, 1.5, 1.7, 1.75, 1.6],
        font_size=10,
        header_size=10.5,
        row_h=1.30,
        header_h=0.36,
        bold_col=0,
    )
    _, tf = add_textbox(s, 0.62, 4.85, 12.1, 1.9)
    paras = [
        [
            ("结论：", True, RED),
            (
                "减剂次策略采用分人群路径——18-59岁确定2剂程序，Ⅲ期确认第2剂接种时间点；≥60岁维持经典3剂程序，保障应答水平；",
                False,
                BLACK,
            ),
        ],
        [
            ("  ", False, BLACK),
            ("Ⅲ期以非劣效设计（Δ=-5%）同步验证，为最终程序推荐提供注册研究证据。", False, BLACK),
        ],
    ]
    for i, segs in enumerate(paras):
        p = add_para(tf, first=(i == 0))
        para_runs(p, segs, size=11, line=1.12, space_after=5)
    add_footnote(
        s,
        "注：Ⅱ期M12数据为基于已有数据的分析，免疫持久性随访仍在进行；引用文献为观察性证据，外推需结合目标人群特点。",
    )
    return s


# ---------- 主流程 ----------


def main():
    prs = Presentation(SRC)
    n_before = len(prs.slides)
    delete_old_module(prs, DELETE_FROM)
    n_mid = len(prs.slides)
    print(f"删除旧模块3：{n_before} -> {n_mid} 页（删除 {n_before - n_mid} 页）")

    section_cover(
        prs,
        "03",
        "新佐剂乙肝疫苗与减剂次的临床意义",
        "基于TVAX-009（重组乙型肝炎疫苗（汉逊酵母，CpG和铝佐剂））Ⅰ/Ⅱ期数据与Ⅲ期设计",
        [
            ("18-59岁人群：", "试验疫苗确定采用2剂程序，Ⅲ期确认第2剂接种时间点（0,1月 vs 0,2月）"),
            ("60岁及以上人群：", "维持经典3剂程序（0,1,6月），早保护、高应答、免疫持久性良好"),
            ("数据来源：", "TVAX-009 Ⅲ期启动前沟通交流材料（2026.8）及5篇依从性支持文献"),
        ],
    )
    slide_framework(prs)
    slide_adherence_gap(prs)
    slide_2dose_benefit(prs)
    slide_1859_early(prs)
    slide_1859_durability(prs)
    slide_60_rationale(prs)
    slide_60_data(prs)
    slide_summary(prs)

    prs.save(DST)
    n_after = len(prs.slides)
    print(f"新模块共 {n_after - n_mid} 页；输出 {n_after} 页 -> {DST}")


if __name__ == "__main__":
    main()
