"""完整验证 V4 Slide 7/8/9 的表格内容"""

import sys

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation

V4 = r"E:/Cursor Project/2-Scientific-Skills-for-Clinical_Trial/review_materials/60岁以上乙肝流调-CpG佐剂安全性-新佐剂减剂次临床意义-20260827-v4.pptx"
prs = Presentation(V4)
for idx in (6, 8):  # Slide 7 (Phase1/2), Slide 9 (Phase3)
    slide = prs.slides[idx]
    print("=" * 90)
    print("SLIDE %d" % (idx + 1))
    print("=" * 90)
    for shp in slide.shapes:
        if shp.has_table:
            tbl = shp.table
            for ri in range(len(tbl.rows)):
                print("--- 行 %d ---" % ri)
                for ci in range(len(tbl.columns)):
                    txt = tbl.cell(ri, ci).text_frame.text
                    print("[col%d] %s" % (ci, txt))
                    print()
