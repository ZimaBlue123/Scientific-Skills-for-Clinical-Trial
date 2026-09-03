# -*- coding: utf-8 -*-
"""校验新增糖尿病亚组 6 页 PPT 的页序、标题、描述与数字一致性。"""
import json, sys
from pptx import Presentation
from pptx.util import Emu

PPT = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\TVAX-009项目3期临床试验启动前沟通交流ppt-20260903（临床部分-新增糖尿病亚组）.pptx"
JSON = r"E:\Cursor Project\2-Scientific-Skills-for-Clinical_Trial\review_materials\diabetes_immuno_result.json"

prs = Presentation(PPT)
data = json.load(open(JSON, encoding="utf-8"))
print("total slides:", len(prs.slides))

def cell_text(c):
    return c.text.strip()

def dump_slide(idx):
    s = prs.slides[idx]
    print("\n" + "=" * 70)
    print(f"[{idx}] (0-based) / 1-based {idx+1}")
    for sh in s.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt:
                print(f"  TEXT[{sh.shape_type}] ({Emu(sh.left).inches:.2f},{Emu(sh.top).inches:.2f}): {txt[:200]!r}")
        if sh.has_table:
            tbl = sh.table
            print(f"  TABLE {len(tbl.rows)}r x {len(tbl.columns)}c")
            for r_i, row in enumerate(tbl.rows):
                cells = [cell_text(c) for c in row.cells]
                print(f"    R{r_i}: {cells}")

# 打印插入位置前后 + 新增 6 页
for idx in [40, 41, 42, 43, 44, 45, 46, 47, 48]:
    if idx < len(prs.slides):
        dump_slide(idx)
